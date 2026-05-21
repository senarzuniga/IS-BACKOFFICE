from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.util import Inches, Pt


BASE_DIR = Path(r"C:\Users\Inaki Senar\Documents\GitHub\IS-BACKOFFICE")
TEMPLATE_PATH = BASE_DIR / "ingecart-marketing-kit" / "assets" / "presentations" / "Ingecart_Corporate_Slide.pptx"
OUTPUT_PATH = Path(r"C:\Users\Inaki Senar\Documents\INGECART\MARKETING\CONTENT\INGECART_CORPORATE_PRESENTATION_PROFESSIONAL.pptx")

INGECART_RED = RGBColor(204, 0, 0)
INGECART_DARK = RGBColor(37, 37, 37)
INGECART_GRAY = RGBColor(102, 102, 102)
INGECART_LIGHT = RGBColor(245, 245, 245)
INGECART_BORDER = RGBColor(220, 220, 220)
INGECART_SOFT_RED = RGBColor(247, 232, 232)
INGECART_SOFT_BLUE = RGBColor(233, 240, 250)
INGECART_SOFT_GREEN = RGBColor(233, 247, 238)

TITLE_FONT = "Montserrat"
BODY_FONT = "Arial"


def remove_slide(prs: Presentation, index: int) -> None:
    slide_id = prs.slides._sldIdLst[index]
    rel_id = slide_id.rId
    prs.part.drop_rel(rel_id)
    del prs.slides._sldIdLst[index]


def clear_presentation(prs: Presentation) -> None:
    for index in range(len(prs.slides) - 1, -1, -1):
        remove_slide(prs, index)


def set_background(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, shape_type, left, top, width, height, fill=None, line=None, radius=None):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    if fill is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line is not None:
        shape.line.color.rgb = line
    else:
        shape.line.color.rgb = fill if fill is not None else INGECART_BORDER
    if radius is not None:
        shape.adjustments[0] = radius
    return shape


def add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text,
    *,
    font_name=BODY_FONT,
    font_size=18,
    color=INGECART_DARK,
    bold=False,
    italic=False,
    align=PP_ALIGN.LEFT,
    fit=True,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.clear()

    lines = text.split("\n")
    for index, line in enumerate(lines):
        paragraph = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        paragraph.alignment = align
        run = paragraph.add_run()
        run.text = line
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.italic = italic

    if fit:
        tf.auto_size = True
    return box


def add_header(slide, title, subtitle=None):
    add_textbox(
        slide,
        Inches(0.55),
        Inches(0.35),
        Inches(8.9),
        Inches(0.55),
        title,
        font_name=TITLE_FONT,
        font_size=24,
        color=INGECART_DARK,
        bold=True,
    )
    accent = add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(1.0), Inches(1.15), Inches(0.07), fill=INGECART_RED, line=INGECART_RED)
    accent.line.width = Pt(0)
    if subtitle:
        add_textbox(
            slide,
            Inches(0.58),
            Inches(1.08),
            Inches(9.0),
            Inches(0.35),
            subtitle,
            font_name=BODY_FONT,
            font_size=11,
            color=INGECART_GRAY,
        )


def add_note(slide, text: str):
    slide.notes_slide.notes_text_frame.text = text


def add_card(slide, left, top, width, height, title, body, fill=INGECART_LIGHT, title_color=INGECART_DARK, body_color=INGECART_GRAY, title_size=16, body_size=12):
    card = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height, fill=fill, line=INGECART_BORDER)
    card.line.width = Pt(1.2)
    add_textbox(slide, left + Inches(0.16), top + Inches(0.12), width - Inches(0.32), Inches(0.28), title, font_name=TITLE_FONT, font_size=title_size, color=title_color, bold=True)
    add_textbox(slide, left + Inches(0.16), top + Inches(0.48), width - Inches(0.32), height - Inches(0.56), body, font_name=BODY_FONT, font_size=body_size, color=body_color)
    return card


def add_title_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, INGECART_DARK)

    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(9.15), Inches(0), Inches(0.45), Inches(7.5), fill=INGECART_RED, line=INGECART_RED)
    add_shape(slide, MSO_SHAPE.OVAL, Inches(7.6), Inches(0.55), Inches(2.1), Inches(2.1), fill=RGBColor(60, 60, 60), line=RGBColor(90, 90, 90))
    add_shape(slide, MSO_SHAPE.OVAL, Inches(8.05), Inches(1.0), Inches(1.2), Inches(1.2), fill=INGECART_RED, line=INGECART_RED)

    add_textbox(
        slide,
        Inches(0.8),
        Inches(1.35),
        Inches(7.4),
        Inches(1.8),
        "Ingecart: Independent Industrial Engineering",
        font_name=TITLE_FONT,
        font_size=30,
        color=RGBColor(255, 255, 255),
        bold=True,
    )
    add_textbox(
        slide,
        Inches(0.83),
        Inches(3.05),
        Inches(7.15),
        Inches(1.0),
        "Transforming Corrugated & Packaging Operations into Predictable Profit Centers",
        font_name=BODY_FONT,
        font_size=18,
        color=RGBColor(232, 232, 232),
    )
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.83), Inches(4.0), Inches(1.5), Inches(0.06), fill=INGECART_RED, line=INGECART_RED)

    add_textbox(
        slide,
        Inches(0.82),
        Inches(5.0),
        Inches(7.3),
        Inches(0.55),
        '"Knowledge is the foundation for growth" – Diego Garcia',
        font_name=BODY_FONT,
        font_size=14,
        color=RGBColor(244, 244, 244),
        italic=True,
    )
    add_textbox(slide, Inches(8.05), Inches(6.65), Inches(1.2), Inches(0.25), "May 13, 2026", font_size=10, color=RGBColor(240, 240, 240), align=PP_ALIGN.RIGHT)
    add_note(slide, "Good morning and thank you for your time. Today I will show how Ingecart turns industrial operations into predictable profit centers through applied knowledge, engineering discipline, and integrated technology. We do not sell machines; we solve operational problems with measurable results.")


def add_venn_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, RGBColor(255, 255, 255))
    add_header(slide, "More Than an Integrator. Your Strategic Partner.")

    left_circle = add_shape(slide, MSO_SHAPE.OVAL, Inches(1.3), Inches(1.75), Inches(3.2), Inches(3.2), fill=INGECART_SOFT_BLUE, line=INGECART_BORDER)
    middle_circle = add_shape(slide, MSO_SHAPE.OVAL, Inches(3.4), Inches(1.75), Inches(3.2), Inches(3.2), fill=INGECART_SOFT_GREEN, line=INGECART_BORDER)
    right_circle = add_shape(slide, MSO_SHAPE.OVAL, Inches(2.35), Inches(3.15), Inches(3.2), Inches(3.2), fill=INGECART_SOFT_RED, line=INGECART_BORDER)
    for shape in (left_circle, middle_circle, right_circle):
        shape.line.width = Pt(1.2)

    add_textbox(slide, Inches(1.7), Inches(2.35), Inches(2.3), Inches(0.55), "Independent\nEngineering", font_name=TITLE_FONT, font_size=20, color=INGECART_DARK, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(4.0), Inches(2.35), Inches(2.2), Inches(0.55), "Automation\nIntegration", font_name=TITLE_FONT, font_size=20, color=INGECART_DARK, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(2.55), Inches(4.0), Inches(2.6), Inches(0.6), "Industrial\nIntelligence", font_name=TITLE_FONT, font_size=20, color=INGECART_DARK, bold=True, align=PP_ALIGN.CENTER)
    add_shape(slide, MSO_SHAPE.OVAL, Inches(3.55), Inches(3.1), Inches(1.35), Inches(1.35), fill=INGECART_RED, line=INGECART_RED)
    add_textbox(slide, Inches(3.22), Inches(3.47), Inches(2.0), Inches(0.5), "Knowledge-Driven\nIndustrial Transformation Platform", font_name=BODY_FONT, font_size=13, color=RGBColor(255, 255, 255), bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1.3), Inches(1.62), Inches(3.0), Inches(0.3), "Unbiased technical auditing", font_size=10, color=INGECART_GRAY, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(3.55), Inches(1.62), Inches(3.0), Inches(0.3), "Ingetrans, Retal, robotics, best-of-breed", font_size=10, color=INGECART_GRAY, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(2.45), Inches(5.95), Inches(2.9), Inches(0.3), "Ing_PRO analytics, AI agents", font_size=10, color=INGECART_GRAY, align=PP_ALIGN.CENTER)
    add_note(slide, "Ingecart occupies a unique hybrid position: independent engineering, automation integration, and industrial intelligence. No competitor combines all three. This is the strategic advantage we bring to corrugated and packaging operators.")


def add_footprint_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, RGBColor(255, 255, 255))
    add_header(slide, "28 Years of Verifiable Results in Your Industry")

    card_w = Inches(2.25)
    card_h = Inches(1.35)
    x1, x2 = Inches(0.75), Inches(5.05)
    y1, y2 = Inches(1.7), Inches(3.35)
    data = [
        (x1, y1, "+1,268", "Projects Completed"),
        (x2, y1, "194", "Active Installations"),
        (x1, y2, "26", "International Agreements"),
        (x2, y2, "65%", "of our base is Corrugated"),
    ]
    for x, y, number, label in data:
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, card_w, card_h, fill=INGECART_LIGHT, line=INGECART_BORDER)
        add_textbox(slide, x + Inches(0.15), y + Inches(0.1), Inches(2.0), Inches(0.55), number, font_name=TITLE_FONT, font_size=24, color=INGECART_RED, bold=True)
        add_textbox(slide, x + Inches(0.15), y + Inches(0.65), Inches(2.0), Inches(0.35), label, font_name=BODY_FONT, font_size=12, color=INGECART_GRAY)

    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(5.05), Inches(6.55), Inches(1.35), fill=RGBColor(250, 250, 250), line=INGECART_BORDER)
    add_textbox(slide, Inches(1.0), Inches(5.3), Inches(6.0), Inches(0.35), "[IMAGE PLACEHOLDER: Stylized world map with pins on Europe, North Africa, and LATAM]", font_size=13, color=INGECART_GRAY, italic=True, align=PP_ALIGN.CENTER)
    add_note(slide, "Numbers matter. We bring 28 years of continuous operation, more than 1,200 projects completed, 194 active installations, 26 international agreements, and deep specialization in corrugated and packaging. This is proven sector credibility.")


def add_dual_strategy_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, RGBColor(255, 255, 255))
    add_header(slide, "Two Engines for Your Growth and Ours")

    left_x = Inches(0.75)
    right_x = Inches(5.1)
    top = Inches(1.75)
    card_w = Inches(3.85)
    card_h = Inches(4.65)

    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left_x, top, card_w, card_h, fill=RGBColor(250, 250, 250), line=INGECART_BORDER)
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, right_x, top, card_w, card_h, fill=INGECART_SOFT_RED, line=INGECART_BORDER)

    add_textbox(slide, left_x + Inches(0.2), top + Inches(0.15), Inches(3.2), Inches(0.35), "Account Development (High Value)", font_name=TITLE_FONT, font_size=18, color=INGECART_DARK, bold=True)
    add_textbox(slide, right_x + Inches(0.2), top + Inches(0.15), Inches(3.2), Inches(0.35), "Prospecting (High Velocity)", font_name=TITLE_FONT, font_size=18, color=INGECART_RED, bold=True)

    left_body = "Target: Large groups such as Smurfit Westrock, International Paper, and DS Smith\n\nApproach: Consultative, long-term partnership\n\nObjective: Full transformations from €500K to €1.5M\n\nTimeline: 6-18 months\n\nOutcome: Strategic account growth and deeper wallet share"
    right_body = "Target: Mid-market and regional operators\n\nApproach: Standardized, repeatable solutions\n\nObjective: Fast ROI and agile deployment from €80K to €200K\n\nTimeline: 3-9 months\n\nOutcome: High-velocity pipeline and quick wins"
    add_textbox(slide, left_x + Inches(0.2), top + Inches(0.65), Inches(3.4), Inches(3.5), left_body, font_size=13, color=INGECART_GRAY)
    add_textbox(slide, right_x + Inches(0.2), top + Inches(0.65), Inches(3.4), Inches(3.5), right_body, font_size=13, color=INGECART_DARK)

    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(4.77), Inches(2.15), Inches(0.06), Inches(3.65), fill=INGECART_RED, line=INGECART_RED)
    add_textbox(slide, Inches(2.65), Inches(6.55), Inches(4.4), Inches(0.25), "Dual engine model: strategic depth for large accounts, repeatability for the broader market.", font_size=10, color=INGECART_GRAY, align=PP_ALIGN.CENTER)
    add_note(slide, "We serve two markets with two distinct approaches. For global leaders we provide consultative transformation; for mid-market operators we deliver fast, repeatable ROI. Both engines are essential to growth and resilience.")


def add_pain_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, RGBColor(255, 255, 255))
    add_header(slide, "What Plant Directors and CFOs Tell Us")

    pain_items = [
        ("Unplanned downtime", "2-5 times per week. Cost: €5K-15K per stop."),
        ("Forklift congestion", "Adds 15-30% to the production cycle."),
        ("Waste and trim cost", "€50K-200K annually lost."),
        ("End-of-line bottleneck", "3-5 FTEs plus capacity limitation."),
        ("No actionable visibility", "Reactive decisions and no real optimization."),
    ]
    positions = [
        (Inches(0.75), Inches(1.75)),
        (Inches(0.75), Inches(2.7)),
        (Inches(0.75), Inches(3.65)),
        (Inches(5.1), Inches(1.75)),
        (Inches(5.1), Inches(2.7)),
    ]
    for (title, body), (x, y) in zip(pain_items, positions):
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(4.15), Inches(0.78), fill=INGECART_LIGHT, line=INGECART_BORDER)
        add_shape(slide, MSO_SHAPE.OVAL, x + Inches(0.15), y + Inches(0.16), Inches(0.38), Inches(0.38), fill=INGECART_RED, line=INGECART_RED)
        add_textbox(slide, x + Inches(0.68), y + Inches(0.08), Inches(3.25), Inches(0.22), title, font_name=TITLE_FONT, font_size=14, color=INGECART_DARK, bold=True)
        add_textbox(slide, x + Inches(0.68), y + Inches(0.32), Inches(3.25), Inches(0.3), body, font_name=BODY_FONT, font_size=11, color=INGECART_GRAY)

    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.1), Inches(3.65), Inches(4.15), Inches(1.7), fill=INGECART_RED, line=INGECART_RED)
    add_textbox(slide, Inches(5.35), Inches(4.02), Inches(3.65), Inches(0.55), "Does any of this sound familiar?", font_name=TITLE_FONT, font_size=18, color=RGBColor(255, 255, 255), bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(5.35), Inches(4.45), Inches(3.65), Inches(0.35), "These are profit leaks, not normal operating costs.", font_name=BODY_FONT, font_size=11, color=RGBColor(255, 255, 255), align=PP_ALIGN.CENTER)
    add_note(slide, "We start where the customer feels the pain. Unplanned stops, forklift congestion, waste, bottlenecks, and poor visibility are not normal. They are operational leaks that hurt margin, service, and safety.")


def add_comparison_table(slide, title, rows, col_widths, left_header, right_header):
    add_header(slide, title)
    table = slide.shapes.add_table(len(rows) + 1, 2, Inches(0.75), Inches(1.75), Inches(8.5), Inches(4.4)).table
    table.columns[0].width = Inches(4.15)
    table.columns[1].width = Inches(4.35)
    headers = [left_header, right_header]
    for col_index, header_text in enumerate(headers):
        cell = table.cell(0, col_index)
        cell.text = header_text
        fill = cell.fill
        fill.solid()
        fill.fore_color.rgb = INGECART_RED if col_index == 1 else INGECART_DARK
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        for run in cell.text_frame.paragraphs[0].runs:
            run.font.name = TITLE_FONT
            run.font.size = Pt(13)
            run.font.bold = True
            run.font.color.rgb = RGBColor(255, 255, 255)

    for row_index, (left_text, right_text) in enumerate(rows, start=1):
        for col_index, text in enumerate([left_text, right_text]):
            cell = table.cell(row_index, col_index)
            cell.text = text
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(255, 255, 255) if row_index % 2 else INGECART_LIGHT
            cell.text_frame.word_wrap = True
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            for run in p.runs:
                run.font.name = BODY_FONT
                run.font.size = Pt(11)
                run.font.color.rgb = INGECART_DARK if col_index == 1 else INGECART_GRAY


def add_value_prop_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, RGBColor(255, 255, 255))
    add_header(slide, "The Promise We Keep", "We do not sell machines. We solve operational pain points through applied knowledge and integrated technology.")

    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(1.8), Inches(7.9), Inches(2.0), fill=INGECART_LIGHT, line=INGECART_BORDER)
    add_textbox(slide, Inches(1.2), Inches(2.25), Inches(7.3), Inches(0.9), '"We do not sell machines. We solve operational pain points through applied knowledge and integrated technology."', font_name=TITLE_FONT, font_size=24, color=INGECART_DARK, bold=True, italic=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1.2), Inches(3.3), Inches(7.3), Inches(0.25), "— Ingecart Industrial Engineering", font_name=BODY_FONT, font_size=12, color=INGECART_RED, align=PP_ALIGN.CENTER)

    pillars = [
        (Inches(0.85), "Audit + Diagnosis", "Find the real bottleneck"),
        (Inches(3.3), "Design + Deployment", "Turnkey integrated automation"),
        (Inches(5.75), "Intelligence + Optimization", "Ing_PRO, predictive maintenance, AI"),
    ]
    for x, title, body in pillars:
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(4.35), Inches(2.15), Inches(1.3), fill=RGBColor(250, 250, 250), line=INGECART_BORDER)
        add_shape(slide, MSO_SHAPE.RECTANGLE, x, Inches(4.35), Inches(2.15), Inches(0.07), fill=INGECART_RED, line=INGECART_RED)
        add_textbox(slide, x + Inches(0.12), Inches(4.55), Inches(1.9), Inches(0.25), title, font_name=TITLE_FONT, font_size=13, color=INGECART_DARK, bold=True, align=PP_ALIGN.CENTER)
        add_textbox(slide, x + Inches(0.12), Inches(4.9), Inches(1.9), Inches(0.35), body, font_name=BODY_FONT, font_size=10, color=INGECART_GRAY, align=PP_ALIGN.CENTER)
    add_note(slide, "Read the promise carefully. We do not lead with products. We lead with problem solving. First we audit and diagnose, then we design and deploy, and finally we provide intelligence to keep improving the operation.")


def add_financial_impact_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, RGBColor(255, 255, 255))
    add_header(slide, "Documented & Verified ROI. No Speculation.")

    metrics = [
        ("Reduce unplanned downtime", "25-40%", 500),
        ("Eliminate manual intralogistics", "€80K-120K", 120),
        ("Waste reduction savings", "€50K-100K", 100),
    ]
    max_width = 5.25
    y_start = 1.95
    for index, (label, value, max_val) in enumerate(metrics):
        y = y_start + index * 0.95
        add_textbox(slide, Inches(0.8), y + Inches(0.02), Inches(2.4), Inches(0.25), label, font_name=BODY_FONT, font_size=12, color=INGECART_DARK, bold=True)
        add_textbox(slide, Inches(0.8), y + Inches(0.28), Inches(1.0), Inches(0.25), value, font_name=TITLE_FONT, font_size=12, color=INGECART_RED, bold=True)
        add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(3.25), y + Inches(0.02), Inches(max_width), Inches(0.35), fill=INGECART_LIGHT, line=INGECART_BORDER)
        bar_width = Inches(max_width * (max_val / 500))
        add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(3.25), y + Inches(0.02), bar_width, Inches(0.35), fill=INGECART_RED, line=INGECART_RED)

    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(4.95), Inches(7.7), Inches(0.95), fill=INGECART_DARK, line=INGECART_DARK)
    add_textbox(slide, Inches(0.95), Inches(5.18), Inches(7.35), Inches(0.25), "Total Annual Impact: €480K-€1.15M operational improvement", font_name=TITLE_FONT, font_size=18, color=RGBColor(255, 255, 255), bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.95), Inches(5.45), Inches(7.35), Inches(0.2), "We model your scenario before you invest.", font_name=BODY_FONT, font_size=11, color=RGBColor(240, 240, 240), italic=True, align=PP_ALIGN.CENTER)
    add_note(slide, "This is the CFO language. We quantify the value before the first euro is spent. The headline impact comes from reduced downtime, lower manual handling, and material savings from waste reduction.")


def add_operational_impact_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, RGBColor(255, 255, 255))
    add_header(slide, "Results You Can See and Measure Every Day")

    kpis = [
        ("+15-30%", "Throughput", "without major new CAPEX"),
        ("-50%", "Safety Incidents", "manual handling reduction"),
        ("20-40%", "Defect Reduction", "in quality"),
        ("95%+", "On-time Delivery", "reliability"),
    ]
    positions = [
        (Inches(0.85), Inches(1.85)),
        (Inches(4.45), Inches(1.85)),
        (Inches(0.85), Inches(3.45)),
        (Inches(4.45), Inches(3.45)),
    ]
    fills = [INGECART_SOFT_RED, INGECART_SOFT_BLUE, INGECART_SOFT_GREEN, INGECART_LIGHT]
    for (big, label, body), (x, y), fill in zip(kpis, positions, fills):
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.0), Inches(1.25), fill=fill, line=INGECART_BORDER)
        add_textbox(slide, x + Inches(0.18), y + Inches(0.12), Inches(1.05), Inches(0.42), big, font_name=TITLE_FONT, font_size=24, color=INGECART_RED, bold=True)
        add_textbox(slide, x + Inches(1.25), y + Inches(0.16), Inches(1.55), Inches(0.22), label, font_name=TITLE_FONT, font_size=14, color=INGECART_DARK, bold=True)
        add_textbox(slide, x + Inches(0.18), y + Inches(0.72), Inches(2.6), Inches(0.25), body, font_name=BODY_FONT, font_size=10, color=INGECART_GRAY, align=PP_ALIGN.CENTER)

    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.0), Inches(2.0), Inches(2.0), Inches(2.6), fill=RGBColor(250, 250, 250), line=INGECART_BORDER)
    add_textbox(slide, Inches(7.16), Inches(2.25), Inches(1.7), Inches(0.4), "Ing_PRO Dashboard", font_name=TITLE_FONT, font_size=14, color=INGECART_DARK, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(7.15), Inches(2.65), Inches(1.7), Inches(1.55), "[IMAGE PLACEHOLDER: real-time operational dashboard mockup showing throughput, alarms, and KPI trends]", font_name=BODY_FONT, font_size=11, color=INGECART_GRAY, italic=True, align=PP_ALIGN.CENTER)
    add_note(slide, "For plant leadership, the results are visible every day: more throughput, fewer safety incidents, better quality, and more reliable delivery. The dashboard keeps the improvement visible and measurable.")


def add_evidence_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, RGBColor(255, 255, 255))
    add_header(slide, "Not Theory. Practice.")

    cards = [
        ("Spain", "Corrugated Converter", "+35% throughput post-automation", "Payback in 14 months"),
        ("Portugal", "Converting Facility", "€180K investment → €320K annual gain", "ROI 78% in year 1"),
        ("France", "End-of-Line Solution", "Eliminated 2.5 FTEs", "€85K savings per year"),
    ]
    xs = [Inches(0.75), Inches(3.55), Inches(6.35)]
    colors = [INGECART_SOFT_BLUE, INGECART_SOFT_GREEN, INGECART_SOFT_RED]
    for x, (country, title, metric, result), fill in zip(xs, cards, colors):
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.9), Inches(2.55), Inches(3.55), fill=fill, line=INGECART_BORDER)
        add_textbox(slide, x + Inches(0.16), Inches(2.05), Inches(2.2), Inches(0.22), country.upper(), font_name=TITLE_FONT, font_size=11, color=INGECART_RED, bold=True)
        add_textbox(slide, x + Inches(0.16), Inches(2.35), Inches(2.2), Inches(0.55), title, font_name=TITLE_FONT, font_size=15, color=INGECART_DARK, bold=True)
        add_textbox(slide, x + Inches(0.16), Inches(3.15), Inches(2.2), Inches(0.8), metric, font_name=BODY_FONT, font_size=13, color=INGECART_DARK)
        add_shape(slide, MSO_SHAPE.RECTANGLE, x + Inches(0.16), Inches(4.05), Inches(2.2), Inches(0.05), fill=INGECART_RED, line=INGECART_RED)
        add_textbox(slide, x + Inches(0.16), Inches(4.15), Inches(2.2), Inches(0.45), result, font_name=TITLE_FONT, font_size=12, color=INGECART_RED, bold=True)

    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.45), Inches(5.75), Inches(2.9), Inches(0.45), fill=INGECART_DARK, line=INGECART_DARK)
    add_textbox(slide, Inches(3.55), Inches(5.86), Inches(2.7), Inches(0.2), "References available on request", font_name=BODY_FONT, font_size=11, color=RGBColor(255, 255, 255), bold=True, align=PP_ALIGN.CENTER)
    add_note(slide, "Theory is cheap. Results matter. The examples shown are real commercial outcomes: improved throughput, clear financial gains, and measurable payback. Additional references are available on request.")


def add_competitive_advantage_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, RGBColor(255, 255, 255))
    add_header(slide, "Competitors vs. Ingecart. The Difference is Clear.")
    rows = [
        ("Vertical OEMs", "Sell their own line, locked ecosystem", "Evaluate the best option, even if not new"),
        ("Generalist Integrators", "Learn your industry as they go", "Have 28 years of corrugated expertise"),
        ("Point Solutions", "Optimize one part, ignore the system", "Optimize the complete system for maximum ROI"),
    ]
    table = slide.shapes.add_table(len(rows) + 1, 3, Inches(0.75), Inches(1.8), Inches(8.6), Inches(3.95)).table
    table.columns[0].width = Inches(2.0)
    table.columns[1].width = Inches(3.2)
    table.columns[2].width = Inches(3.4)
    headers = ["Against...", "They...", "We..."]
    for col_index, header_text in enumerate(headers):
        cell = table.cell(0, col_index)
        cell.text = header_text
        cell.fill.solid()
        cell.fill.fore_color.rgb = INGECART_RED if col_index == 2 else INGECART_DARK
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.alignment = PP_ALIGN.CENTER
        for run in paragraph.runs:
            run.font.name = TITLE_FONT
            run.font.size = Pt(12)
            run.font.bold = True
            run.font.color.rgb = RGBColor(255, 255, 255)
    for row_index, (against, they, we) in enumerate(rows, start=1):
        values = [against, they, we]
        for col_index, value in enumerate(values):
            cell = table.cell(row_index, col_index)
            cell.text = value
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(255, 255, 255) if row_index % 2 else INGECART_LIGHT
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.alignment = PP_ALIGN.LEFT if col_index == 0 else PP_ALIGN.CENTER
            for run in paragraph.runs:
                run.font.name = BODY_FONT
                run.font.size = Pt(11)
                run.font.color.rgb = INGECART_DARK if col_index == 2 else INGECART_GRAY

    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.45), Inches(5.95), Inches(4.3), Inches(0.38), fill=INGECART_RED, line=INGECART_RED)
    add_textbox(slide, Inches(2.7), Inches(6.02), Inches(3.8), Inches(0.2), "Independent. Specialized. System-focused.", font_name=TITLE_FONT, font_size=12, color=RGBColor(255, 255, 255), bold=True, align=PP_ALIGN.CENTER)
    add_note(slide, "Ingecart wins because it is independent, specialized, and system-focused. We do not push a vendor line. We do not learn your industry on your time. We solve the complete operational system.")


def add_timeline_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, RGBColor(255, 255, 255))
    add_header(slide, "A Clear, Phased Process to Minimize Your Risk")

    steps = [
        ("1", "Audit", "2 days onsite"),
        ("2", "Modeling", "Simulation"),
        ("3", "Design", "Best-of-breed"),
        ("4", "Approval", "Clear ROI"),
        ("5", "Implementation", "Turnkey"),
        ("6", "Optimization", "Continuous improvement"),
    ]
    start_x = 0.7
    gap = 1.48
    y_circle = 2.45
    for index, (number, title, subtitle) in enumerate(steps):
        x = Inches(start_x + index * gap)
        if index < len(steps) - 1:
            add_shape(slide, MSO_SHAPE.RECTANGLE, x + Inches(0.42), Inches(2.67), Inches(1.06), Inches(0.05), fill=INGECART_RED, line=INGECART_RED)
        add_shape(slide, MSO_SHAPE.OVAL, x, Inches(y_circle), Inches(0.9), Inches(0.9), fill=INGECART_RED if index in (0, 3, 5) else INGECART_LIGHT, line=INGECART_BORDER)
        add_textbox(slide, x, Inches(y_circle + 0.1), Inches(0.9), Inches(0.25), number, font_name=TITLE_FONT, font_size=18, color=RGBColor(255, 255, 255) if index in (0, 3, 5) else INGECART_DARK, bold=True, align=PP_ALIGN.CENTER)
        add_textbox(slide, x - Inches(0.25), Inches(3.5), Inches(1.4), Inches(0.25), title, font_name=TITLE_FONT, font_size=12, color=INGECART_DARK, bold=True, align=PP_ALIGN.CENTER)
        add_textbox(slide, x - Inches(0.35), Inches(3.75), Inches(1.6), Inches(0.35), subtitle, font_name=BODY_FONT, font_size=10, color=INGECART_GRAY, align=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(2.0), Inches(6.05), Inches(5.8), Inches(0.25), "You control the pace and the investment.", font_name=BODY_FONT, font_size=11, color=INGECART_GRAY, italic=True, align=PP_ALIGN.CENTER)
    add_note(slide, "Our process reduces risk. Audit first, then model, then design, then approve, then implement, and finally optimize continuously. The customer controls each phase and each investment decision.")


def add_cta_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, RGBColor(255, 255, 255))
    add_header(slide, "The Time to Act is Now")

    options = [
        ("Preliminary Diagnosis", "30-minute call to understand your main pain point", "Low"),
        ("Quick Audit", "2 days at your plant to identify the real bottleneck", "Low to Medium"),
        ("Custom Roadmap", "We deliver the roadmap and potential ROI with no obligation", "No commitment"),
    ]
    xs = [Inches(0.75), Inches(3.55), Inches(6.35)]
    fills = [INGECART_LIGHT, INGECART_SOFT_BLUE, INGECART_SOFT_RED]
    widths = [Inches(2.55), Inches(2.55), Inches(2.55)]
    for x, (title, body, effort), fill, width in zip(xs, options, fills, widths):
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.95), width, Inches(2.65), fill=fill, line=INGECART_BORDER)
        add_shape(slide, MSO_SHAPE.OVAL, x + Inches(0.9), Inches(2.1), Inches(0.75), Inches(0.75), fill=INGECART_RED, line=INGECART_RED)
        add_textbox(slide, x + Inches(0.25), Inches(2.95), Inches(2.05), Inches(0.45), title, font_name=TITLE_FONT, font_size=14, color=INGECART_DARK, bold=True, align=PP_ALIGN.CENTER)
        add_textbox(slide, x + Inches(0.2), Inches(3.4), Inches(2.1), Inches(0.7), body, font_name=BODY_FONT, font_size=11, color=INGECART_GRAY, align=PP_ALIGN.CENTER)
        add_textbox(slide, x + Inches(0.3), Inches(4.15), Inches(1.95), Inches(0.25), f"Effort: {effort}", font_name=BODY_FONT, font_size=10, color=INGECART_RED, bold=True, align=PP_ALIGN.CENTER)

    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.55), Inches(5.35), Inches(4.65), Inches(0.6), fill=INGECART_RED, line=INGECART_RED)
    add_textbox(slide, Inches(2.8), Inches(5.5), Inches(4.15), Inches(0.2), "Request your Opportunity Audit", font_name=TITLE_FONT, font_size=15, color=RGBColor(255, 255, 255), bold=True, align=PP_ALIGN.CENTER)
    add_note(slide, "There are three ways to start. A 30-minute call, a two-day audit, or a custom roadmap with ROI analysis. The only wrong move is not starting. Your competitors are already optimizing.")


def add_contact_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, RGBColor(255, 255, 255))
    add_header(slide, "Let's Talk with Data")

    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(1.8), Inches(4.4), Inches(4.3), fill=INGECART_LIGHT, line=INGECART_BORDER)
    add_textbox(slide, Inches(1.0), Inches(2.1), Inches(1.3), Inches(0.25), "Name", font_name=TITLE_FONT, font_size=13, color=INGECART_DARK, bold=True)
    add_textbox(slide, Inches(2.05), Inches(2.1), Inches(2.7), Inches(0.25), "[Insert Ingecart Commercial Contact Name]", font_name=BODY_FONT, font_size=12, color=INGECART_GRAY)
    add_textbox(slide, Inches(1.0), Inches(2.55), Inches(1.3), Inches(0.25), "Title", font_name=TITLE_FONT, font_size=13, color=INGECART_DARK, bold=True)
    add_textbox(slide, Inches(2.05), Inches(2.55), Inches(2.7), Inches(0.25), "[Insert Title]", font_name=BODY_FONT, font_size=12, color=INGECART_GRAY)
    add_textbox(slide, Inches(1.0), Inches(3.0), Inches(1.3), Inches(0.25), "Email", font_name=TITLE_FONT, font_size=13, color=INGECART_DARK, bold=True)
    add_textbox(slide, Inches(2.05), Inches(3.0), Inches(2.7), Inches(0.25), "[Insert Email]", font_name=BODY_FONT, font_size=12, color=INGECART_GRAY)
    add_textbox(slide, Inches(1.0), Inches(3.45), Inches(1.3), Inches(0.25), "Phone", font_name=TITLE_FONT, font_size=13, color=INGECART_DARK, bold=True)
    add_textbox(slide, Inches(2.05), Inches(3.45), Inches(2.7), Inches(0.25), "[Insert Direct Phone]", font_name=BODY_FONT, font_size=12, color=INGECART_GRAY)
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(4.0), Inches(3.85), Inches(1.55), fill=RGBColor(250, 250, 250), line=INGECART_BORDER)
    add_textbox(slide, Inches(1.25), Inches(4.47), Inches(3.35), Inches(0.25), "[QR CODE PLACEHOLDER: Link to Ingecart Contact Form]", font_name=BODY_FONT, font_size=11, color=INGECART_GRAY, italic=True, align=PP_ALIGN.CENTER)

    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.55), Inches(1.8), Inches(3.65), Inches(4.3), fill=INGECART_DARK, line=INGECART_DARK)
    add_textbox(slide, Inches(5.9), Inches(2.65), Inches(2.95), Inches(0.7), '"Don\'t adapt your factory. We adapt to you."', font_name=TITLE_FONT, font_size=22, color=RGBColor(255, 255, 255), bold=True, italic=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(6.25), Inches(4.35), Inches(2.2), Inches(0.3), "Ingecart", font_name=TITLE_FONT, font_size=18, color=INGECART_RED, bold=True, align=PP_ALIGN.CENTER)
    add_note(slide, "No empty presentations. No vague promises. Here is exactly who to contact. Scan the QR code, send an email, or make a call. We will respond within 24 hours with a diagnostic proposal tailored to the plant.")


def add_thank_you_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, RGBColor(250, 250, 250))
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(5.95), Inches(8.6), Inches(0.06), fill=INGECART_RED, line=INGECART_RED)
    add_textbox(slide, Inches(2.2), Inches(1.8), Inches(4.5), Inches(0.8), "Thank You", font_name=TITLE_FONT, font_size=28, color=INGECART_DARK, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1.65), Inches(2.7), Inches(5.7), Inches(0.5), "Knowledge is the foundation for growth", font_name=TITLE_FONT, font_size=20, color=INGECART_RED, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1.65), Inches(3.45), Inches(5.7), Inches(0.4), "Whitepapers | Case Studies | ROI Calculator", font_name=BODY_FONT, font_size=12, color=INGECART_GRAY, align=PP_ALIGN.CENTER)
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.0), Inches(4.25), Inches(2.0), Inches(0.65), fill=INGECART_DARK, line=INGECART_DARK)
    add_textbox(slide, Inches(3.15), Inches(4.4), Inches(1.7), Inches(0.2), "INGECART", font_name=TITLE_FONT, font_size=18, color=RGBColor(255, 255, 255), bold=True, align=PP_ALIGN.CENTER)
    add_note(slide, "Thank you for your time. We have shown the problem, the solution, the benefits, and the proof. We are ready when you are.")


def build_presentation() -> Presentation:
    prs = Presentation(str(TEMPLATE_PATH))
    clear_presentation(prs)

    add_title_slide(prs)
    add_venn_slide(prs)
    add_footprint_slide(prs)
    add_dual_strategy_slide(prs)
    add_pain_slide(prs)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, RGBColor(255, 255, 255))
    add_comparison_table(
        slide,
        "Don't Adapt Your Factory to a Standard Machine",
        [
            ("I sell the machine I have", "We design the solution you need"),
            ("One-size-fits-all solution", "Tailored to your unique operational context"),
            ("Locked into one brand or technology", "Best-of-breed, multi-vendor"),
            ("Product-focused (hardware)", "Outcome-focused (ROI & KPIs)"),
        ],
        (Inches(4.15), Inches(4.35)),
        "Traditional Approach (OEM or Generalist)",
        "Ingecart's Approach",
    )
    add_note(slide, "The traditional approach forces a plant to adapt to a machine. Ingecart reverses the logic: we adapt technology to the customer\'s operational reality and financial goals.")

    add_value_prop_slide(prs)
    add_financial_impact_slide(prs)
    add_operational_impact_slide(prs)
    add_evidence_slide(prs)
    add_competitive_advantage_slide(prs)
    add_timeline_slide(prs)
    add_cta_slide(prs)
    add_contact_slide(prs)
    add_thank_you_slide(prs)

    return prs


if __name__ == "__main__":
    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    presentation = build_presentation()
    presentation.save(str(OUTPUT_PATH))
    print(f"Presentation saved to {OUTPUT_PATH}")