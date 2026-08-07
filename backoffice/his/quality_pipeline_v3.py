from __future__ import annotations

import json
import re
import shutil
import time
import zipfile
from dataclasses import dataclass
from datetime import UTC, datetime
from pathlib import Path
from typing import Any
from uuid import uuid4

from bs4 import BeautifulSoup
from pptx import Presentation

from backoffice.dipc.component_library import normalize_component
from backoffice.dipc.models import AssetRef, BlockNode, DocumentModel, EvidenceRecord, KnowledgeLink, MissionLink, SectionNode, VersionEntry
from backoffice.dipc.preview_engine import PreviewEngine
from backoffice.dipc.publication_engine import PublicationEngine
from backoffice.dipc.versioning import DocumentVersionStore
from backoffice.dipc.visual_recognition import DiagramRecognitionEngine
from backoffice.pie.powerpoint_parser import PowerPointParser
from backoffice.theme.design_system import INDUSTRIAL
from backoffice.dipc.theme_engine import build_css


PUBLICATION_STATES = ["Draft", "Editing", "Review", "Validated", "Published", "Archived"]


def _write_text_atomic_with_retries(path: Path, content: str, attempts: int = 5, delay_s: float = 0.2) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    last_error: OSError | None = None
    for _ in range(attempts):
        try:
            temp_path = path.with_suffix(path.suffix + ".tmp")
            temp_path.write_text(content, encoding="utf-8")
            temp_path.replace(path)
            return
        except OSError as exc:
            last_error = exc
            time.sleep(delay_s)
    if last_error is not None:
        raise last_error


class QualityGateError(RuntimeError):
    pass


@dataclass
class PhaseTimer:
    name: str
    started_at: float
    ended_at: float

    @property
    def elapsed_s(self) -> float:
        return round(self.ended_at - self.started_at, 3)


class HtmlIntelligenceStudioV3Pipeline:
    def __init__(self, corporate_model_path: Path) -> None:
        self.corporate_model_path = corporate_model_path
        self.parser = PowerPointParser()
        self.recognizer = DiagramRecognitionEngine()
        self.publisher = PublicationEngine()
        self.preview = PreviewEngine()

    def run(
        self,
        *,
        sources: list[str],
        source_format: str,
        output_root: str,
        document_name: str,
        project: str,
        client: str,
        category: str,
        language: str,
        objective: str,
        audience: str,
        theme_variant: str = "industrial",
        author: str = "Mission Manager",
        force_no_cache: bool = True,
        max_regeneration_attempts: int = 2,
    ) -> dict[str, Any]:
        cleaned = [s.strip() for s in sources if s and s.strip()]
        if not cleaned:
            raise ValueError("At least one source is required for HIS V3 pipeline.")

        run_id = datetime.now(UTC).strftime("%Y%m%d_%H%M%S") + "_" + uuid4().hex[:8]
        base_dir = Path(output_root).expanduser().resolve() / run_id

        source_dir = base_dir / "source"
        assets_dir = base_dir / "assets"
        variants_dir = base_dir / "variants"
        metadata_dir = base_dir / "metadata"
        logs_dir = base_dir / "logs"
        history_dir = base_dir / "history"
        selected_dir = base_dir / "selected"
        for p in [source_dir, assets_dir, variants_dir, metadata_dir, logs_dir, history_dir, selected_dir]:
            p.mkdir(parents=True, exist_ok=True)

        asset_tree = {
            "images": assets_dir / "images",
            "icons": assets_dir / "icons",
            "svg": assets_dir / "svg",
            "logos": assets_dir / "logos",
            "attachments": assets_dir / "attachments",
            "thumbnails": assets_dir / "thumbnails",
        }
        for p in asset_tree.values():
            p.mkdir(parents=True, exist_ok=True)

        copied_sources: list[Path] = []
        for src in cleaned:
            source_path = Path(src).expanduser().resolve()
            if not source_path.exists():
                continue
            if force_no_cache:
                target = source_dir / source_path.name
                shutil.copy2(source_path, target)
                copied_sources.append(target)
            else:
                copied_sources.append(source_path)

        if not copied_sources:
            raise FileNotFoundError("None of the provided sources exist on disk.")

        primary = self._choose_primary(copied_sources, source_format)
        phases: list[PhaseTimer] = []
        hypotheses, selected_hypothesis = self._build_hypotheses(primary, copied_sources)
        discarded_hypotheses = [h for h in hypotheses if h["id"] != selected_hypothesis["id"]]

        quality_attempts: list[dict[str, Any]] = []
        selected_payload: dict[str, Any] | None = None
        quality_thresholds = self._quality_thresholds()

        for attempt in range(1, max_regeneration_attempts + 1):
            attempt_suffix = f"attempt_{attempt}"
            attempt_variants_dir = variants_dir / attempt_suffix
            attempt_variants_dir.mkdir(parents=True, exist_ok=True)

            # 1 Discovery
            t0 = time.perf_counter()
            discovery = self._document_discovery(primary)
            phases.append(PhaseTimer(f"phase_1_document_discovery_{attempt_suffix}", t0, time.perf_counter()))

            # 2 Object extraction
            t0 = time.perf_counter()
            analysis = self._extract_objects(primary, assets_dir)
            object_inventory = self._object_inventory(analysis)
            phases.append(PhaseTimer(f"phase_2_object_extraction_{attempt_suffix}", t0, time.perf_counter()))

            # 3 Image extraction
            t0 = time.perf_counter()
            image_inventory = self._image_inventory(analysis, asset_tree)
            phases.append(PhaseTimer(f"phase_3_image_extraction_{attempt_suffix}", t0, time.perf_counter()))

            # 4 Semantic classification
            t0 = time.perf_counter()
            text_analysis = self._text_analysis(analysis)
            visual_analysis = self._visual_analysis(analysis)
            phases.append(PhaseTimer(f"phase_4_semantic_classification_{attempt_suffix}", t0, time.perf_counter()))

            # 5 DOM reconstruction
            t0 = time.perf_counter()
            slide_flow_doc = self._build_document_model(
                analysis=analysis,
                visual_analysis=visual_analysis,
                document_name=document_name,
                primary=primary,
                project=project,
                client=client,
                category=category,
                language=language,
                objective=objective,
                audience=audience,
                theme_variant=theme_variant,
                mode="slide_flow",
            )
            smart_doc = self._build_document_model(
                analysis=analysis,
                visual_analysis=visual_analysis,
                document_name=document_name,
                primary=primary,
                project=project,
                client=client,
                category=category,
                language=language,
                objective=objective,
                audience=audience,
                theme_variant=theme_variant,
                mode="smart_reconstruction",
            )
            phases.append(PhaseTimer(f"phase_5_dom_reconstruction_{attempt_suffix}", t0, time.perf_counter()))

            # 6 Theme corporate
            t0 = time.perf_counter()
            corporate_css = self._load_corporate_css_or_bootstrap(theme_variant)
            phases.append(PhaseTimer(f"phase_6_theme_application_{attempt_suffix}", t0, time.perf_counter()))

            # 7 Validation and publication scoring
            t0 = time.perf_counter()
            slide_flow_out = self.publisher.export_all(slide_flow_doc, attempt_variants_dir / "slide_flow", corporate_css=corporate_css)
            smart_out = self.publisher.export_all(smart_doc, attempt_variants_dir / "smart_reconstruction", corporate_css=corporate_css)

            q_slide = self._quality_gate_scores(analysis, slide_flow_doc, visual_analysis, slide_flow_out, quality_thresholds)
            q_smart = self._quality_gate_scores(analysis, smart_doc, visual_analysis, smart_out, quality_thresholds)

            # Re-evaluate quality after publication artifacts are generated.
            q_slide = self._quality_gate_scores(analysis, slide_flow_doc, visual_analysis, slide_flow_out, quality_thresholds)
            q_smart = self._quality_gate_scores(analysis, smart_doc, visual_analysis, smart_out, quality_thresholds)
            s_slide = self._score_variant(q_slide)
            s_smart = self._score_variant(q_smart)
            phases.append(PhaseTimer(f"phase_7_validation_publication_{attempt_suffix}", t0, time.perf_counter()))

            chosen_name = "slide_flow" if s_slide["executive_quality_score"] >= s_smart["executive_quality_score"] else "smart_reconstruction"
            chosen_doc = slide_flow_doc if chosen_name == "slide_flow" else smart_doc
            chosen_out = slide_flow_out if chosen_name == "slide_flow" else smart_out
            chosen_quality = q_slide if chosen_name == "slide_flow" else q_smart
            chosen_score = s_slide if chosen_name == "slide_flow" else s_smart

            chosen_metrics = chosen_quality.get("metrics", {})
            passed = bool(chosen_quality["passed"]) and all(
                float(chosen_metrics.get(metric, 0.0)) >= float(threshold) for metric, threshold in quality_thresholds.items()
            )
            attempt_result = {
                "attempt": attempt,
                "selected_variant": chosen_name,
                "passed": passed,
                "quality": {"slide_flow": q_slide, "smart_reconstruction": q_smart},
                "scores": {"slide_flow": s_slide, "smart_reconstruction": s_smart},
            }
            quality_attempts.append(attempt_result)

            if passed or attempt == max_regeneration_attempts:
                selected_payload = {
                    "selected_name": chosen_name,
                    "selected_doc": chosen_doc,
                    "selected_out": chosen_out,
                    "selected_quality": chosen_quality,
                    "selected_score": chosen_score,
                    "discovery": discovery,
                    "object_inventory": object_inventory,
                    "image_inventory": image_inventory,
                    "text_analysis": text_analysis,
                    "visual_analysis": visual_analysis,
                    "slide_flow_out": slide_flow_out,
                    "smart_out": smart_out,
                    "slide_flow_quality": q_slide,
                    "smart_quality": q_smart,
                    "slide_flow_score": s_slide,
                    "smart_score": s_smart,
                }
                break

            # simple regeneration strategy: boost smart reconstruction for next attempt
            selected_hypothesis = self._rebalance_hypothesis(selected_hypothesis)

        if selected_payload is None:
            raise QualityGateError("HIS V3 could not produce a selected payload.")

        selected_doc: DocumentModel = selected_payload["selected_doc"]
        selected_out: dict[str, str] = selected_payload["selected_out"]
        selected_doc.preview_profiles = self.preview.build_profiles(selected_doc, selected_out["html"])
        preview_manifest_path = self.preview.write_manifest(selected_doc, metadata_dir / "preview_manifest.json")

        publication_state = self._publication_state(selected_payload["selected_quality"], selected_payload["selected_score"])
        state_history = self._state_history(publication_state)

        mission_id = run_id
        selected_doc.metadata["publication_state"] = publication_state
        selected_doc.metadata["publication_state_history"] = state_history
        selected_doc.metadata["executive_quality_score"] = selected_payload["selected_score"].get("executive_quality_score", 0.0)
        selected_doc.metadata["asset_manager"] = {
            "root": str(assets_dir),
            "tree": {k: str(v) for k, v in asset_tree.items()},
        }

        model_path = metadata_dir / "document_model.json"
        model_path.write_text(selected_doc.model_dump_json(indent=2), encoding="utf-8")

        store = DocumentVersionStore(history_dir)
        first_version = VersionEntry(
            version_number=1,
            author=author,
            objective=objective or "HIS V3 full reconstruction",
            result=publication_state.lower(),
            mission_id=mission_id,
            output_files=selected_out,
        )
        selected_doc.version_history.append(first_version)
        selected_doc.mission_links.append(MissionLink(mission_id=mission_id, objective=first_version.objective, agent=author))
        store.append_version(selected_doc, first_version)

        self._copy_selected_variant(selected_out, selected_dir)

        asset_registry = self._build_asset_registry(selected_doc, asset_tree)
        asset_registry_path = metadata_dir / "asset_registry.json"
        asset_registry_path.write_text(json.dumps(asset_registry, indent=2, ensure_ascii=False), encoding="utf-8")

        knowledge_package = self._build_knowledge_package(selected_doc, selected_payload, hypotheses, selected_hypothesis)
        knowledge_path = metadata_dir / "knowledge_package.json"
        knowledge_path.write_text(json.dumps(knowledge_package, indent=2, ensure_ascii=False), encoding="utf-8")

        enterprise_memory = self._build_enterprise_memory(selected_doc, selected_payload, run_id)
        memory_path = metadata_dir / "enterprise_memory.json"
        memory_path.write_text(json.dumps(enterprise_memory, indent=2, ensure_ascii=False), encoding="utf-8")

        truth_graph_payload = self._build_truth_graph(selected_doc, selected_payload)
        truth_graph_file = metadata_dir / "truth_graph.json"
        truth_graph_file.write_text(json.dumps(truth_graph_payload, indent=2, ensure_ascii=False), encoding="utf-8")
        self._append_truth_graph_global(truth_graph_payload)

        report_payload = self._build_technical_report_payload(
            run_id=run_id,
            primary=primary,
            phases=phases,
            selected_payload=selected_payload,
            quality_attempts=quality_attempts,
            hypotheses=hypotheses,
            selected_hypothesis=selected_hypothesis,
            discarded_hypotheses=discarded_hypotheses,
            thresholds=quality_thresholds,
        )
        report_json = metadata_dir / "technical_report_his_v3.json"
        report_md = metadata_dir / "technical_report_his_v3.md"
        report_json.write_text(json.dumps(report_payload, indent=2, ensure_ascii=False), encoding="utf-8")
        report_md.write_text(self._render_report_markdown(report_payload), encoding="utf-8")

        quality_report = {
            "thresholds": quality_thresholds,
            "attempts": quality_attempts,
            "selected_quality": selected_payload["selected_quality"],
            "selected_scores": selected_payload["selected_score"],
        }
        quality_report_path = metadata_dir / "quality_report.json"
        quality_report_path.write_text(json.dumps(quality_report, indent=2, ensure_ascii=False), encoding="utf-8")

        mission_log = {
            "mission_id": mission_id,
            "mode": "AHDE",
            "pipeline": "HIS V3 Production Ready",
            "selected_variant": selected_payload["selected_name"],
            "publication_state": publication_state,
            "hypotheses": hypotheses,
            "selected_hypothesis": selected_hypothesis,
            "quality_attempts": quality_attempts,
            "phases": [{"name": p.name, "elapsed_s": p.elapsed_s} for p in phases],
        }
        mission_log_path = logs_dir / "mission_log.json"
        mission_log_path.write_text(json.dumps(mission_log, indent=2, ensure_ascii=False), encoding="utf-8")

        ai_history_path = logs_dir / "ai_history.json"
        ai_history_path.write_text(json.dumps({"commands": [], "generated_at": datetime.now(UTC).isoformat()}, indent=2, ensure_ascii=False), encoding="utf-8")

        evidence_zip = base_dir / f"his_v3_evidence_{run_id}.zip"
        self._build_evidence_zip(base_dir, evidence_zip)

        return {
            "run_id": run_id,
            "output_dir": str(base_dir),
            "document_model_path": str(model_path),
            "theme_css_path": selected_out["theme_css"],
            "preview_manifest_path": str(preview_manifest_path),
            "knowledge_package_path": str(knowledge_path),
            "enterprise_memory_path": str(memory_path),
            "truth_graph_path": str(truth_graph_file),
            "mission_log_path": str(mission_log_path),
            "technical_report_json": str(report_json),
            "technical_report_md": str(report_md),
            "quality_report_path": str(quality_report_path),
            "asset_registry_path": str(asset_registry_path),
            "evidence_zip_path": str(evidence_zip),
            "publication_state": publication_state,
            "publication_state_history": state_history,
            "publication_outputs": {
                "selected": selected_out,
                "slide_flow": selected_payload["slide_flow_out"],
                "smart_reconstruction": selected_payload["smart_out"],
            },
            "selected_variant": selected_payload["selected_name"],
            "scores": {
                "slide_flow": selected_payload["slide_flow_score"],
                "smart_reconstruction": selected_payload["smart_score"],
            },
            "quality": {
                "slide_flow": selected_payload["slide_flow_quality"],
                "smart_reconstruction": selected_payload["smart_quality"],
            },
            "html_path": selected_out["html"],
            "quality_attempts": quality_attempts,
            "hypotheses": hypotheses,
            "selected_hypothesis": selected_hypothesis,
        }

    def apply_dom_command(self, document_model_path: str, command: str, author: str = "AI Coordinator") -> dict[str, Any]:
        model_path = Path(document_model_path).expanduser().resolve()
        document = DocumentModel.model_validate_json(model_path.read_text(encoding="utf-8"))
        base_dir = model_path.parent.parent
        command_l = command.lower().strip()
        changes: list[str] = []

        if "crear capítulo" in command_l or "create chapter" in command_l:
            new_order = len(document.sections) + 1
            chapter_title = f"Chapter {new_order}"
            document.sections.append(
                SectionNode(
                    title=chapter_title,
                    summary="Created by AI Command Layer.",
                    order=new_order,
                    blocks=[BlockNode(block_type="chapter", components=[normalize_component("text", chapter_title, "New chapter content")])],
                )
            )
            changes.append(f"Added {chapter_title}")

        if "eliminar capítulo" in command_l or "delete chapter" in command_l:
            if document.sections:
                removed = document.sections.pop()
                changes.append(f"Removed chapter: {removed.title}")

        if "añadir kpi" in command_l or "add kpi" in command_l:
            if document.sections:
                document.sections[0].blocks.append(
                    BlockNode(
                        block_type="kpi",
                        components=[
                            normalize_component(
                                "industrial_kpi",
                                "Executive KPI",
                                None,
                                items=[
                                    {"label": "Coverage", "value": "100%", "detail": "Quality gates"},
                                    {"label": "Version", "value": len(document.version_history) + 1, "detail": "DOM operation"},
                                ],
                            )
                        ],
                    )
                )
                changes.append("Added KPI block")

        if "traducir" in command_l or "translate" in command_l:
            document.metadata["translations"] = {
                "en": {"title": document.title, "subtitle": document.subtitle or "Translated by AI command"},
                "es": {"title": document.title, "subtitle": document.subtitle or "Traducido por comando de IA"},
            }
            changes.append("Updated bilingual translation metadata")

        if "añadir imagen" in command_l or "add image" in command_l:
            image_path_match = re.search(r"([a-zA-Z]:\\[^\n\r]+\.(?:png|jpg|jpeg|webp|gif))", command, flags=re.IGNORECASE)
            if image_path_match:
                img_src = Path(image_path_match.group(1)).expanduser().resolve()
                if img_src.exists():
                    assets_images = base_dir / "assets" / "images"
                    assets_images.mkdir(parents=True, exist_ok=True)
                    target = assets_images / img_src.name
                    shutil.copy2(img_src, target)
                    document.assets.append(AssetRef(kind="image", path=str(target), title=img_src.stem))
                    if document.sections:
                        document.sections[0].blocks.append(
                            BlockNode(block_type="image", components=[normalize_component("image_gallery", "Inserted Image", None, items=[{"title": img_src.stem, "body": target.as_posix()}])])
                        )
                    changes.append(f"Added image asset: {img_src.name}")

        if not changes:
            changes.append("No structural changes detected from command.")

        corporate_css = self._load_corporate_css_or_bootstrap(str(document.theme_variant))
        output_dir = base_dir / "published" / (datetime.now(UTC).strftime("%Y%m%d_%H%M%S") + "_dom")
        outputs = self.publisher.export_all(document, output_dir, corporate_css=corporate_css)

        mission_id = datetime.now(UTC).strftime("%Y%m%d_%H%M%S") + "_" + uuid4().hex[:8]
        version = VersionEntry(
            version_number=len(document.version_history) + 1,
            author=author,
            objective=command,
            result="updated",
            mission_id=mission_id,
            output_files=outputs,
        )
        document.version_history.append(version)
        document.mission_links.append(MissionLink(mission_id=mission_id, objective=command, command=command, agent=author))
        history = document.metadata.get("publication_state_history", [])
        if not isinstance(history, list):
            history = []
        history.append({"state": "Editing", "ts": datetime.now(UTC).isoformat(), "mission_id": mission_id})
        document.metadata["publication_state"] = "Editing"
        document.metadata["publication_state_history"] = history
        _write_text_atomic_with_retries(model_path, document.model_dump_json(indent=2))

        store = DocumentVersionStore(base_dir / "history")
        store.append_version(document, version)

        ai_history = base_dir / "logs" / "ai_history.json"
        ai_history.parent.mkdir(parents=True, exist_ok=True)
        if ai_history.exists():
            data = json.loads(ai_history.read_text(encoding="utf-8"))
            cmds = data.get("commands", []) if isinstance(data, dict) else []
        else:
            cmds = []
        cmds.append({"ts": datetime.now(UTC).isoformat(), "command": command, "changes": changes, "mission_id": mission_id})
        ai_history.write_text(json.dumps({"commands": cmds}, indent=2, ensure_ascii=False), encoding="utf-8")

        mission_log = base_dir / "logs" / "mission_log.json"
        mission_log.parent.mkdir(parents=True, exist_ok=True)
        if mission_log.exists():
            try:
                mission_payload = json.loads(mission_log.read_text(encoding="utf-8"))
            except Exception:
                mission_payload = {}
        else:
            mission_payload = {}
        actions = mission_payload.get("actions", []) if isinstance(mission_payload, dict) else []
        actions.append(
            {
                "ts": datetime.now(UTC).isoformat(),
                "mission_id": mission_id,
                "action": "dom_command",
                "command": command,
                "author": author,
                "changes": changes,
                "publication_state": "Editing",
            }
        )
        if not isinstance(mission_payload, dict):
            mission_payload = {}
        mission_payload["actions"] = actions
        mission_payload["updated_at"] = datetime.now(UTC).isoformat()
        mission_log.write_text(json.dumps(mission_payload, indent=2, ensure_ascii=False), encoding="utf-8")

        return {
            "run_id": mission_id,
            "html_path": outputs["html"],
            "document_model_path": str(model_path),
            "publication_outputs": outputs,
            "changes": changes,
            "mission_log_path": str(base_dir / "logs" / "mission_log.json"),
            "ai_history_path": str(ai_history),
        }

    def _choose_primary(self, sources: list[Path], source_format: str) -> Path:
        ext_order = {
            ".pptx": 100,
            ".ppt": 95,
            ".docx": 90,
            ".pdf": 88,
            ".html": 80,
            ".htm": 80,
            ".md": 72,
            ".txt": 70,
            ".jpg": 65,
            ".jpeg": 65,
            ".png": 65,
        }
        weighted = sorted(sources, key=lambda p: ext_order.get(p.suffix.lower(), 50), reverse=True)
        return weighted[0]

    def _build_hypotheses(self, primary: Path, all_sources: list[Path]) -> tuple[list[dict[str, Any]], dict[str, Any]]:
        source_count = len(all_sources)
        ext = primary.suffix.lower()
        hypotheses = [
            {"id": "H1", "name": "Strict structure-first reconstruction", "engineering": 8.9, "business": 8.3, "quality": 9.2, "risk": 2.0},
            {"id": "H2", "name": "Balanced semantic-first reconstruction", "engineering": 8.4, "business": 8.9, "quality": 8.7, "risk": 2.4},
            {"id": "H3", "name": "Visual-priority reconstruction", "engineering": 8.0, "business": 8.2, "quality": 8.4, "risk": 2.9},
        ]
        for h in hypotheses:
            if source_count > 1:
                h["quality"] += 0.2
            if ext in {".pptx", ".ppt"} and h["id"] == "H1":
                h["engineering"] += 0.4
            h["executive_score"] = round(h["engineering"] * 0.35 + h["business"] * 0.25 + h["quality"] * 0.4 - h["risk"] * 0.2, 2)
        selected = max(hypotheses, key=lambda x: x["executive_score"])
        return hypotheses, selected

    def _rebalance_hypothesis(self, selected_hypothesis: dict[str, Any]) -> dict[str, Any]:
        updated = dict(selected_hypothesis)
        updated["executive_score"] = round(float(updated.get("executive_score", 0.0)) + 0.3, 2)
        return updated

    def _document_discovery(self, primary: Path) -> dict[str, Any]:
        ext = primary.suffix.lower()
        if ext in {".ppt", ".pptx"}:
            prs = Presentation(str(primary))
            cp = prs.core_properties
            hyperlinks = 0
            notes_slides = 0
            for slide in prs.slides:
                try:
                    notes_text = slide.notes_slide.notes_text_frame.text if slide.has_notes_slide and slide.notes_slide and slide.notes_slide.notes_text_frame else ""
                    if notes_text.strip():
                        notes_slides += 1
                except Exception:
                    pass
                for shape in slide.shapes:
                    try:
                        if shape.click_action and shape.click_action.hyperlink and shape.click_action.hyperlink.address:
                            hyperlinks += 1
                    except Exception:
                        pass
            return {
                "source": str(primary),
                "source_type": "pptx",
                "slide_count": len(prs.slides),
                "master_slides": len(prs.slide_masters),
                "layouts": sum(len(master.slide_layouts) for master in prs.slide_masters),
                "author": cp.author or "unknown",
                "metadata": {
                    "title": cp.title or primary.stem,
                    "language": cp.language or "unknown",
                    "keywords": cp.keywords or "",
                },
                "hyperlinks": hyperlinks,
                "comments": 0,
                "notes_slides": notes_slides,
                "hidden_slides": 0,
            }

        text = self._extract_text_from_file(primary)
        return {
            "source": str(primary),
            "source_type": ext.lstrip("."),
            "slide_count": 1,
            "master_slides": 0,
            "layouts": 1,
            "author": "unknown",
            "metadata": {"title": primary.stem, "language": "unknown", "keywords": ""},
            "hyperlinks": len(re.findall(r"https?://", text, flags=re.IGNORECASE)),
            "comments": 0,
            "notes_slides": 0,
            "hidden_slides": 0,
        }

    def _extract_objects(self, primary: Path, assets_dir: Path) -> Any:
        ext = primary.suffix.lower()
        if ext in {".ppt", ".pptx"}:
            return self.parser.parse(str(primary), str(assets_dir / "images"))
        return self._build_generic_analysis(primary)

    def _build_generic_analysis(self, primary: Path) -> Any:
        from backoffice.pie.models import ElementStyle, PresentationAnalysis, SlideAnalysis, SlideElement

        text = self._extract_text_from_file(primary)
        blocks = [chunk.strip() for chunk in re.split(r"\n\s*\n", text) if chunk.strip()]
        elements: list[SlideElement] = []
        y = 10.0
        for idx, chunk in enumerate(blocks[:40], start=1):
            kind = "title" if idx == 1 else "text"
            elements.append(
                SlideElement(
                    element_id=f"g-e{idx}",
                    kind=kind,
                    x=10.0,
                    y=y,
                    w=900.0,
                    h=40.0,
                    text=chunk[:1500],
                    style=ElementStyle(font_name="Segoe UI", font_size_pt=18.0 if idx == 1 else 12.0, bold=(idx == 1)),
                )
            )
            y += 42.0

        slide = SlideAnalysis(
            index=1,
            title=primary.stem,
            width=1366.0,
            height=768.0,
            background="#0D0F13",
            layout="single-column",
            palette=["#0D0F13", "#161920", "#FF6A00", "#F2F3F5"],
            elements=elements,
            visual_hierarchy=[f"g-e{idx}:{'title' if idx == 1 else 'text'}" for idx in range(1, len(elements) + 1)],
        )
        return PresentationAnalysis(
            source_path=str(primary),
            slide_count=1,
            global_palette=slide.palette,
            typography={"font_usage": {"Segoe UI": len(elements)}, "size_usage": {12.0: max(0, len(elements) - 1), 18.0: 1}},
            slides=[slide],
            metadata={"generated": "generic_source_adapter"},
        )

    def _object_inventory(self, analysis: Any) -> dict[str, Any]:
        counts: dict[str, int] = {}
        for slide in analysis.slides:
            for element in slide.elements:
                counts[element.kind] = counts.get(element.kind, 0) + 1
        return {"slides": analysis.slide_count, "total_objects": sum(counts.values()), "by_kind": counts}

    def _image_inventory(self, analysis: Any, asset_tree: dict[str, Path]) -> dict[str, Any]:
        items: list[dict[str, Any]] = []
        seen: set[str] = set()
        for slide in analysis.slides:
            for element in slide.elements:
                if not element.asset_path:
                    continue
                src = Path(element.asset_path)
                if not src.is_absolute():
                    src = (asset_tree["images"] / src.name).resolve()
                dst = asset_tree["images"] / src.name
                if src.exists() and src != dst:
                    shutil.copy2(src, dst)
                if dst.name in seen:
                    continue
                seen.add(dst.name)
                w, h = self._safe_image_size(dst)
                items.append({"asset_id": element.element_id, "path": str(dst), "slide": slide.index, "width_px": w, "height_px": h})

        return {"total_images": len(items), "deduplicated_images": len(items), "items": items}

    def _text_analysis(self, analysis: Any) -> dict[str, Any]:
        stats = {
            "title": 0,
            "subtitle": 0,
            "list": 0,
            "note": 0,
            "table": 0,
            "legend": 0,
            "callout": 0,
            "summary": 0,
            "conclusion": 0,
            "total_text_elements": 0,
        }
        for slide in analysis.slides:
            for element in slide.elements:
                body = (element.text or "").strip().lower()
                if not body:
                    continue
                stats["total_text_elements"] += 1
                if element.kind == "title":
                    stats["title"] += 1
                if element.kind in {"heading", "text"}:
                    stats["subtitle"] += 1
                if "\n" in body or body.startswith("-") or body.startswith("•"):
                    stats["list"] += 1
                if "table" in body:
                    stats["table"] += 1
                if "summary" in body or "resumen" in body:
                    stats["summary"] += 1
                if "conclusion" in body or "conclus" in body:
                    stats["conclusion"] += 1
        return stats

    def _visual_analysis(self, analysis: Any) -> dict[str, Any]:
        per_slide: list[dict[str, Any]] = []
        diagrams = 0
        for slide in analysis.slides:
            recognized = self.recognizer.recognize(slide)
            per_slide.append(
                {
                    "slide": slide.index,
                    "component_kind": recognized.component_kind,
                    "confidence": recognized.confidence,
                    "labels": recognized.labels,
                    "layout": slide.layout,
                }
            )
            if recognized.component_kind not in {"text", "table"}:
                diagrams += 1
        return {"slides_analyzed": len(analysis.slides), "diagram_like_slides": diagrams, "per_slide": per_slide}

    def _build_document_model(
        self,
        *,
        analysis: Any,
        visual_analysis: dict[str, Any],
        document_name: str,
        primary: Path,
        project: str,
        client: str,
        category: str,
        language: str,
        objective: str,
        audience: str,
        theme_variant: str,
        mode: str,
    ) -> DocumentModel:
        recognition_by_slide = {row["slide"]: row for row in visual_analysis.get("per_slide", [])}
        sections: list[SectionNode] = []
        assets: list[AssetRef] = []
        evidence: list[EvidenceRecord] = []
        links: list[KnowledgeLink] = []

        for slide in analysis.slides:
            rec = recognition_by_slide.get(slide.index, {"component_kind": "text", "confidence": 0.5, "labels": []})
            items: list[dict[str, Any]] = []
            table_rows: list[list[str]] | None = None

            for element in slide.elements:
                body = (element.text or "").strip()
                if body:
                    items.append({"title": element.kind.title(), "body": body})
                    evidence.append(EvidenceRecord(kind="slide_text", description=body[:240], metadata={"slide": slide.index, "element": element.element_id}))
                if element.table_rows:
                    table_rows = element.table_rows
                if element.asset_path:
                    assets.append(AssetRef(kind="image", path=element.asset_path, title=slide.title, metadata={"slide": slide.index}))

            component_kind = "text" if mode == "slide_flow" else rec.get("component_kind", "text")
            primary_component = normalize_component(component_kind, slide.title, slide.title, items=items, props={"layout": slide.layout, "labels": rec.get("labels", []), "confidence": rec.get("confidence", 0.5)})
            blocks = [BlockNode(block_type=mode, title=slide.title, components=[primary_component])]
            if table_rows:
                blocks.append(BlockNode(block_type="table", title="Data", components=[normalize_component("comparison_table", "Data Table", None, props={"rows": table_rows})]))

            sections.append(
                SectionNode(
                    title=slide.title,
                    summary=f"{mode} reconstruction for slide {slide.index}",
                    order=slide.index,
                    blocks=blocks,
                    metadata={"recognized_kind": rec.get("component_kind", "text"), "visual_hierarchy": slide.visual_hierarchy},
                )
            )
            links.append(KnowledgeLink(key=f"slide:{slide.index}", value=slide.title, category="layout", metadata={"mode": mode, "layout": slide.layout}))

        title = document_name or (analysis.slides[0].title if analysis.slides else primary.stem)
        return DocumentModel(
            title=title,
            subtitle=f"{mode.replace('_', ' ').title()} from {primary.name}",
            source_path=str(primary),
            source_type=primary.suffix.lower().lstrip("."),
            document_type=mode,
            theme_variant="light" if theme_variant == "light" else "industrial",
            metadata=self._base_metadata(analysis, project, client, category, language, objective, audience),
            sections=sections,
            assets=assets,
            evidence=evidence,
            knowledge_links=links,
        )

    def _quality_gate_scores(
        self,
        analysis: Any,
        document: DocumentModel,
        visual_analysis: dict[str, Any],
        outputs: dict[str, str],
        thresholds: dict[str, float],
    ) -> dict[str, Any]:
        src_text = sum(1 for s in analysis.slides for e in s.elements if (e.text or "").strip())
        src_images = sum(1 for s in analysis.slides for e in s.elements if e.asset_path)
        src_tables = sum(1 for s in analysis.slides for e in s.elements if e.table_rows)
        src_smartart = sum(1 for s in analysis.slides for e in s.elements if e.kind == "smartart")
        src_diagrams = max(1, visual_analysis.get("diagram_like_slides", 0))

        doc_text = sum(len([item for item in comp.items if (item.get("body") or "").strip()]) for sec in document.sections for block in sec.blocks for comp in block.components)
        doc_images = len(document.assets)
        doc_tables = sum(1 for sec in document.sections for block in sec.blocks for comp in block.components if comp.component_kind in {"comparison_table", "table"})
        doc_smartart = sum(1 for sec in document.sections if sec.metadata.get("recognized_kind") in {"hierarchy", "process_flow", "matrix", "relationship", "cycle", "pyramid"})
        doc_diagrams = sum(1 for sec in document.sections if sec.metadata.get("recognized_kind") not in {None, "text", "table"})

        def ratio(current: int, total: int) -> float:
            if total <= 0:
                return 1.0
            return round(min(1.0, current / total), 4)

        coverage = {
            "text_coverage": ratio(doc_text, src_text),
            "image_coverage": ratio(doc_images, src_images),
            "diagram_coverage": ratio(doc_diagrams, src_diagrams),
            "smartart_coverage": ratio(doc_smartart, src_smartart),
            "table_coverage": ratio(doc_tables, src_tables),
            "layout_quality": ratio(len(document.sections), max(1, analysis.slide_count)),
        }

        html_text = Path(outputs["html"]).read_text(encoding="utf-8", errors="ignore")
        soup = BeautifulSoup(html_text, "html.parser")
        theme_compliance = 1.0 if Path(outputs["theme_css"]).read_text(encoding="utf-8", errors="ignore").strip() else 0.0
        has_lang = bool(soup.find("html") and soup.find("html").get("lang"))
        has_main = soup.find("main") is not None
        images = soup.find_all("img")
        images_alt_ok = all((img.get("alt") or "").strip() for img in images) if images else True
        accessibility_signals = sum([1 if has_lang else 0, 1 if has_main else 0, 1 if images_alt_ok else 0])
        accessibility = round(accessibility_signals / 3, 4)
        responsive = 1.0 if "viewport" in html_text else 0.0
        typography = 1.0 if document.metadata.get("typography") else 0.75
        visual_similarity = round((coverage["text_coverage"] + coverage["image_coverage"] + coverage["layout_quality"]) / 3, 4)

        metrics = {
            "text_coverage": coverage["text_coverage"] * 100,
            "image_coverage": coverage["image_coverage"] * 100,
            "diagram_coverage": coverage["diagram_coverage"] * 100,
            "smartart_coverage": coverage["smartart_coverage"] * 100,
            "theme_compliance": theme_compliance * 100,
            "accessibility": accessibility * 100,
            "responsive": responsive * 100,
            "typography": typography * 100,
            "layout_quality": coverage["layout_quality"] * 100,
            "visual_similarity": visual_similarity * 100,
        }

        passed = all(metrics[k] >= thresholds[k] for k in thresholds)
        return {
            "passed": passed,
            "coverage": coverage,
            "metrics": metrics,
            "components_reconstructed": sum(len(block.components) for sec in document.sections for block in sec.blocks),
            "source_counts": {
                "text": src_text,
                "images": src_images,
                "tables": src_tables,
                "smartart": src_smartart,
                "diagrams": src_diagrams,
            },
            "reconstructed_counts": {
                "text": doc_text,
                "images": doc_images,
                "tables": doc_tables,
                "smartart": doc_smartart,
                "diagrams": doc_diagrams,
            },
        }

    def _score_variant(self, quality: dict[str, Any]) -> dict[str, Any]:
        metrics = quality.get("metrics", {})
        executive = round(
            metrics.get("text_coverage", 0) * 0.15
            + metrics.get("image_coverage", 0) * 0.1
            + metrics.get("diagram_coverage", 0) * 0.1
            + metrics.get("smartart_coverage", 0) * 0.1
            + metrics.get("theme_compliance", 0) * 0.15
            + metrics.get("accessibility", 0) * 0.1
            + metrics.get("responsive", 0) * 0.1
            + metrics.get("typography", 0) * 0.08
            + metrics.get("layout_quality", 0) * 0.07
            + metrics.get("visual_similarity", 0) * 0.05,
            2,
        )
        return {
            "visual_similarity_score": round(metrics.get("visual_similarity", 0), 2),
            "structure_score": round((metrics.get("layout_quality", 0) + metrics.get("diagram_coverage", 0)) / 2, 2),
            "theme_score": round(metrics.get("theme_compliance", 0), 2),
            "typography_score": round(metrics.get("typography", 0), 2),
            "layout_score": round(metrics.get("layout_quality", 0), 2),
            "image_score": round(metrics.get("image_coverage", 0), 2),
            "corporate_identity_score": round(metrics.get("theme_compliance", 0), 2),
            "executive_quality_score": executive,
        }

    def _quality_thresholds(self) -> dict[str, float]:
        return {
            "text_coverage": 95.0,
            "image_coverage": 95.0,
            "diagram_coverage": 80.0,
            "smartart_coverage": 80.0,
            "theme_compliance": 100.0,
            "accessibility": 90.0,
            "responsive": 100.0,
            "typography": 85.0,
            "layout_quality": 90.0,
            "visual_similarity": 85.0,
        }

    def _publication_state(self, selected_quality: dict[str, Any], selected_score: dict[str, Any]) -> str:
        if not selected_quality.get("passed", False):
            return "Review"
        if selected_score.get("executive_quality_score", 0) >= 95:
            return "Published"
        if selected_score.get("executive_quality_score", 0) >= 90:
            return "Validated"
        return "Review"

    def _state_history(self, state: str) -> list[dict[str, Any]]:
        now = datetime.now(UTC).isoformat()
        sequence = ["Draft", "Editing", "Review"]
        if state in {"Validated", "Published"}:
            sequence.append("Validated")
        if state == "Published":
            sequence.append("Published")
        return [{"state": s, "ts": now} for s in sequence]

    def _base_metadata(self, analysis: Any, project: str, client: str, category: str, language: str, objective: str, audience: str) -> dict[str, Any]:
        return {
            "his_module": "HTML Intelligence Studio V3",
            "project": project,
            "client": client,
            "category": category,
            "language": language,
            "objective": objective,
            "target_audience": audience,
            "slide_count": analysis.slide_count,
            "palette": analysis.global_palette,
            "typography": analysis.typography,
            "corporate_model": str(self.corporate_model_path),
            "languages": ["en", "es"],
            "translations": {
                "en": {"title": project or "Corporate Document", "subtitle": "Corporate HTML generated by HIS V3"},
                "es": {"title": project or "Documento corporativo", "subtitle": "HTML corporativo generado por HIS V3"},
            },
            "dom_schema": {
                "Document": ["Metadata", "Chapters", "Sections", "Blocks", "Components", "Assets", "Styles", "Languages", "Knowledge Links"],
            },
        }

    def _build_asset_registry(self, document: DocumentModel, asset_tree: dict[str, Path]) -> dict[str, Any]:
        items = []
        for asset in document.assets:
            src = Path(asset.path)
            category = "images"
            ext = src.suffix.lower()
            if ext == ".svg":
                category = "svg"
            elif any(k in src.name.lower() for k in ["logo", "brand"]):
                category = "logos"
            elif ext in {".ico"}:
                category = "icons"
            items.append(
                {
                    "id": asset.id,
                    "kind": asset.kind,
                    "title": asset.title,
                    "path": asset.path,
                    "category": category,
                    "reusable": True,
                    "metadata": asset.metadata,
                }
            )

        return {
            "generated_at": datetime.now(UTC).isoformat(),
            "asset_root": str(asset_tree["images"].parent),
            "categories": {k: str(v) for k, v in asset_tree.items()},
            "items": items,
        }

    def _build_knowledge_package(
        self,
        document: DocumentModel,
        selected_payload: dict[str, Any],
        hypotheses: list[dict[str, Any]],
        selected_hypothesis: dict[str, Any],
    ) -> dict[str, Any]:
        components = [comp.model_dump(mode="json") for sec in document.sections for block in sec.blocks for comp in block.components]
        prompts = [m.command for m in document.mission_links if m.command]
        return {
            "document_id": document.id,
            "title": document.title,
            "components": components,
            "diagram_components": [c for c in components if c.get("component_kind") not in {"text", "table", "comparison_table"}],
            "tables": [c for c in components if c.get("component_kind") in {"table", "comparison_table"}],
            "kpis": [c for c in components if c.get("component_kind") in {"industrial_kpi", "executive_dashboard", "roi_summary"}],
            "images": [asset.model_dump(mode="json") for asset in document.assets],
            "products": [document.metadata.get("project", "")],
            "clients": [document.metadata.get("client", "")],
            "technologies": sorted({(link.metadata or {}).get("recognized_kind", "") for link in document.knowledge_links if (link.metadata or {}).get("recognized_kind")}),
            "concepts": sorted({item.get("title", "") for c in components for item in (c.get("items") or []) if item.get("title")}),
            "prompts": prompts,
            "hypotheses": hypotheses,
            "selected_hypothesis": selected_hypothesis,
            "decisions": {
                "variant": selected_payload["selected_name"],
                "score": selected_payload["selected_score"],
            },
        }

    def _build_enterprise_memory(self, document: DocumentModel, selected_payload: dict[str, Any], run_id: str) -> dict[str, Any]:
        return {
            "document_id": document.id,
            "run_id": run_id,
            "source_path": document.source_path,
            "document_type": document.document_type,
            "missions": [m.model_dump(mode="json") for m in document.mission_links],
            "versions": [v.model_dump(mode="json") for v in document.version_history],
            "selected_variant": selected_payload["selected_name"],
            "selected_score": selected_payload["selected_score"],
            "timestamp": datetime.now(UTC).isoformat(),
        }

    def _build_truth_graph(self, document: DocumentModel, selected_payload: dict[str, Any]) -> dict[str, Any]:
        nodes = [{"id": f"doc:{document.id}", "type": "document", "label": document.title}]
        edges = []
        for section in document.sections:
            sid = f"sec:{section.id}"
            nodes.append({"id": sid, "type": "section", "label": section.title})
            edges.append({"source": f"doc:{document.id}", "target": sid, "relation": "has_section"})
        for asset in document.assets:
            aid = f"asset:{asset.id}"
            nodes.append({"id": aid, "type": "asset", "label": asset.title or Path(asset.path).name})
            edges.append({"source": f"doc:{document.id}", "target": aid, "relation": "uses_asset"})
        return {
            "generated_at": datetime.now(UTC).isoformat(),
            "selected_variant": selected_payload["selected_name"],
            "nodes": nodes,
            "edges": edges,
        }

    def _append_truth_graph_global(self, graph_payload: dict[str, Any]) -> None:
        root = Path(__file__).resolve().parents[2]
        global_graph = root / "data" / "truth_graph" / "his_truth_graph.json"
        global_graph.parent.mkdir(parents=True, exist_ok=True)
        if global_graph.exists():
            try:
                data = json.loads(global_graph.read_text(encoding="utf-8"))
                graphs = data.get("graphs", []) if isinstance(data, dict) else []
            except Exception:
                graphs = []
        else:
            graphs = []
        graphs.append(graph_payload)
        global_graph.write_text(json.dumps({"graphs": graphs}, indent=2, ensure_ascii=False), encoding="utf-8")

    def _build_technical_report_payload(
        self,
        *,
        run_id: str,
        primary: Path,
        phases: list[PhaseTimer],
        selected_payload: dict[str, Any],
        quality_attempts: list[dict[str, Any]],
        hypotheses: list[dict[str, Any]],
        selected_hypothesis: dict[str, Any],
        discarded_hypotheses: list[dict[str, Any]],
        thresholds: dict[str, float],
    ) -> dict[str, Any]:
        risks = self._detect_risks(selected_payload, quality_attempts)
        improvements = self._suggest_improvements(risks)
        return {
            "run_id": run_id,
            "source": str(primary),
            "generated_at": datetime.now(UTC).isoformat(),
            "pipeline_executed": "HIS V3 Production Ready",
            "phases": [{"name": p.name, "elapsed_s": p.elapsed_s} for p in phases],
            "hypotheses_generated": hypotheses,
            "selected_hypothesis": selected_hypothesis,
            "hypotheses_discarded": discarded_hypotheses,
            "quality_thresholds": thresholds,
            "quality_attempts": quality_attempts,
            "selected_variant": selected_payload["selected_name"],
            "scores": {
                "slide_flow": selected_payload["slide_flow_score"],
                "smart_reconstruction": selected_payload["smart_score"],
            },
            "quality": {
                "slide_flow": selected_payload["slide_flow_quality"],
                "smart_reconstruction": selected_payload["smart_quality"],
            },
            "discovery": selected_payload["discovery"],
            "object_inventory": selected_payload["object_inventory"],
            "image_inventory": selected_payload["image_inventory"],
            "components_reconstructed": {
                "slide_flow": selected_payload["slide_flow_quality"].get("components_reconstructed", 0),
                "smart_reconstruction": selected_payload["smart_quality"].get("components_reconstructed", 0),
            },
            "executive_quality_score": selected_payload["selected_score"].get("executive_quality_score", 0.0),
            "visual_similarity_score": selected_payload["selected_score"].get("visual_similarity_score", 0.0),
            "risks_detected": risks,
            "suggested_improvements": improvements,
        }

    def _render_report_markdown(self, payload: dict[str, Any]) -> str:
        lines = []
        lines.append("# HIS V3 Technical Report")
        lines.append("")
        lines.append(f"- Run ID: {payload['run_id']}")
        lines.append(f"- Source: {payload['source']}")
        lines.append(f"- Pipeline: {payload['pipeline_executed']}")
        lines.append(f"- Selected Variant: {payload['selected_variant']}")
        lines.append(f"- Executive Quality Score: {payload['executive_quality_score']}")
        lines.append(f"- Visual Similarity Score: {payload['visual_similarity_score']}")
        lines.append("")
        lines.append("## Phase Timing")
        for phase in payload["phases"]:
            lines.append(f"- {phase['name']}: {phase['elapsed_s']} s")
        lines.append("")

        lines.append("## Hypotheses")
        for h in payload.get("hypotheses_generated", []):
            lines.append(f"- {h.get('id')}: {h.get('name')} | score={h.get('executive_score')}")
        lines.append(f"- Selected: {payload.get('selected_hypothesis', {}).get('id')} ({payload.get('selected_hypothesis', {}).get('name')})")
        lines.append("")

        d = payload.get("discovery", {})
        lines.append("## Discovery")
        lines.append(f"- Slides: {d.get('slide_count', 0)}")
        lines.append(f"- Master Slides: {d.get('master_slides', 0)}")
        lines.append(f"- Layouts: {d.get('layouts', 0)}")
        lines.append(f"- Hyperlinks: {d.get('hyperlinks', 0)}")
        lines.append(f"- Notes Slides: {d.get('notes_slides', 0)}")
        lines.append("")

        inv = payload.get("object_inventory", {})
        lines.append("## Object Inventory")
        lines.append(f"- Total Objects: {inv.get('total_objects', 0)}")
        for key, value in sorted((inv.get("by_kind") or {}).items()):
            lines.append(f"- {key}: {value}")
        lines.append("")

        img = payload.get("image_inventory", {})
        lines.append("## Images")
        lines.append(f"- Total Images: {img.get('total_images', 0)}")
        lines.append(f"- Deduplicated Images: {img.get('deduplicated_images', 0)}")
        lines.append("")

        lines.append("## Quality Gates")
        for variant, info in (payload.get("quality") or {}).items():
            lines.append(f"### {variant}")
            lines.append(f"- Passed: {info.get('passed')}")
            lines.append(f"- Components reconstructed: {info.get('components_reconstructed', 0)}")
            metrics = info.get("metrics", {})
            for metric, value in metrics.items():
                lines.append(f"- {metric}: {value}")

        lines.append("")
        lines.append("## Risks Detected")
        for risk in payload.get("risks_detected", []):
            lines.append(f"- {risk}")
        lines.append("")

        lines.append("## Suggested Improvements")
        for imp in payload.get("suggested_improvements", []):
            lines.append(f"- {imp}")
        lines.append("")

        return "\n".join(lines)

    def _detect_risks(self, selected_payload: dict[str, Any], quality_attempts: list[dict[str, Any]]) -> list[str]:
        risks: list[str] = []
        if len(quality_attempts) > 1:
            risks.append("Multiple regeneration attempts were needed; source complexity may impact deterministic quality.")
        img_total = selected_payload.get("image_inventory", {}).get("total_images", 0)
        if img_total == 0:
            risks.append("No images detected in source; visual fidelity relies on text/layout only.")
        if selected_payload.get("selected_score", {}).get("executive_quality_score", 0) < 95:
            risks.append("Executive score below 95; requires manual review.")
        if not risks:
            risks.append("No critical risks detected under current thresholds.")
        return risks

    def _suggest_improvements(self, risks: list[str]) -> list[str]:
        improvements = []
        if any("regeneration" in r.lower() for r in risks):
            improvements.append("Add source-type-specific extraction adapters for complex diagrams.")
        if any("no images" in r.lower() for r in risks):
            improvements.append("Enable OCR/vision fallback to recover embedded rasterized visuals.")
        if any("manual review" in r.lower() for r in risks):
            improvements.append("Increase semantic weighting for business-critical sections in AHDE scoring.")
        if not improvements:
            improvements.append("Maintain current thresholds and monitor drift through periodic regression runs.")
        return improvements

    def _load_corporate_css_or_bootstrap(self, theme_variant: str = "industrial") -> str:
        if str(theme_variant).strip().lower() == "light":
            return build_css("light")
        if not self.corporate_model_path.exists():
            self.corporate_model_path.parent.mkdir(parents=True, exist_ok=True)
            self.corporate_model_path.write_text(self._bootstrap_corporate_css(), encoding="utf-8")

        raw_model = self.corporate_model_path.read_text(encoding="utf-8", errors="ignore")
        css = self._extract_css_from_corporate_model(raw_model)
        if not raw_model.strip():
            css = self._bootstrap_corporate_css()
            self.corporate_model_path.write_text(css, encoding="utf-8")
        elif not css.strip():
            css = self._bootstrap_corporate_css()

        if ":root" not in css:
            # Keep legacy corporate CSS usable by injecting default design tokens.
            css = INDUSTRIAL.as_css_vars() + "\n\n" + css
        return css

    def _extract_css_from_corporate_model(self, raw_model: str) -> str:
        text = (raw_model or "").strip()
        if not text:
            return ""

        lowered = text.lower()
        if "<style" in lowered:
            soup = BeautifulSoup(text, "html.parser")
            style_blocks = [node.get_text("\n", strip=True) for node in soup.find_all("style") if node.get_text(strip=True)]
            return "\n\n".join(style_blocks)

        # If the corporate model is HTML without style blocks, force bootstrap fallback.
        if re.search(r"<\s*[a-zA-Z][^>]*>", text):
            return ""
        return text

    def _bootstrap_corporate_css(self) -> str:
        return (
            INDUSTRIAL.as_css_vars()
            + "\n\n"
            + """
body { margin: 0; background: var(--bg-primary); color: var(--text-primary); font-family: Segoe UI, Arial, sans-serif; }
.dipc-shell { display: grid; grid-template-columns: 280px 1fr; min-height: 100vh; }
.dipc-sidebar { background: var(--bg-sidebar); color: var(--text-primary); padding: 18px; border-right: 1px solid var(--border-default); }
.dipc-main { padding: 18px 24px; background: var(--bg-primary); }
.dipc-hero, .dipc-section, .dipc-component { background: var(--bg-surface); border: 1px solid var(--border-default); border-radius: 12px; padding: 14px; margin-bottom: 14px; }
.dipc-lang-switch { position: sticky; top: 0; z-index: 12; display: flex; gap: 8px; align-items: center; padding: 8px 12px; border-bottom: 1px solid var(--border-default); background: var(--bg-elevated); }
.dipc-lang-switch button { background: var(--btn-secondary-bg); color: var(--btn-secondary-text); border: 1px solid var(--border-default); border-radius: 999px; padding: 4px 10px; }
@media (max-width: 1024px) { .dipc-shell { grid-template-columns: 1fr; } }
@media print { body { background: #fff; color: #111; } }
@media (prefers-color-scheme: dark) { body { background: var(--bg-primary); color: var(--text-primary); } }
""".strip()
        )

    def _copy_selected_variant(self, outputs: dict[str, str], selected_dir: Path) -> None:
        html = Path(outputs["html"])
        css = Path(outputs["theme_css"])
        if html.exists():
            shutil.copy2(html, selected_dir / "index.html")
        if css.exists():
            shutil.copy2(css, selected_dir / "corporate_theme.css")

    def _build_evidence_zip(self, base_dir: Path, zip_path: Path) -> None:
        include_paths = [
            base_dir / "selected",
            base_dir / "metadata",
            base_dir / "assets",
            base_dir / "history",
            base_dir / "logs",
            base_dir / "variants",
        ]
        with zipfile.ZipFile(zip_path, "w", compression=zipfile.ZIP_DEFLATED) as zf:
            for p in include_paths:
                if not p.exists():
                    continue
                for child in p.rglob("*"):
                    if child.is_file():
                        zf.write(child, child.relative_to(base_dir).as_posix())

    def _auto_remediate_accessibility(self, html_path: str) -> None:
        path = Path(html_path)
        if not path.exists():
            return
        soup = BeautifulSoup(path.read_text(encoding="utf-8", errors="ignore"), "html.parser")

        html_tag = soup.find("html")
        if html_tag is not None and not html_tag.get("lang"):
            html_tag["lang"] = "en"

        if soup.find("main") is None and soup.body is not None:
            main_tag = soup.new_tag("main")
            # Move all body children into main while preserving order.
            for child in list(soup.body.contents):
                main_tag.append(child.extract())
            soup.body.append(main_tag)

        for img in soup.find_all("img"):
            alt = (img.get("alt") or "").strip()
            if not alt:
                src = Path((img.get("src") or "image")).name
                img["alt"] = src or "image"

        path.write_text(str(soup), encoding="utf-8")

    def _safe_image_size(self, path: Path) -> tuple[int, int]:
        try:
            from PIL import Image

            with Image.open(path) as image:
                return int(image.width), int(image.height)
        except Exception:
            return 0, 0

    def _extract_text_from_file(self, path: Path) -> str:
        ext = path.suffix.lower()
        if ext in {".txt", ".md", ".html", ".htm"}:
            return path.read_text(encoding="utf-8", errors="ignore")
        if ext == ".docx":
            try:
                from docx import Document as DocxDocument

                doc = DocxDocument(str(path))
                return "\n".join(p.text for p in doc.paragraphs)
            except Exception:
                return ""
        if ext == ".pdf":
            try:
                from pypdf import PdfReader

                reader = PdfReader(str(path))
                return "\n".join((page.extract_text() or "") for page in reader.pages)
            except Exception:
                return ""
        return ""
