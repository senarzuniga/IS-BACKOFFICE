from __future__ import annotations

import html
import json
import zipfile
from datetime import UTC, datetime
from pathlib import Path

from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill
from pptx import Presentation

from backoffice.dipc.component_library import render_component
from backoffice.dipc.models import DocumentModel
from backoffice.dipc.publication_engine import PublicationEngine

from .bilingual_html_importer import BilingualHtmlImporter
from .corporate_models import (
    CorporateDocument,
    DeliveryPolicy,
    PublicationArtifact,
    PublicationFormat,
    ValidationResult,
)
from .corporate_registry import CorporateDocumentRegistry
from .repository_adapters import CorporateRepositorySettings, RepositoryAdapter, build_repository_adapters


class CorporatePublishingService:
    def __init__(self, settings: CorporateRepositorySettings | None = None) -> None:
        self.settings = settings or CorporateRepositorySettings.from_environment()
        self.adapters = build_repository_adapters(self.settings)
        self.registry = CorporateDocumentRegistry(self.settings.registry_path)
        self.importer = BilingualHtmlImporter()
        self.publisher = PublicationEngine()

    def publish_bilingual_html(
        self,
        *,
        repository_id: str,
        relative_path: str,
        title: str,
        client: str,
        project: str,
        delivery_policy: DeliveryPolicy | None = None,
        destination: str | None = None,
    ) -> CorporateDocument:
        adapter = self._adapter(repository_id)
        dependency = adapter.snapshot(relative_path)
        source = adapter.resolve_source(relative_path)
        models = self.importer.import_file(source, title=title)
        policy = delivery_policy or DeliveryPolicy()
        missing_languages = [language for language in policy.required_languages if language not in models]
        if missing_languages:
            raise ValueError(f"Missing required language branches: {', '.join(missing_languages)}")

        document = CorporateDocument(
            title=title,
            client=client,
            project=project,
            source_language="en",
            languages=policy.required_languages,
            source_document_model="",
            destination=destination or str(self.settings.output_root),
            delivery_policy=policy,
            dependencies=[dependency],
            status="validating",
        )
        root = self._version_root(document)
        model_dir = root / "model"
        model_dir.mkdir(parents=True, exist_ok=True)
        for language in policy.required_languages:
            model_path = model_dir / f"document_model_{language}.json"
            model_path.write_text(models[language].model_dump_json(indent=2), encoding="utf-8")
        document.source_document_model = str(model_dir / "document_model_en.json")

        artifacts: list[PublicationArtifact] = []
        for language in policy.required_languages:
            artifacts.extend(self._render_language(models[language], language, root, policy.allowed_formats, document))
        document.artifacts = artifacts
        document.validation = self._validate(document, models)
        document.status = "ready" if all(item.passed or item.severity != "error" for item in document.validation) else "failed"
        document.updated_at = datetime.now(UTC)
        self._write_manifest(document, root / "manifest.json")
        self.registry.save(document)
        return document

    def create_delivery_package(self, document_id: str) -> str:
        document = self.registry.refresh_staleness(document_id, self.adapters)
        if document.status == "stale":
            raise RuntimeError("Document dependencies changed; regenerate before packaging")
        if document.status not in {"ready", "published"}:
            raise RuntimeError(f"Document is not ready for delivery: {document.status}")
        allowed = set(document.delivery_policy.allowed_formats)
        selected = [artifact for artifact in document.artifacts if artifact.format in allowed]
        if not selected:
            raise RuntimeError("Delivery profile selects no generated artifacts")
        root = self._version_root(document)
        package_dir = root / "packages"
        package_dir.mkdir(parents=True, exist_ok=True)
        package_path = package_dir / f"{document.document_id}_v{document.version}_{document.delivery_policy.profile_id}.zip"
        validation_path = root / "validation" / "validation_report.json"
        validation_path.parent.mkdir(parents=True, exist_ok=True)
        validation_path.write_text(
            json.dumps([item.model_dump(mode="json") for item in document.validation], indent=2, ensure_ascii=False),
            encoding="utf-8",
        )
        manifest_path = root / "manifest.json"
        with zipfile.ZipFile(package_path, "w", compression=zipfile.ZIP_DEFLATED) as archive:
            for artifact in selected:
                path = Path(artifact.path)
                archive.write(path, f"documents/{artifact.language}/{path.name}")
            archive.write(manifest_path, "manifest.json")
            archive.write(validation_path, "validation_report.json")
        return str(package_path)

    def _render_language(
        self,
        model: DocumentModel,
        language: str,
        root: Path,
        formats: list[PublicationFormat],
        document: CorporateDocument,
    ) -> list[PublicationArtifact]:
        output_dir = root / language
        output_dir.mkdir(parents=True, exist_ok=True)
        stem = f"{document.document_id}_v{document.version}_{language}"
        paths: dict[PublicationFormat, Path] = {}
        if "html" in formats:
            path = output_dir / f"{stem}.html"
            path.write_text(self._render_corporate_html(model, language), encoding="utf-8")
            paths["html"] = path
        if "pdf" in formats:
            path = output_dir / f"{stem}.pdf"
            self.publisher.render_pdf(model, path)
            paths["pdf"] = path
        if "docx" in formats:
            path = output_dir / f"{stem}.docx"
            self.publisher.render_docx(model, path)
            paths["docx"] = path
        if "xlsx" in formats:
            path = output_dir / f"{stem}.xlsx"
            self._render_xlsx(model, path)
            paths["xlsx"] = path
        if "pptx" in formats:
            path = output_dir / f"{stem}.pptx"
            self._render_pptx(model, path)
            paths["pptx"] = path
        unsupported = set(formats).difference({"html", "pdf", "docx", "xlsx", "pptx"})
        if unsupported:
            raise ValueError(f"Renderers not yet available for: {', '.join(sorted(unsupported))}")
        return [
            PublicationArtifact(
                language=language,
                format=format_name,
                path=str(path),
                sha256=RepositoryAdapter.hash_file(path),
                size_bytes=path.stat().st_size,
            )
            for format_name, path in paths.items()
        ]

    def _render_xlsx(self, model: DocumentModel, path: Path) -> None:
        workbook = Workbook()
        summary = workbook.active
        summary.title = "Summary"
        summary.append([model.title])
        summary.append([model.subtitle or ""])
        summary.append([])
        summary.append(["Section", "Component", "Content"])
        summary["A1"].font = Font(bold=True, size=18, color="FFFFFF")
        summary["A1"].fill = PatternFill("solid", fgColor="0C0C0C")
        for section in model.sections:
            for block in section.blocks:
                for component in block.components:
                    summary.append([section.title, component.title or component.component_kind, component.body or ""])
        summary.freeze_panes = "A5"
        summary.column_dimensions["A"].width = 34
        summary.column_dimensions["B"].width = 34
        summary.column_dimensions["C"].width = 80

        tables = workbook.create_sheet("Tables")
        row_index = 1
        for section in model.sections:
            for block in section.blocks:
                for component in block.components:
                    rows = component.props.get("rows", [])
                    if not rows:
                        continue
                    tables.cell(row=row_index, column=1, value=section.title).font = Font(bold=True, color="FF5A10")
                    row_index += 1
                    for row in rows:
                        for column_index, value in enumerate(row, start=1):
                            tables.cell(row=row_index, column=column_index, value=value)
                        row_index += 1
                    row_index += 1

        sources = workbook.create_sheet("Sources")
        sources.append(["Evidence ID", "Kind", "Description", "Source"])
        for evidence in model.evidence:
            sources.append([evidence.id, evidence.kind, evidence.description, evidence.source_ref or ""])
        workbook.save(path)

    def _render_pptx(self, model: DocumentModel, path: Path) -> None:
        presentation = Presentation()
        title_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
        title_slide.shapes.title.text = model.title
        title_slide.placeholders[1].text = model.subtitle or "Corporate Intelligence"
        for section in model.sections:
            slide = presentation.slides.add_slide(presentation.slide_layouts[1])
            slide.shapes.title.text = section.title
            lines: list[str] = []
            for block in section.blocks:
                for component in block.components:
                    if component.title:
                        lines.append(component.title)
                    if component.body:
                        lines.append(component.body)
                    for item in component.items:
                        lines.append(f"{item.get('title', '')}: {item.get('body', '')}".strip(": "))
                    if len(lines) >= 8:
                        break
                if len(lines) >= 8:
                    break
            slide.placeholders[1].text = "\n".join(lines[:8]) or section.summary or ""
        presentation.save(path)

    def _render_corporate_html(self, model: DocumentModel, language: str) -> str:
        nav = []
        sections = []
        for section in model.sections:
            anchor = str(section.metadata.get("source_section_id", f"section-{section.order}"))
            nav.append(f'<a href="#{html.escape(anchor)}">{section.order:02d} · {html.escape(section.title)}</a>')
            components = "".join(
                render_component(component)
                for block in section.blocks
                for component in block.components
            )
            sections.append(
                f'<section id="{html.escape(anchor)}"><div class="section-no">{section.order:02d}</div>'
                f'<h2>{html.escape(section.title)}</h2>{components}</section>'
            )
        css = self._pcg_css()
        return f"""<!doctype html>
<html lang="{language}"><head><meta charset="utf-8"><meta name="viewport" content="width=device-width,initial-scale=1">
<title>{html.escape(model.title)}</title><meta name="description" content="{html.escape(model.subtitle or model.title)}"><style>{css}</style></head>
<body><header class="hero"><div class="hero-inner"><div class="eyebrow">Corporate intelligence report</div>
<h1>{html.escape(model.title)}</h1><p>{html.escape(model.subtitle or '')}</p><div class="meta"><span>{language.upper()}</span><span>Versioned delivery</span></div></div></header>
<div class="ribbon"><div class="ribbon-inner"><strong>Decision basis</strong><span>Evidence</span><span>Analysis</span><span>Actions</span><span>Validation</span></div></div>
<div class="layout"><aside><nav class="toc"><div class="toc-title">Contents</div>{''.join(nav)}</nav></aside><main>{''.join(sections)}</main></div>
<footer>Generated by IS_BACKOFFICE HTML Intelligence · See manifest for provenance.</footer></body></html>"""

    def _pcg_css(self) -> str:
        css_source = self.settings.ai_factory_root / "calgary_report_theme.css"
        if css_source.exists():
            css = css_source.read_text(encoding="utf-8")
        else:
            css = "body{font-family:Segoe UI,sans-serif;margin:0;color:#181817}main{max-width:1000px;margin:auto;padding:32px}"
        return css + "\n.dipc-component{margin:18px 0} footer{padding:28px;text-align:center;color:var(--muted)}"

    def _validate(self, document: CorporateDocument, models: dict[str, DocumentModel]) -> list[ValidationResult]:
        section_counts = {language: len(model.sections) for language, model in models.items()}
        results = [
            ValidationResult(
                check="language_structure_parity",
                passed=len(set(section_counts.values())) == 1,
                message=f"Section counts: {section_counts}",
            ),
            ValidationResult(
                check="artifacts_non_empty",
                passed=bool(document.artifacts) and all(item.size_bytes > 100 for item in document.artifacts),
                message=f"Validated {len(document.artifacts)} artifacts",
            ),
            ValidationResult(
                check="delivery_formats",
                passed=all(item.format in document.delivery_policy.allowed_formats for item in document.artifacts),
                message="Generated artifacts comply with the delivery profile",
            ),
            ValidationResult(
                check="source_provenance",
                passed=bool(document.dependencies) and all(len(item.sha256) == 64 for item in document.dependencies),
                message="Source dependencies include SHA-256 provenance",
            ),
        ]
        return results

    def _write_manifest(self, document: CorporateDocument, path: Path) -> None:
        path.parent.mkdir(parents=True, exist_ok=True)
        path.write_text(document.model_dump_json(indent=2), encoding="utf-8")

    def _version_root(self, document: CorporateDocument) -> Path:
        client_slug = "".join(character.lower() if character.isalnum() else "-" for character in document.client).strip("-")
        return self.settings.output_root / client_slug / document.document_id / f"v{document.version}"

    def _adapter(self, repository_id: str) -> RepositoryAdapter:
        try:
            return self.adapters[repository_id]
        except KeyError as exc:
            raise ValueError(f"Unknown repository: {repository_id}") from exc
