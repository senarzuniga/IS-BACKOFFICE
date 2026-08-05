from __future__ import annotations

import hashlib
from pathlib import Path
from typing import Any

from .layout_analysis import build_visual_hierarchy, detect_layout
from .models import ElementStyle, PresentationAnalysis, SlideAnalysis, SlideElement


class PowerPointParser:
    def parse(self, pptx_path: str, assets_dir: str) -> PresentationAnalysis:
        from pptx import Presentation

        source = Path(pptx_path).expanduser().resolve()
        asset_root = Path(assets_dir)
        asset_root.mkdir(parents=True, exist_ok=True)

        prs = Presentation(str(source))
        slides: list[SlideAnalysis] = []
        image_hash_to_path: dict[str, str] = {}

        slide_width = float(prs.slide_width.pt)
        slide_height = float(prs.slide_height.pt)

        for idx, slide in enumerate(prs.slides, start=1):
            title = self._extract_slide_title(slide) or f"Slide {idx}"
            parsed_elements: list[SlideElement] = []
            local_palette: list[str] = []

            for sidx, shape in enumerate(slide.shapes, start=1):
                element = self._parse_shape(
                    shape=shape,
                    slide_idx=idx,
                    shape_idx=sidx,
                    asset_root=asset_root,
                    dedupe=image_hash_to_path,
                )
                if not element:
                    continue
                parsed_elements.append(element)
                if element.style.color:
                    local_palette.append(element.style.color)

            element_boxes = [{"x": e.x, "y": e.y, "w": e.w, "h": e.h} for e in parsed_elements]
            layout = detect_layout(slide_width, element_boxes)
            background = self._extract_background(slide)

            if background:
                local_palette.append(background)

            slide_model = SlideAnalysis(
                index=idx,
                title=title,
                width=slide_width,
                height=slide_height,
                background=background,
                layout=layout,
                palette=[c for c in dict.fromkeys(local_palette) if c],
                elements=parsed_elements,
            )
            slide_model.visual_hierarchy = build_visual_hierarchy(slide_model)
            slides.append(slide_model)

        all_colors: list[str] = []
        font_usage: dict[str, int] = {}
        size_usage: dict[float, int] = {}
        for slide in slides:
            all_colors.extend(slide.palette)
            for el in slide.elements:
                if el.style.font_name:
                    font_usage[el.style.font_name] = font_usage.get(el.style.font_name, 0) + 1
                if el.style.font_size_pt:
                    key = round(float(el.style.font_size_pt), 1)
                    size_usage[key] = size_usage.get(key, 0) + 1

        return PresentationAnalysis(
            source_path=str(source),
            slide_count=len(slides),
            global_palette=[c for c in dict.fromkeys(all_colors) if c],
            typography={"font_usage": font_usage, "size_usage": size_usage},
            slides=slides,
            metadata={
                "slide_width_pt": slide_width,
                "slide_height_pt": slide_height,
                "deduplicated_images": len(image_hash_to_path),
            },
        )

    def _parse_shape(
        self,
        shape: Any,
        slide_idx: int,
        shape_idx: int,
        asset_root: Path,
        dedupe: dict[str, str],
    ) -> SlideElement | None:
        element_id = f"s{slide_idx}-e{shape_idx}"
        x = float(shape.left.pt) if hasattr(shape, "left") else 0.0
        y = float(shape.top.pt) if hasattr(shape, "top") else 0.0
        w = float(shape.width.pt) if hasattr(shape, "width") else 0.0
        h = float(shape.height.pt) if hasattr(shape, "height") else 0.0

        text, style = self._extract_text_and_style(shape)

        if getattr(shape, "has_table", False):
            rows = []
            for row in shape.table.rows:
                rows.append([cell.text.strip() for cell in row.cells])
            return SlideElement(
                element_id=element_id,
                kind="table",
                x=x,
                y=y,
                w=w,
                h=h,
                text=text,
                style=style,
                table_rows=rows,
            )

        if getattr(shape, "has_chart", False):
            chart = self._extract_chart(shape)
            return SlideElement(
                element_id=element_id,
                kind="chart",
                x=x,
                y=y,
                w=w,
                h=h,
                text=text,
                style=style,
                chart=chart,
            )

        if getattr(shape, "shape_type", None) is not None and str(shape.shape_type).endswith("PICTURE"):
            image_path = self._export_picture(shape, asset_root, dedupe)
            return SlideElement(
                element_id=element_id,
                kind="image",
                x=x,
                y=y,
                w=w,
                h=h,
                style=style,
                asset_path=image_path,
                metadata={"name": getattr(shape, "name", "")},
            )

        if "smartart" in str(getattr(shape, "name", "")).lower():
            return SlideElement(
                element_id=element_id,
                kind="smartart",
                x=x,
                y=y,
                w=w,
                h=h,
                text=text,
                style=style,
            )

        if text.strip():
            kind = "text"
            if shape_idx == 1 or text.strip().count("\n") == 0 and (style.font_size_pt or 0) >= 24:
                kind = "title"
            elif (style.font_size_pt or 0) >= 18:
                kind = "heading"
            return SlideElement(element_id=element_id, kind=kind, x=x, y=y, w=w, h=h, text=text, style=style)

        fill_color = self._shape_fill_color(shape)
        return SlideElement(
            element_id=element_id,
            kind="shape",
            x=x,
            y=y,
            w=w,
            h=h,
            style=ElementStyle(color=fill_color),
            metadata={"name": getattr(shape, "name", "")},
        )

    def _extract_text_and_style(self, shape: Any) -> tuple[str, ElementStyle]:
        text = ""
        style = ElementStyle()
        if not getattr(shape, "has_text_frame", False):
            return text, style

        text = shape.text_frame.text or ""
        for paragraph in shape.text_frame.paragraphs:
            if style.align is None and paragraph.alignment is not None:
                style.align = str(paragraph.alignment)
            for run in paragraph.runs:
                font = run.font
                if style.font_name is None and font.name:
                    style.font_name = font.name
                if style.font_size_pt is None and font.size is not None:
                    style.font_size_pt = float(font.size.pt)
                if style.bold is None and font.bold is not None:
                    style.bold = bool(font.bold)
                if style.italic is None and font.italic is not None:
                    style.italic = bool(font.italic)
                if style.color is None:
                    rgb = self._font_rgb(font)
                    if rgb:
                        style.color = rgb
        if style.color is None:
            style.color = self._shape_fill_color(shape)
        return text, style

    def _extract_chart(self, shape: Any) -> dict[str, Any]:
        data: dict[str, Any] = {"series": [], "categories": []}
        try:
            chart = shape.chart
            data["chart_type"] = str(chart.chart_type)
            categories = []
            if chart.plots and chart.plots[0].categories:
                categories = [str(c.label) for c in chart.plots[0].categories]
            data["categories"] = categories
            for series in chart.series:
                points = []
                for point in series.points:
                    points.append(point.value)
                data["series"].append({"name": getattr(series, "name", "series"), "values": points})
        except Exception as exc:
            data["error"] = str(exc)
        return data

    def _export_picture(self, shape: Any, asset_root: Path, dedupe: dict[str, str]) -> str | None:
        try:
            blob = shape.image.blob
            digest = hashlib.sha256(blob).hexdigest()
            if digest in dedupe:
                return dedupe[digest]

            ext = shape.image.ext or "png"
            filename = f"img_{digest[:14]}.{ext}"
            output = asset_root / filename
            output.write_bytes(blob)
            rel = output.as_posix()
            dedupe[digest] = rel
            return rel
        except Exception:
            return None

    def _extract_slide_title(self, slide: Any) -> str:
        if getattr(slide.shapes, "title", None) is not None and slide.shapes.title is not None:
            title = slide.shapes.title.text or ""
            if title.strip():
                return title.strip()

        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text = (shape.text_frame.text or "").strip()
                if text:
                    return text.split("\n", 1)[0][:120]
        return ""

    def _extract_background(self, slide: Any) -> str | None:
        try:
            fill = slide.background.fill
            if fill and fill.fore_color and fill.fore_color.rgb:
                return f"#{fill.fore_color.rgb}"
        except Exception:
            return None
        return None

    def _shape_fill_color(self, shape: Any) -> str | None:
        try:
            fill = shape.fill
            if fill and fill.fore_color and fill.fore_color.rgb:
                return f"#{fill.fore_color.rgb}"
        except Exception:
            return None
        return None

    def _font_rgb(self, font: Any) -> str | None:
        try:
            color = font.color
            if color and color.rgb:
                return f"#{color.rgb}"
        except Exception:
            return None
        return None
