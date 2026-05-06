from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

NAVY = RGBColor(21, 56, 94)
BLUE = RGBColor(41, 98, 153)
RED = RGBColor(214, 46, 56)
ORANGE = RGBColor(232, 111, 55)
INK = RGBColor(34, 42, 52)
SLATE = RGBColor(90, 103, 118)
STONE = RGBColor(240, 235, 226)
WHITE = RGBColor(255, 255, 255)
PALE_BLUE = RGBColor(227, 236, 244)
PALE_RED = RGBColor(252, 236, 236)
PALE_ORANGE = RGBColor(252, 242, 235)
LIGHT_GRAY = RGBColor(227, 229, 232)

TITLE_FONT = "Bahnschrift SemiBold"
BODY_FONT = "Aptos"

ROOT = Path(__file__).resolve().parent
ASSET_DIR = ROOT / "ingecart_assets_v2"
OUTPUT_PATH = ROOT / "PSC_VISALIA_INGECART_LINETEX_MASTER.pptx"

LOGO_PATH = ASSET_DIR / "Copia%20de%20sello%202%20(1)_edited.png"
MACHINE_COVER = ASSET_DIR / "secondhand_page1.png"
MACHINE_OVERVIEW = ASSET_DIR / "secondhand_page3.png"
REEL_IMAGE = ASSET_DIR / "secondhand_page4.png"
BHS_LETTER = ASSET_DIR / "bhs_support_letter_page2.png"
ROBOT_IMAGE = ASSET_DIR / "Robot%20FFG-2.jpg"


def set_background(slide, color=STONE):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_box(slide, left, top, width, height, fill_color, line_color=None, radius=False):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = line_color or fill_color
    return shape


def add_text(slide, left, top, width, height, text, size=18, color=INK, bold=False,
             font_name=BODY_FONT, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = text
    font = run.font
    font.name = font_name
    font.size = Pt(size)
    font.bold = bold
    font.color.rgb = color
    paragraph.alignment = align
    return box


def add_bullets(slide, left, top, width, height, items, size=18, color=INK, bullet_color=None,
                font_name=BODY_FONT, space_after=6):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    for index, item in enumerate(items):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = item
        paragraph.level = 0
        paragraph.bullet = True
        paragraph.space_after = Pt(space_after)
        font = paragraph.font
        font.name = font_name
        font.size = Pt(size)
        font.color.rgb = color
        if bullet_color:
            paragraph.runs[0].font.color.rgb = bullet_color
    return box


def add_picture(slide, path, left, top, width=None, height=None):
    if path.exists():
        if width is not None and height is not None:
            return slide.shapes.add_picture(str(path), left, top, width=width, height=height)
        if width is not None:
            return slide.shapes.add_picture(str(path), left, top, width=width)
        if height is not None:
            return slide.shapes.add_picture(str(path), left, top, height=height)
        return slide.shapes.add_picture(str(path), left, top)
    placeholder = add_box(slide, left, top, width or Inches(2), height or Inches(1.2), LIGHT_GRAY, LIGHT_GRAY)
    add_text(slide, left + Inches(0.15), top + Inches(0.15), (width or Inches(2)) - Inches(0.3), Inches(0.6),
             "Asset unavailable", size=12, color=SLATE)
    return placeholder


def add_header_footer(slide, page_number, section_label):
    add_box(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.18), RED)
    add_box(slide, Inches(0), Inches(7.08), SLIDE_W, Inches(0.42), NAVY)
    add_text(slide, Inches(0.45), Inches(0.22), Inches(5.0), Inches(0.3), section_label.upper(), size=10, color=SLATE, bold=True)
    add_text(slide, Inches(0.45), Inches(7.15), Inches(8.8), Inches(0.2),
             "INGECART ENGINEERING & AUDITING  |  LINETEX COMMERCIAL CHANNEL USA", size=10, color=WHITE, bold=True)
    add_text(slide, Inches(12.2), Inches(7.13), Inches(0.7), Inches(0.22), str(page_number), size=10, color=WHITE, bold=True, align=PP_ALIGN.RIGHT)
    if LOGO_PATH.exists():
        add_picture(slide, LOGO_PATH, Inches(11.78), Inches(0.2), width=Inches(1.08), height=Inches(0.42))


def add_title(slide, title, subtitle=None):
    add_text(slide, Inches(0.6), Inches(0.62), Inches(8.6), Inches(0.65), title, size=28, color=NAVY, bold=True, font_name=TITLE_FONT)
    if subtitle:
        add_text(slide, Inches(0.6), Inches(1.1), Inches(8.5), Inches(0.35), subtitle, size=14, color=SLATE)


def add_card(slide, left, top, width, height, heading, value, caption, fill_color=WHITE):
    card = add_box(slide, left, top, width, height, fill_color, LIGHT_GRAY, radius=True)
    add_text(slide, left + Inches(0.18), top + Inches(0.12), width - Inches(0.36), Inches(0.28), heading.upper(), size=10, color=SLATE, bold=True)
    add_text(slide, left + Inches(0.18), top + Inches(0.42), width - Inches(0.36), Inches(0.42), value, size=21, color=NAVY, bold=True, font_name=TITLE_FONT)
    add_text(slide, left + Inches(0.18), top + Inches(0.87), width - Inches(0.36), height - Inches(0.98), caption, size=12, color=INK)
    return card


def add_chevron(slide, left, top, width, height, text, fill_color):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = fill_color
    add_text(slide, left + Inches(0.16), top + Inches(0.1), width - Inches(0.28), height - Inches(0.2), text,
             size=15, color=WHITE, bold=True, font_name=TITLE_FONT, valign=MSO_ANCHOR.MIDDLE)
    return shape


def cover_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, STONE)
    add_header_footer(slide, 1, "Proyecto PSC Visalia")
    add_box(slide, Inches(0.65), Inches(1.0), Inches(5.1), Inches(5.55), WHITE, WHITE, radius=True)
    add_text(slide, Inches(0.95), Inches(1.25), Inches(4.2), Inches(1.15), "PSC VISALIA\nPROYECTO CORRUGADORA 98", size=26, color=NAVY, bold=True, font_name=TITLE_FONT)
    add_text(slide, Inches(0.95), Inches(2.45), Inches(4.1), Inches(0.65), "Modernizacion integral + reinstalacion + automatizacion", size=19, color=ORANGE, bold=True, font_name=TITLE_FONT)
    add_text(slide, Inches(0.95), Inches(3.22), Inches(4.1), Inches(1.1),
             "Reposicionamiento de una linea BHS 2500 mm en USA con nuevo flujo logístico, nueva arquitectura de handling y una base comercial reforzada por Ingecart y Linetex.",
             size=15, color=INK)
    add_box(slide, Inches(0.95), Inches(4.7), Inches(3.8), Inches(1.2), PALE_RED, PALE_RED, radius=True)
    add_text(slide, Inches(1.15), Inches(4.95), Inches(3.4), Inches(0.25), "NARRATIVA COMERCIAL", size=11, color=RED, bold=True)
    add_text(slide, Inches(1.15), Inches(5.2), Inches(3.2), Inches(0.55),
             "No es solo mover una maquina: es reconstruir capacidad, fiabilidad y automatizacion alrededor de un activo probado.",
             size=14, color=INK)
    add_picture(slide, MACHINE_COVER, Inches(5.95), Inches(0.95), width=Inches(6.65), height=Inches(5.8))
    return slide


def executive_summary_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_header_footer(slide, 2, "Resumen ejecutivo")
    add_title(slide, "Resumen ejecutivo", "Una propuesta para transformar disponibilidad, soporte y flujo operativo")
    add_text(slide, Inches(0.72), Inches(1.62), Inches(5.3), Inches(1.55),
             "El proyecto integra una corrugadora BHS 2500 mm de segunda mano con mejoras mecanicas, electricas y logísticas para relanzarla en Visalia con una operacion mucho mas ordenada y escalable.",
             size=20, color=INK)
    add_bullets(slide, Inches(0.75), Inches(3.0), Inches(5.2), Inches(2.5), [
        "Base instalada BHS con working width 2.500 mm y plataforma hasta 300 m/min.",
        "Reinstalacion completa con desmontaje documentado, transporte, montaje y puesta en marcha.",
        "Nueva capa de handling: conveyors, WIP, AMR-ready y home area automatizada.",
        "Enfoque USA: Ingecart ejecuta; Linetex actua como canal comercial y cara local del proyecto."
    ], size=18)
    add_card(slide, Inches(6.25), Inches(1.62), Inches(2.0), Inches(1.5), "Plataforma", "BHS 2500 mm", "Wet-end y dry-end duplex con inventario completo de equipos")
    add_card(slide, Inches(8.47), Inches(1.62), Inches(2.0), Inches(1.5), "Velocidad", "300 m/min", "Referencia del offer sheet y del roll stand IIM")
    add_card(slide, Inches(10.69), Inches(1.62), Inches(1.95), Inches(1.5), "Calendario", "May 26 - Feb 27", "Curva de cobros y pagos ya modelada en los Excel de proyecto")
    add_card(slide, Inches(6.25), Inches(3.45), Inches(3.0), Inches(1.7), "Subtotales visibles", "> USD 12.0M", "Suma de hojas revisadas: services, peripherals, conveyors e Ingetrans")
    add_card(slide, Inches(9.48), Inches(3.45), Inches(3.16), Inches(1.7), "Pieza critica", "Roll handling", "Nuevo IIM shaftless hydraulic roll stand y reorganizacion de alimentacion")
    add_box(slide, Inches(6.25), Inches(5.48), Inches(6.39), Inches(0.82), NAVY, NAVY, radius=True)
    add_text(slide, Inches(6.5), Inches(5.67), Inches(5.85), Inches(0.3),
             "Mensaje de venta: de-risking tecnico + continuidad industrial + automatizacion interna", size=16, color=WHITE, bold=True)
    return slide


def objectives_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_header_footer(slide, 3, "Objetivos")
    add_title(slide, "Objetivos del proyecto", "Que debe conseguir PSC Visalia con esta inversion")
    add_box(slide, Inches(0.72), Inches(1.65), Inches(5.35), Inches(4.95), WHITE, LIGHT_GRAY, radius=True)
    add_bullets(slide, Inches(1.0), Inches(2.0), Inches(4.85), Inches(4.25), [
        "Reinstalar una linea probada con trazabilidad mecanica y electrica desde origen hasta SAT.",
        "Reducir la dependencia de equipamiento obsoleto mediante actualizaciones selectivas en handling, control y soporte.",
        "Preparar la operacion para un flujo mas limpio entre corrugadora, WIP, conversion y residuos.",
        "Alinear la propuesta con una presencia comercial local en USA a traves de Linetex.",
        "Construir una base operativa lista para entrenamiento, spare parts y servicio postventa."
    ], size=19)
    add_chevron(slide, Inches(6.4), Inches(2.0), Inches(1.75), Inches(0.8), "MOVER", BLUE)
    add_chevron(slide, Inches(8.15), Inches(2.0), Inches(1.75), Inches(0.8), "REHACER", RED)
    add_chevron(slide, Inches(9.9), Inches(2.0), Inches(1.75), Inches(0.8), "AUTOMATIZAR", ORANGE)
    add_box(slide, Inches(6.45), Inches(3.05), Inches(5.15), Inches(3.35), PALE_BLUE, PALE_BLUE, radius=True)
    add_text(slide, Inches(6.78), Inches(3.35), Inches(4.55), Inches(0.45), "Lectura estrategica", size=18, color=NAVY, bold=True, font_name=TITLE_FONT)
    add_text(slide, Inches(6.78), Inches(3.88), Inches(4.45), Inches(1.95),
             "La propuesta combina tres niveles de valor: rescatar un activo industrial robusto, elevarlo con nuevas capas de ingenieria y convertir el material handling en una ventaja competitiva real.",
             size=18, color=INK)
    return slide


def scope_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_header_footer(slide, 4, "Alcance")
    add_title(slide, "Alcance de suministro y servicios", "Seis paquetes que convierten una compra industrial en una solucion integral")
    cards = [
        ("Core line", "BHS corrugator", "Base 2.500 mm, wet-end y dry-end duplex, activos 1997-2017", PALE_BLUE),
        ("Services", "Desmontaje + montaje", "Despiece, marcado, documentacion, cranes, viajes e instalacion", WHITE),
        ("Roll handling", "IIM roll stand", "Shaftless hydraulic roll stand 305 m/min con RE chucks y RE brakes", PALE_ORANGE),
        ("Peripherals", "Kitchen + utilities", "Flexamix, tuberias, boiler, steam tubes y auxiliares", WHITE),
        ("Flow", "Conveyors + WIP", "Double stacker with strappers, AMR-ready reception y home area", PALE_BLUE),
        ("Controls", "Software + integration", "Cuadros electricos, software, FAT/SAT y entrenamiento", PALE_RED),
    ]
    lefts = [Inches(0.75), Inches(4.52), Inches(8.29)]
    tops = [Inches(1.8), Inches(4.1)]
    index = 0
    for top in tops:
        for left in lefts:
            heading, value, caption, fill = cards[index]
            add_card(slide, left, top, Inches(3.2), Inches(1.9), heading, value, caption, fill)
            index += 1
    return slide


def machine_base_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_header_footer(slide, 5, "Base instalada")
    add_title(slide, "Base instalada y activo de partida", "Inventario de corrugadora revisado en el dossier de segunda mano")
    add_picture(slide, MACHINE_OVERVIEW, Inches(0.75), Inches(1.55), width=Inches(7.45), height=Inches(5.25))
    add_card(slide, Inches(8.55), Inches(1.75), Inches(3.3), Inches(1.2), "Width", "2.500 mm", "Working width declarada en el dossier")
    add_card(slide, Inches(8.55), Inches(3.1), Inches(3.3), Inches(1.2), "Configuration", "Duplex / Duplex", "Wet-end y dry-end duplex")
    add_card(slide, Inches(8.55), Inches(4.45), Inches(3.3), Inches(1.2), "Flutes", "B/E + C/B", "Dos single facers con combinaciones de flauta en inventario")
    add_text(slide, Inches(8.55), Inches(6.0), Inches(3.45), Inches(0.5),
             "El inventario muestra reel stands, splicers, preheater, glue unit, bridge, double facer, rotary shear, slitter scorer, cutoff y stacker.",
             size=13, color=INK)
    return slide


def modernization_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_header_footer(slide, 6, "Modernizacion")
    add_title(slide, "Modernizacion y upgrade por obsolescencia", "La logica del proyecto no es mantener lo viejo: es reconstruir una plataforma fiable")
    add_box(slide, Inches(0.8), Inches(1.65), Inches(5.7), Inches(4.95), WHITE, LIGHT_GRAY, radius=True)
    add_text(slide, Inches(1.05), Inches(1.95), Inches(2.2), Inches(0.3), "BASE LEGACY", size=15, color=SLATE, bold=True)
    add_bullets(slide, Inches(1.05), Inches(2.35), Inches(4.9), Inches(3.8), [
        "Equipos de origen 1997 con upgrades puntuales 2007, 2009, 2012, 2015 y 2017.",
        "Reel stands y splicers BHS robustos pero dependientes de una integracion mas actual.",
        "Riesgo natural de dispersion de soporte, repuestos y curva de aprendizaje si no se rearquitectura el conjunto."
    ], size=18)
    add_box(slide, Inches(6.8), Inches(1.65), Inches(5.7), Inches(4.95), PALE_BLUE, PALE_BLUE, radius=True)
    add_text(slide, Inches(7.05), Inches(1.95), Inches(2.4), Inches(0.3), "TARGET STATE", size=15, color=BLUE, bold=True)
    add_bullets(slide, Inches(7.05), Inches(2.35), Inches(4.9), Inches(3.8), [
        "Roll handling nuevo con IIM shaftless hydraulic roll stand: 2.500 mm, 305 m/min, Allen Bradley PLC y Atos hydraulic pack.",
        "Nueva capa de conveyors, rails, home area, software y cuadros electricos para ordenar flujo y mantenimiento.",
        "Modelo orientado a FAT, SAT, formacion y soporte comercial local en USA a traves de Linetex."
    ], size=18)
    return slide


def relocation_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_header_footer(slide, 7, "Reinstalacion")
    add_title(slide, "Layout de reinstalacion y roadmap fisico", "Un traslado controlado, documentado y orientado a un arranque con menos friccion")
    add_picture(slide, MACHINE_OVERVIEW, Inches(0.78), Inches(1.7), width=Inches(6.2), height=Inches(4.05))
    add_chevron(slide, Inches(7.3), Inches(2.0), Inches(1.65), Inches(0.8), "TURKEY", BLUE)
    add_chevron(slide, Inches(8.93), Inches(2.0), Inches(1.65), Inches(0.8), "SHIPPING", RED)
    add_chevron(slide, Inches(10.56), Inches(2.0), Inches(1.65), Inches(0.8), "VISALIA", ORANGE)
    add_bullets(slide, Inches(7.3), Inches(3.0), Inches(4.9), Inches(2.2), [
        "Desmontaje completo con marcado documental de instalaciones mecanicas y electricas.",
        "Preparacion para container loading y coordinacion de crane operations.",
        "Montaje en USA con equipo local, commissioning y acceptance testing."
    ], size=17)
    add_box(slide, Inches(7.3), Inches(5.35), Inches(4.85), Inches(0.95), PALE_ORANGE, PALE_ORANGE, radius=True)
    add_text(slide, Inches(7.6), Inches(5.65), Inches(4.2), Inches(0.3),
             "La hoja de servicios incluye desmontaje, cranes, instalacion y gastos asociados.", size=15, color=INK, bold=True)
    return slide


def roll_handling_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_header_footer(slide, 8, "Roll handling")
    add_title(slide, "Roll handling: nueva pieza critica del sistema", "El nuevo roll stand define fiabilidad, seguridad de cambio y continuidad de papel")
    add_picture(slide, REEL_IMAGE, Inches(0.82), Inches(1.7), width=Inches(6.2), height=Inches(4.8))
    add_bullets(slide, Inches(7.25), Inches(1.9), Inches(4.7), Inches(2.7), [
        "2500 mm width, shaftless hydraulic roll stand.",
        "Reel max diameter 1525 mm y reel weight 2000 kg.",
        "Machine speed 305 m/min con RE chucks 3\"/6\" y RE brakes.",
        "Allen Bradley PLC sin programa y Atos hydraulic pack M-18."
    ], size=18)
    add_card(slide, Inches(7.25), Inches(4.95), Inches(2.35), Inches(1.15), "Offer price", "71.685 EUR", "Ex works del IIM roll stand")
    add_card(slide, Inches(9.77), Inches(4.95), Inches(2.35), Inches(1.15), "Freight", "7.560 EUR", "DDP Incoterms transport")
    add_text(slide, Inches(7.25), Inches(6.3), Inches(4.9), Inches(0.28),
             "Lead time declarado: 120 dias desde el primer hito de pago. Garantia: 1 ano.", size=13, color=SLATE)
    return slide


def conveyors_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_header_footer(slide, 9, "Conveyors")
    add_title(slide, "Conveyors, WIP y home area automatizada", "El verdadero salto operativo aparece en el material handling")
    add_box(slide, Inches(0.8), Inches(1.8), Inches(12.0), Inches(0.85), NAVY, NAVY, radius=True)
    add_text(slide, Inches(1.0), Inches(2.07), Inches(11.4), Inches(0.25),
             "Corrugator output  ->  double conveyor stacker  ->  AMR reception  ->  WIP conversion  ->  waste movement  ->  home area",
             size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_card(slide, Inches(0.9), Inches(3.0), Inches(2.8), Inches(1.6), "End of line", "USD 1.170M", "Double conveyor line output stacker with strappers")
    add_card(slide, Inches(3.95), Inches(3.0), Inches(2.8), Inches(1.6), "AMR / WIP", "USD 0.715M", "Reception conveyor, WIP placement y waste movement")
    add_card(slide, Inches(7.0), Inches(3.0), Inches(2.8), Inches(1.6), "Rails", "USD 0.241M", "New corrugar line rails 12 meters")
    add_card(slide, Inches(10.05), Inches(3.0), Inches(2.35), Inches(1.6), "Home area", "USD 1.014M", "Automatic loading and unloading ramp")
    add_box(slide, Inches(0.9), Inches(5.0), Inches(11.5), Inches(1.2), PALE_BLUE, PALE_BLUE, radius=True)
    add_text(slide, Inches(1.15), Inches(5.35), Inches(10.8), Inches(0.55),
             "El mensaje para el cliente es claro: no se vende solo capacidad de corrugado, se vende una operacion con menos trafico manual, buffers mejor controlados y una transicion limpia hacia automatizacion interna.",
             size=17, color=INK)
    return slide


def automation_concept_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_header_footer(slide, 10, "Operacion")
    add_title(slide, "Concepto operativo WIP / AMR", "Como debe leerse el flujo futuro dentro de la planta")
    stages = [
        ("1", "Paper infeed", BLUE),
        ("2", "Corrugator", NAVY),
        ("3", "Stacking", RED),
        ("4", "AMR handoff", ORANGE),
        ("5", "WIP buffer", BLUE),
        ("6", "Conversion", NAVY),
    ]
    current_left = Inches(0.9)
    for number, label, color in stages:
        add_box(slide, current_left, Inches(2.2), Inches(1.75), Inches(1.0), color, color, radius=True)
        add_text(slide, current_left + Inches(0.1), Inches(2.42), Inches(1.55), Inches(0.3), f"{number}. {label}", size=15, color=WHITE, bold=True, font_name=TITLE_FONT, align=PP_ALIGN.CENTER)
        current_left += Inches(1.95)
    add_box(slide, Inches(1.0), Inches(4.0), Inches(3.55), Inches(1.8), WHITE, LIGHT_GRAY, radius=True)
    add_text(slide, Inches(1.3), Inches(4.33), Inches(2.9), Inches(0.28), "QUE GANA PSC", size=13, color=SLATE, bold=True)
    add_bullets(slide, Inches(1.25), Inches(4.68), Inches(3.0), Inches(0.95), [
        "Menos trafico manual en la salida de corrugadora.",
        "Mayor control visual del WIP y del desperdicio.",
        "Base lista para integrar logica de trazabilidad."
    ], size=15)
    add_box(slide, Inches(4.85), Inches(4.0), Inches(3.55), Inches(1.8), PALE_ORANGE, PALE_ORANGE, radius=True)
    add_text(slide, Inches(5.15), Inches(4.33), Inches(2.9), Inches(0.28), "QUE GANA INGENIERIA", size=13, color=SLATE, bold=True)
    add_bullets(slide, Inches(5.1), Inches(4.68), Inches(3.0), Inches(0.95), [
        "Puntos de integracion claros entre mecanica, electrica y software.",
        "Secuencias FAT/SAT mas defendibles.",
        "Un alcance de automatizacion facilmente escalable."
    ], size=15)
    add_box(slide, Inches(8.7), Inches(4.0), Inches(3.55), Inches(1.8), PALE_RED, PALE_RED, radius=True)
    add_text(slide, Inches(9.0), Inches(4.33), Inches(2.9), Inches(0.28), "MENSAJE COMERCIAL", size=13, color=SLATE, bold=True)
    add_bullets(slide, Inches(8.95), Inches(4.68), Inches(3.0), Inches(0.95), [
        "Operacion mas limpia.",
        "Capacidad de crecimiento sin rehacer todo el layout.",
        "Menor friccion entre corrugado y conversion."
    ], size=15)
    return slide


def bhs_support_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_header_footer(slide, 11, "BHS support")
    add_title(slide, "BHS integration y condicion de soporte", "Lo que la carta revisada pide formalizar para cerrar el proyecto con seguridad")
    add_box(slide, Inches(0.82), Inches(1.72), Inches(5.5), Inches(4.9), WHITE, LIGHT_GRAY, radius=True)
    add_text(slide, Inches(1.1), Inches(2.0), Inches(4.8), Inches(0.3), "SUPPORT FRAMEWORK REQUESTED", size=16, color=NAVY, bold=True, font_name=TITLE_FONT)
    add_bullets(slide, Inches(1.05), Inches(2.42), Inches(4.85), Inches(3.8), [
        "Technical supervision during the machine disassembly in Turkey.",
        "Crane operations management and support.",
        "Local or nearby service in the United States after installation and acceptance.",
        "Maintenance training and process training.",
        "Spare parts supply plus technical service and after-sales support."
    ], size=18)
    add_picture(slide, BHS_LETTER, Inches(6.65), Inches(1.7), width=Inches(5.55), height=Inches(4.95))
    add_text(slide, Inches(6.75), Inches(6.35), Inches(5.3), Inches(0.28),
             "Nota de rigor: el documento revisado es una solicitud de confirmacion escrita, no una carta final emitida por BHS.",
             size=12, color=SLATE)
    return slide


def engineering_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_header_footer(slide, 12, "Ingenieria y software")
    add_title(slide, "Ingenieria y arquitectura de control", "El upgrade necesita una narrativa tecnica coherente, no solo una lista de equipos")
    add_box(slide, Inches(0.95), Inches(2.0), Inches(2.0), Inches(1.0), BLUE, BLUE, radius=True)
    add_text(slide, Inches(1.1), Inches(2.36), Inches(1.7), Inches(0.25), "SENSORS &\nMOTION", size=18, color=WHITE, bold=True, font_name=TITLE_FONT, align=PP_ALIGN.CENTER)
    add_box(slide, Inches(3.35), Inches(2.0), Inches(2.0), Inches(1.0), NAVY, NAVY, radius=True)
    add_text(slide, Inches(3.53), Inches(2.36), Inches(1.62), Inches(0.25), "PLC / HMI\nLAYER", size=18, color=WHITE, bold=True, font_name=TITLE_FONT, align=PP_ALIGN.CENTER)
    add_box(slide, Inches(5.75), Inches(2.0), Inches(2.0), Inches(1.0), RED, RED, radius=True)
    add_text(slide, Inches(5.93), Inches(2.36), Inches(1.62), Inches(0.25), "CONVEYORS\n& WIP", size=18, color=WHITE, bold=True, font_name=TITLE_FONT, align=PP_ALIGN.CENTER)
    add_box(slide, Inches(8.15), Inches(2.0), Inches(2.0), Inches(1.0), ORANGE, ORANGE, radius=True)
    add_text(slide, Inches(8.35), Inches(2.36), Inches(1.55), Inches(0.25), "MES /\nREPORTING", size=18, color=WHITE, bold=True, font_name=TITLE_FONT, align=PP_ALIGN.CENTER)
    add_box(slide, Inches(10.55), Inches(2.0), Inches(2.0), Inches(1.0), BLUE, BLUE, radius=True)
    add_text(slide, Inches(10.75), Inches(2.36), Inches(1.55), Inches(0.25), "REMOTE\nSERVICE", size=18, color=WHITE, bold=True, font_name=TITLE_FONT, align=PP_ALIGN.CENTER)
    add_bullets(slide, Inches(1.0), Inches(4.0), Inches(5.25), Inches(2.0), [
        "El offer del roll stand ya cita Allen Bradley PLC y sensores M-18 con logica de diametro de bobina.",
        "La hoja Ingetrans incorpora cuadros electricos y software como parte del alcance.",
        "La presentacion debe vender integracion total: mecanica, control, WIP y mantenimiento."
    ], size=18)
    add_picture(slide, ROBOT_IMAGE, Inches(7.25), Inches(3.55), width=Inches(4.9), height=Inches(2.45))
    return slide


def timeline_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_header_footer(slide, 13, "Calendario")
    add_title(slide, "Calendario de ejecucion y curva de caja", "El flujo financiero anticipa los picos criticos del proyecto")
    months = [
        ("May 26", "Cobro inicial\n+ ingenieria", BLUE),
        ("Jun 26", "Ingenieria", NAVY),
        ("Jul 26", "Maquina +\ndesmontaje", RED),
        ("Aug 26", "Transporte", ORANGE),
        ("Sep 26", "Upgrade elec/mec", RED),
        ("Oct 26", "Contingencia", BLUE),
        ("Nov 26", "Contingencia", NAVY),
        ("Dec 26", "Embalaje", BLUE),
        ("Jan 27", "Transporte", ORANGE),
        ("Feb 27", "Cierre", NAVY),
    ]
    left = Inches(0.72)
    for label, note, color in months:
        add_box(slide, left, Inches(2.45), Inches(1.16), Inches(1.55), WHITE, LIGHT_GRAY, radius=True)
        add_box(slide, left, Inches(2.45), Inches(1.16), Inches(0.28), color, color, radius=True)
        add_text(slide, left + Inches(0.14), Inches(2.82), Inches(0.88), Inches(0.2), label, size=12, color=NAVY, bold=True, font_name=TITLE_FONT, align=PP_ALIGN.CENTER)
        add_text(slide, left + Inches(0.1), Inches(3.17), Inches(0.96), Inches(0.55), note, size=11, color=INK, align=PP_ALIGN.CENTER)
        left += Inches(1.24)
    add_box(slide, Inches(0.82), Inches(4.55), Inches(4.0), Inches(1.35), PALE_RED, PALE_RED, radius=True)
    add_text(slide, Inches(1.1), Inches(4.88), Inches(3.4), Inches(0.27), "Pico de tension", size=13, color=RED, bold=True)
    add_text(slide, Inches(1.1), Inches(5.15), Inches(3.4), Inches(0.45),
             "Julio concentra maquina + desmontaje. Septiembre absorbe upgrade mecanico/electrico.", size=15, color=INK)
    add_box(slide, Inches(4.95), Inches(4.55), Inches(3.95), Inches(1.35), PALE_BLUE, PALE_BLUE, radius=True)
    add_text(slide, Inches(5.22), Inches(4.88), Inches(3.2), Inches(0.27), "Cobros", size=13, color=BLUE, bold=True)
    add_text(slide, Inches(5.22), Inches(5.15), Inches(3.25), Inches(0.45),
             "El modelo de flujo refleja cobros de USD 200k por mes a lo largo del periodo principal.", size=15, color=INK)
    add_box(slide, Inches(9.02), Inches(4.55), Inches(3.25), Inches(1.35), PALE_ORANGE, PALE_ORANGE, radius=True)
    add_text(slide, Inches(9.28), Inches(4.88), Inches(2.8), Inches(0.27), "Lectura comercial", size=13, color=ORANGE, bold=True)
    add_text(slide, Inches(9.28), Inches(5.15), Inches(2.7), Inches(0.45),
             "El calendario favorece una venta por hitos, no un simple CAPEX aislado.", size=15, color=INK)
    return slide


def fat_sat_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_header_footer(slide, 14, "FAT / SAT")
    add_title(slide, "FAT, SAT y puesta en marcha", "La credibilidad del proyecto depende de cerrar bien la secuencia de validacion")
    steps = [
        ("FAT mechanical", "Revision de equipos reacondicionados, checklists y trazabilidad documental", BLUE),
        ("FAT controls", "Secuencias de seguridad, I/O, recipe logic y handshakes", NAVY),
        ("SAT installation", "Alineacion, utilities, cold tests y dry run", RED),
        ("Ramp-up", "Training, acceptance protocol y paso a operacion", ORANGE),
    ]
    top = Inches(1.9)
    for title, caption, color in steps:
        add_box(slide, Inches(0.9), top, Inches(2.15), Inches(1.0), color, color, radius=True)
        add_text(slide, Inches(1.05), top + Inches(0.34), Inches(1.85), Inches(0.25), title, size=18, color=WHITE, bold=True, font_name=TITLE_FONT, align=PP_ALIGN.CENTER)
        add_box(slide, Inches(3.25), top, Inches(8.65), Inches(1.0), WHITE, LIGHT_GRAY, radius=True)
        add_text(slide, Inches(3.55), top + Inches(0.28), Inches(8.1), Inches(0.45), caption, size=17, color=INK)
        top += Inches(1.2)
    add_text(slide, Inches(0.95), Inches(6.75), Inches(11.0), Inches(0.25),
             "La historia correcta para el cliente: FAT/SAT no es burocracia, es la herramienta para volver predecible una reinstalacion compleja.",
             size=14, color=SLATE)
    return slide


def benefits_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_header_footer(slide, 15, "Beneficios")
    add_title(slide, "Beneficios esperados para PSC Visalia", "Lo que realmente compra el cliente cuando compra esta solucion")
    add_card(slide, Inches(0.85), Inches(1.8), Inches(2.8), Inches(1.8), "Continuidad", "Menos riesgo", "Desmontaje documentado, reinstall planificada y soporte estructurado")
    add_card(slide, Inches(3.9), Inches(1.8), Inches(2.8), Inches(1.8), "Produccion", "Mas control", "Base BHS robusta con handling rediseñado alrededor de la salida de linea", fill_color=PALE_BLUE)
    add_card(slide, Inches(6.95), Inches(1.8), Inches(2.8), Inches(1.8), "People & safety", "Menos friccion", "Mejor gestion de bobinas, stacking y transito interno", fill_color=PALE_ORANGE)
    add_card(slide, Inches(10.0), Inches(1.8), Inches(2.3), Inches(1.8), "Scalability", "AMR-ready", "Flujo preparado para siguientes capas de automatizacion", fill_color=PALE_RED)
    add_box(slide, Inches(0.9), Inches(4.15), Inches(11.35), Inches(1.95), NAVY, NAVY, radius=True)
    add_text(slide, Inches(1.25), Inches(4.55), Inches(10.7), Inches(0.9),
             "La clave comercial es hablar de transformacion operacional: mayor orden interno, mejor mantenimiento, una salida de corrugadora menos caotica y una infraestructura lista para crecer con automatizacion futura.",
             size=22, color=WHITE, bold=True)
    return slide


def differentiators_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_header_footer(slide, 16, "Diferenciadores")
    add_title(slide, "Porque Ingecart + Linetex pueden vender este proyecto", "No es una simple compra de equipos usados")
    add_bullets(slide, Inches(0.92), Inches(1.9), Inches(5.35), Inches(4.7), [
        "Ingecart combina vision mecanica, electrica y logistica en un solo discurso tecnico.",
        "Linetex funciona como marca/comercializacion en USA, acercando el proyecto al cliente final.",
        "La propuesta une equipo principal, handling, software y fase de arranque en una sola historia.",
        "El valor diferencial esta en convertir activos dispersos en una solucion industrial cohesionada."
    ], size=20)
    add_box(slide, Inches(6.55), Inches(1.9), Inches(5.45), Inches(4.7), PALE_BLUE, PALE_BLUE, radius=True)
    add_text(slide, Inches(6.88), Inches(2.22), Inches(4.8), Inches(0.28), "PARTNERSHIP MESSAGE", size=15, color=NAVY, bold=True)
    add_text(slide, Inches(6.88), Inches(2.7), Inches(4.55), Inches(2.6),
             "Linetex debe aparecer como la cara comercial en Estados Unidos, mientras Ingecart sostiene el contenido de ingenieria, integracion, coordinacion del desmontaje y puesta en marcha. Eso reduce distancia comercial sin vaciar el valor tecnico.",
             size=20, color=INK)
    return slide


def investment_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_header_footer(slide, 17, "Inversion")
    add_title(slide, "Vista de inversion", "Subtotales visibles en las hojas comerciales revisadas")
    categories = [
        ("Services", 3.041, BLUE),
        ("Peripherals", 2.264, NAVY),
        ("Conveyors", 3.770, RED),
        ("Ingetrans", 2.972, ORANGE),
    ]
    max_value = max(value for _, value, _ in categories)
    top = Inches(2.1)
    for name, value, color in categories:
        add_text(slide, Inches(0.95), top + Inches(0.12), Inches(1.7), Inches(0.25), name, size=16, color=INK, bold=True)
        bar_width = Inches(7.6 * (value / max_value))
        add_box(slide, Inches(2.3), top, bar_width, Inches(0.48), color, color, radius=True)
        add_text(slide, Inches(10.2), top + Inches(0.08), Inches(1.8), Inches(0.25), f"USD {value:.3f}M", size=16, color=NAVY, bold=True, align=PP_ALIGN.RIGHT)
        top += Inches(0.82)
    add_box(slide, Inches(0.95), Inches(5.7), Inches(11.2), Inches(0.9), PALE_RED, PALE_RED, radius=True)
    add_text(slide, Inches(1.25), Inches(5.98), Inches(10.65), Inches(0.28),
             "Nota: esta slide muestra el subtotal visible de las hojas revisadas. Debe contrastarse con el cierre comercial final antes de emitir oferta definitiva.",
             size=15, color=INK)
    return slide


def next_steps_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_header_footer(slide, 18, "Proximos pasos")
    add_title(slide, "Proximos pasos", "Como convertir esta narrativa en cierre comercial y ejecucion")
    actions = [
        "Congelar el scope tecnico y los paquetes economicos a presentar a PSC.",
        "Validar el layout objetivo y utilidades de planta en Visalia.",
        "Cerrar la condicion BHS de supervision, training, spare parts y after-sales.",
        "Fijar el plan FAT/SAT, responsables y criterios de acceptance.",
        "Lanzar la fase comercial final bajo marca Linetex con respaldo tecnico de Ingecart."
    ]
    top = Inches(1.95)
    for idx, action in enumerate(actions, start=1):
        add_box(slide, Inches(1.0), top, Inches(0.65), Inches(0.65), NAVY, NAVY, radius=True)
        add_text(slide, Inches(1.13), top + Inches(0.18), Inches(0.38), Inches(0.18), str(idx), size=18, color=WHITE, bold=True, font_name=TITLE_FONT, align=PP_ALIGN.CENTER)
        add_box(slide, Inches(1.85), top, Inches(10.3), Inches(0.65), WHITE, LIGHT_GRAY, radius=True)
        add_text(slide, Inches(2.15), top + Inches(0.15), Inches(9.7), Inches(0.3), action, size=17, color=INK)
        top += Inches(0.84)
    return slide


def closing_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, NAVY)
    add_box(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.18), RED)
    add_box(slide, Inches(0), Inches(7.08), SLIDE_W, Inches(0.42), RED)
    if LOGO_PATH.exists():
        add_picture(slide, LOGO_PATH, Inches(0.85), Inches(0.6), width=Inches(1.45), height=Inches(0.55))
    add_text(slide, Inches(0.92), Inches(1.65), Inches(5.8), Inches(1.15),
             "Cerrar esta operacion es darle a PSC una linea relanzada, no una mudanza de maquinas.",
             size=28, color=WHITE, bold=True, font_name=TITLE_FONT)
    add_text(slide, Inches(0.95), Inches(3.0), Inches(5.55), Inches(1.4),
             "Ingecart aporta la ingenieria y la integracion. Linetex aporta cercania comercial en USA. Juntas, ambas marcas convierten un activo BHS probado en una plataforma industrial lista para volver a competir.",
             size=20, color=WHITE)
    add_picture(slide, MACHINE_COVER, Inches(7.0), Inches(1.1), width=Inches(5.5), height=Inches(4.7))
    add_text(slide, Inches(0.95), Inches(6.55), Inches(4.8), Inches(0.25), "Prepared from project spreadsheets, commercial offer, BHS support request and machine dossier", size=12, color=WHITE)
    add_text(slide, Inches(12.15), Inches(7.13), Inches(0.6), Inches(0.22), "19", size=10, color=WHITE, bold=True, align=PP_ALIGN.RIGHT)
    return slide


def build_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    cover_slide(prs)
    executive_summary_slide(prs)
    objectives_slide(prs)
    scope_slide(prs)
    machine_base_slide(prs)
    modernization_slide(prs)
    relocation_slide(prs)
    roll_handling_slide(prs)
    conveyors_slide(prs)
    automation_concept_slide(prs)
    bhs_support_slide(prs)
    engineering_slide(prs)
    timeline_slide(prs)
    fat_sat_slide(prs)
    benefits_slide(prs)
    differentiators_slide(prs)
    investment_slide(prs)
    next_steps_slide(prs)
    closing_slide(prs)
    prs.save(OUTPUT_PATH)


if __name__ == "__main__":
    build_presentation()
    print(OUTPUT_PATH)