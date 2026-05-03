"""DocumentParser — extracts raw text and tables from documents of various types."""

from __future__ import annotations

import csv
import json
import re
import subprocess
import tempfile
from datetime import datetime
from pathlib import Path
from typing import Callable

from document_analysis.models import DocumentInfo, DocumentType


class DocumentParser:
    """Parse a file into a :class:`DocumentInfo` object with raw text and tables."""

    def parse(self, file_path: str | Path) -> DocumentInfo:
        path = Path(file_path).expanduser().resolve()
        ext = path.suffix.lower()
        doc_type = _ext_to_type(ext)

        stat = path.stat() if path.exists() else None
        info = DocumentInfo(
            file_path=str(path),
            file_name=path.name,
            doc_type=doc_type,
            size_bytes=stat.st_size if stat else 0,
            modified_at=datetime.fromtimestamp(stat.st_mtime) if stat else None,
        )

        try:
            parser = self._parser_for_extension(ext)
            parser(path, info)
        except Exception as exc:  # noqa: BLE001
            info.error = f"{type(exc).__name__}: {exc}"

        info.word_count = len(info.raw_text.split()) if info.raw_text else 0
        info.char_count = len(info.raw_text) if info.raw_text else 0
        return info

    # ------------------------------------------------------------------
    # Format-specific parsers
    # ------------------------------------------------------------------

    def _parser_for_extension(self, ext: str) -> Callable[[Path, DocumentInfo], None]:
        parsers: dict[str, Callable[[Path, DocumentInfo], None]] = {
            # Documents
            ".pdf": self._parse_pdf,
            ".docx": self._parse_docx,
            ".doc": self._parse_doc,
            ".docm": self._parse_docx,
            ".rtf": self._parse_rtf,
            ".odt": self._parse_odt,
            # Spreadsheets
            ".xlsx": self._parse_xlsx,
            ".xls": self._parse_xls,
            ".xlsm": self._parse_xlsx,
            ".xlsb": self._parse_xlsx,
            ".ods": self._parse_xlsx,
            ".csv": self._parse_csv,
            ".tsv": self._parse_csv,
            ".tab": self._parse_csv,
            # Presentations
            ".pptx": self._parse_pptx,
            ".ppt": self._parse_ppt,
            ".pptm": self._parse_pptx,
            # Text
            ".txt": self._parse_txt,
            ".md": self._parse_txt,
            ".rst": self._parse_txt,
            ".log": self._parse_txt,
            ".text": self._parse_txt,
            # Structured
            ".json": self._parse_json,
            ".xml": self._parse_xml,
            ".yaml": self._parse_yaml,
            ".yml": self._parse_yaml,
            # Web
            ".html": self._parse_html,
            ".htm": self._parse_html,
            ".mht": self._parse_html,
            ".mhtml": self._parse_html,
            # Email
            ".eml": self._parse_email,
            ".msg": self._parse_email,
            ".emlx": self._parse_email,
            # Images (OCR)
            ".png": self._parse_image_ocr,
            ".jpg": self._parse_image_ocr,
            ".jpeg": self._parse_image_ocr,
            ".tiff": self._parse_image_ocr,
            ".tif": self._parse_image_ocr,
            ".bmp": self._parse_image_ocr,
            ".gif": self._parse_image_ocr,
            ".webp": self._parse_image_ocr,
        }
        return parsers.get(ext, self._parse_unknown)

    def _parse_pdf(self, path: Path, info: DocumentInfo) -> None:
        try:
            import fitz  # PyMuPDF

            doc = fitz.open(str(path))
            info.page_count = doc.page_count
            pages: list[str] = []
            for page in doc:
                pages.append(page.get_text())
            info.raw_text = "\n\n".join(pages)
            doc.close()
        except ImportError:
            # Fallback: try pdfplumber
            try:
                import pdfplumber

                with pdfplumber.open(str(path)) as pdf:
                    info.page_count = len(pdf.pages)
                    pages = [p.extract_text() or "" for p in pdf.pages]
                    info.raw_text = "\n\n".join(pages)
                    # Also extract tables
                    for page in pdf.pages:
                        for table in page.extract_tables():
                            info.tables.append({"rows": table})
            except ImportError:
                info.error = "Neither PyMuPDF nor pdfplumber is installed."

    def _parse_docx(self, path: Path, info: DocumentInfo) -> None:
        from docx import Document

        doc = Document(str(path))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        info.raw_text = "\n".join(paragraphs)
        for table in doc.tables:
            rows = [[cell.text for cell in row.cells] for row in table.rows]
            info.tables.append({"rows": rows})

    def _parse_xlsx(self, path: Path, info: DocumentInfo) -> None:
        try:
            import openpyxl

            wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
            all_text: list[str] = []
            for sheet in wb.worksheets:
                rows: list[list[str]] = []
                for row in sheet.iter_rows(values_only=True):
                    cells = [str(c) if c is not None else "" for c in row]
                    rows.append(cells)
                    all_text.append(" ".join(cells))
                info.tables.append({"sheet": sheet.title, "rows": rows})
            info.raw_text = "\n".join(all_text)
            wb.close()
            return
        except Exception:
            self._parse_xls(path, info)

    def _parse_csv(self, path: Path, info: DocumentInfo) -> None:
        encoding = _detect_encoding(path)
        delimiter = "\t" if path.suffix.lower() in {".tsv", ".tab"} else ","
        with open(path, encoding=encoding, newline="", errors="replace") as fh:
            reader = csv.reader(fh, delimiter=delimiter)
            rows = list(reader)
        info.tables.append({"rows": rows})
        info.raw_text = "\n".join(" ".join(r) for r in rows)
        info.encoding = encoding

    def _parse_txt(self, path: Path, info: DocumentInfo) -> None:
        encoding = _detect_encoding(path)
        info.raw_text = path.read_text(encoding=encoding, errors="replace")
        info.encoding = encoding

    def _parse_json(self, path: Path, info: DocumentInfo) -> None:
        encoding = _detect_encoding(path)
        data = json.loads(path.read_text(encoding=encoding, errors="replace"))
        info.raw_text = json.dumps(data, indent=2, ensure_ascii=False)
        info.metadata["json_keys"] = list(data.keys()) if isinstance(data, dict) else []
        info.encoding = encoding

    def _parse_yaml(self, path: Path, info: DocumentInfo) -> None:
        encoding = _detect_encoding(path)
        raw = path.read_text(encoding=encoding, errors="replace")
        try:
            import yaml

            data = yaml.safe_load(raw)
            info.raw_text = json.dumps(data, indent=2, ensure_ascii=False)
            info.metadata["json_keys"] = list(data.keys()) if isinstance(data, dict) else []
        except Exception:
            info.raw_text = raw
        info.encoding = encoding

    def _parse_xml(self, path: Path, info: DocumentInfo) -> None:
        text = path.read_text(encoding="utf-8", errors="replace")
        # Strip tags for raw text extraction
        info.raw_text = re.sub(r"<[^>]+>", " ", text)
        info.raw_text = re.sub(r"\s+", " ", info.raw_text).strip()

    def _parse_html(self, path: Path, info: DocumentInfo) -> None:
        text = path.read_text(encoding="utf-8", errors="replace")
        try:
            from bs4 import BeautifulSoup

            soup = BeautifulSoup(text, "html.parser")
            for script_or_style in soup(["script", "style"]):
                script_or_style.decompose()
            info.raw_text = soup.get_text(separator="\n")
        except ImportError:
            info.raw_text = re.sub(r"<[^>]+>", " ", text)
            info.raw_text = re.sub(r"\s+", " ", info.raw_text).strip()

    def _parse_pptx(self, path: Path, info: DocumentInfo) -> None:
        try:
            from pptx import Presentation

            prs = Presentation(str(path))
            slides: list[str] = []
            for slide in prs.slides:
                texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        texts.append(shape.text)
                slides.append("\n".join(texts))
            info.page_count = len(prs.slides)
            info.raw_text = "\n\n".join(slides)
        except ImportError:
            info.error = "python-pptx is not installed."

    def _parse_doc(self, path: Path, info: DocumentInfo) -> None:
        # 1) antiword
        try:
            result = subprocess.run(["antiword", str(path)], capture_output=True, text=True, timeout=30)
            if result.returncode == 0 and result.stdout.strip():
                info.raw_text = result.stdout.strip()
                return
        except Exception:
            pass

        # 2) textract
        try:
            import textract

            text = textract.process(str(path)).decode("utf-8", errors="ignore")
            if text.strip():
                info.raw_text = text.strip()
                return
        except Exception:
            pass

        # 3) OLE stream scan
        try:
            import olefile

            ole = olefile.OleFileIO(str(path))
            if ole.exists("WordDocument"):
                stream = ole.openstream("WordDocument")
                data = stream.read()
                text = data.decode("utf-8", errors="ignore")
                text = re.sub(r"[^\x20-\x7E\x0A\x0D\xC0-\xFF]", " ", text)
                if len(text.strip()) > 50:
                    info.raw_text = text.strip()
                    ole.close()
                    return
            ole.close()
        except Exception:
            pass

        # 4) python-docx fallback
        try:
            from docx import Document

            doc = Document(str(path))
            text = "\n".join(p.text for p in doc.paragraphs)
            if text.strip():
                info.raw_text = text.strip()
                return
        except Exception:
            pass

        # 5) LibreOffice conversion
        try:
            with tempfile.TemporaryDirectory() as tmpdir:
                result = subprocess.run(
                    ["soffice", "--headless", "--convert-to", "txt:Text", "--outdir", tmpdir, str(path)],
                    capture_output=True,
                    text=True,
                    timeout=60,
                )
                if result.returncode == 0:
                    txt_files = list(Path(tmpdir).glob("*.txt"))
                    if txt_files:
                        text = txt_files[0].read_text(encoding="utf-8", errors="ignore")
                        if text.strip():
                            info.raw_text = text.strip()
                            return
        except Exception:
            pass

        self._parse_unknown(path, info)
        if not info.raw_text:
            info.error = "Could not parse .doc file. Try converting to .docx or .txt first."

    def _parse_rtf(self, path: Path, info: DocumentInfo) -> None:
        try:
            from striprtf.striprtf import rtf_to_text

            rtf_text = path.read_text(encoding="utf-8", errors="ignore")
            text = rtf_to_text(rtf_text)
            if text.strip():
                info.raw_text = text.strip()
                return
        except Exception:
            pass

        self._parse_unknown(path, info)

    def _parse_xls(self, path: Path, info: DocumentInfo) -> None:
        try:
            import pandas as pd

            all_sheets = pd.read_excel(str(path), sheet_name=None, engine="xlrd")
            text_parts: list[str] = []
            for sheet_name, df in all_sheets.items():
                text_parts.append(f"--- Sheet: {sheet_name} ---")
                text_parts.append(df.to_string())
                info.tables.append({"sheet": sheet_name, "rows": df.fillna("").astype(str).values.tolist()})
            info.raw_text = "\n".join(text_parts).strip()
            return
        except Exception:
            pass

        try:
            import xlrd

            wb = xlrd.open_workbook(str(path))
            lines: list[str] = []
            for sheet in wb.sheets():
                lines.append(f"--- Sheet: {sheet.name} ---")
                rows: list[list[str]] = []
                for row in range(sheet.nrows):
                    row_values = [str(sheet.cell_value(row, col)) for col in range(sheet.ncols)]
                    rows.append(row_values)
                    lines.append(" | ".join(row_values))
                info.tables.append({"sheet": sheet.name, "rows": rows})
            info.raw_text = "\n".join(lines).strip()
            return
        except Exception:
            pass

        self._parse_unknown(path, info)

    def _parse_ppt(self, path: Path, info: DocumentInfo) -> None:
        try:
            with tempfile.TemporaryDirectory() as tmpdir:
                result = subprocess.run(
                    ["soffice", "--headless", "--convert-to", "txt:Text", "--outdir", tmpdir, str(path)],
                    capture_output=True,
                    text=True,
                    timeout=60,
                )
                if result.returncode == 0:
                    txt_files = list(Path(tmpdir).glob("*.txt"))
                    if txt_files:
                        text = txt_files[0].read_text(encoding="utf-8", errors="ignore")
                        if text.strip():
                            info.raw_text = text.strip()
                            return
        except Exception:
            pass

        self._parse_unknown(path, info)

    def _parse_email(self, path: Path, info: DocumentInfo) -> None:
        text = ""
        suffix = path.suffix.lower()

        if suffix in {".eml", ".emlx"}:
            try:
                import email
                from email import policy

                raw = path.read_bytes()
                msg = email.message_from_bytes(raw, policy=policy.default)
                headers = [
                    f"From: {msg.get('From', '')}",
                    f"To: {msg.get('To', '')}",
                    f"Date: {msg.get('Date', '')}",
                    f"Subject: {msg.get('Subject', '')}",
                    "",
                ]
                text = "\n".join(headers)

                for part in msg.walk():
                    if part.get_content_type() == "text/plain":
                        payload = part.get_payload(decode=True)
                        if payload:
                            text += payload.decode("utf-8", errors="ignore") + "\n"
            except Exception:
                pass
        elif suffix == ".msg":
            try:
                import extract_msg

                msg = extract_msg.Message(str(path))
                text = (
                    f"From: {msg.sender}\nTo: {msg.to}\nDate: {msg.date}\n"
                    f"Subject: {msg.subject}\n\n{msg.body or ''}"
                )
            except Exception:
                pass

        if text.strip():
            info.raw_text = text.strip()
            return
        self._parse_unknown(path, info)

    def _parse_image_ocr(self, path: Path, info: DocumentInfo) -> None:
        try:
            from PIL import Image
            import pytesseract

            img = Image.open(str(path))
            text = pytesseract.image_to_string(img, lang="spa+eng")
            if text.strip():
                info.raw_text = text.strip()
                info.metadata["ocr_method"] = "pytesseract"
                return
        except Exception:
            pass

        try:
            import easyocr

            reader = easyocr.Reader(["es", "en"])
            result = reader.readtext(str(path))
            text = "\n".join(item[1] for item in result)
            if text.strip():
                info.raw_text = text.strip()
                info.metadata["ocr_method"] = "easyocr"
                return
        except Exception:
            pass

        info.error = "OCR not available"

    def _parse_odt(self, path: Path, info: DocumentInfo) -> None:
        try:
            from odf import text as odf_text, teletype
            from odf.opendocument import load

            doc = load(str(path))
            paragraphs = doc.getElementsByType(odf_text.P)
            text = "\n".join(teletype.extractText(p) for p in paragraphs)
            info.raw_text = text.strip()
            return
        except Exception:
            pass

        self._parse_unknown(path, info)

    def _parse_unknown(self, path: Path, info: DocumentInfo) -> None:
        encodings = ["utf-8", "latin-1", "cp1252", "iso-8859-1", "cp850"]
        for encoding in encodings:
            try:
                text = path.read_text(encoding=encoding, errors="ignore")
                if len(text.strip()) > 20:
                    text = re.sub(r"[^\x20-\x7E\x0A\x0D\xC0-\xFF\xA0-\xFF]", " ", text)
                    text = re.sub(r"\s+", " ", text)
                    if len(text.strip()) > 20:
                        info.raw_text = text.strip()
                        info.encoding = encoding
                        info.metadata["warning"] = "File type not recognized. Content extracted as raw text."
                        return
            except Exception:
                continue

        info.raw_text = ""
        info.error = f"Cannot read file: {path.suffix} format not supported and raw text extraction failed."


# ------------------------------------------------------------------
# Helpers
# ------------------------------------------------------------------

def _ext_to_type(ext: str) -> DocumentType:
    MAPPING = {
        ".pdf": DocumentType.PDF,
        ".docx": DocumentType.DOCX,
        ".doc": DocumentType.DOCX,
        ".docm": DocumentType.DOCX,
        ".rtf": DocumentType.DOCX,
        ".odt": DocumentType.DOCX,
        ".wpd": DocumentType.DOCX,
        ".wps": DocumentType.DOCX,
        ".pages": DocumentType.DOCX,
        ".xlsx": DocumentType.XLSX,
        ".xls": DocumentType.XLSX,
        ".xlsm": DocumentType.XLSX,
        ".xlsb": DocumentType.XLSX,
        ".ods": DocumentType.XLSX,
        ".numbers": DocumentType.XLSX,
        ".csv": DocumentType.CSV,
        ".tsv": DocumentType.CSV,
        ".tab": DocumentType.CSV,
        ".txt": DocumentType.TXT,
        ".md": DocumentType.TXT,
        ".rst": DocumentType.TXT,
        ".log": DocumentType.TXT,
        ".text": DocumentType.TXT,
        ".eml": DocumentType.TXT,
        ".msg": DocumentType.TXT,
        ".emlx": DocumentType.TXT,
        ".json": DocumentType.JSON,
        ".yaml": DocumentType.JSON,
        ".yml": DocumentType.JSON,
        ".xml": DocumentType.XML,
        ".html": DocumentType.HTML,
        ".htm": DocumentType.HTML,
        ".mht": DocumentType.HTML,
        ".mhtml": DocumentType.HTML,
        ".pptx": DocumentType.PPTX,
        ".ppt": DocumentType.PPTX,
        ".pptm": DocumentType.PPTX,
        ".odp": DocumentType.PPTX,
        ".key": DocumentType.PPTX,
    }
    return MAPPING.get(ext.lower(), DocumentType.UNKNOWN)


def _detect_encoding(path: Path) -> str:
    try:
        import chardet

        raw = path.read_bytes()
        result = chardet.detect(raw)
        return result.get("encoding") or "utf-8"
    except ImportError:
        return "utf-8"
