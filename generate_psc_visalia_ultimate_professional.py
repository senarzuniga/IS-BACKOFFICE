from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

# ============================================================
# PRESENTATION: PSC VISALIA PROJECT
# Customer: Pacific Southwest Container, LLC
# Brand: Ingecart + Linetex USA
# Design: Ingecart Booth 2026 + ingecart.eu
# Version: ULTIMATE PROFESSIONAL
# ============================================================

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ============================================================
# EXACT BRAND COLORS (from artwork + website)
# ============================================================
DARK_BLUE = RGBColor(0, 40, 80)       # Primary (#002850)
MEDIUM_BLUE = RGBColor(0, 112, 192)   # Secondary (#0070C0)
ORANGE = RGBColor(255, 140, 0)        # Accent (#FF8C00)
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(248, 248, 248)
MID_GRAY = RGBColor(100, 100, 100)
DARK_GRAY = RGBColor(60, 60, 60)
GREEN = RGBColor(0, 128, 0)           # For positive metrics

# ============================================================
# COMPANY STATISTICS (from ingecart.eu)
# ============================================================
STATS = {
    "projects": "1,268",
    "experience": "28",
    "agreements": "26",
    "installations": "194"
}

# ============================================================
# FINANCIAL DATA (from your Excel)
# ============================================================
FINANCIALS = {
    "bhs_corrugator": 950000,
    "services": 1200000,
    "peripherals": 1100000,
    "conveyors_amr": 1450000,
    "ingetrans": 1320000,
    "civil_work": 900000,
    "total": 6920000
}

# ============================================================
# CORE DESIGN FUNCTIONS
# ============================================================

def add_booth_background(slide, with_footer=True):
    """Professional white background with brand bars"""
    # White background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()
    bg.name = "Background"
    
    # Top blue bar (0.12" height)
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.12))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = MEDIUM_BLUE
    top_bar.line.fill.background()
    top_bar.name = "TopBar"
    
    # Bottom orange accent line (subtle)
    bottom_accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, prs.slide_height - Inches(0.06), prs.slide_width, Inches(0.04))
    bottom_accent.fill.solid()
    bottom_accent.fill.fore_color.rgb = ORANGE
    bottom_accent.line.fill.background()
    bottom_accent.name = "BottomAccent"
    
    if with_footer:
        # Footer line
        footer_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3), prs.slide_height - Inches(0.48), Inches(7.333), Inches(0.02))
        footer_line.fill.solid()
        footer_line.fill.fore_color.rgb = MEDIUM_BLUE
        footer_line.line.fill.background()
        footer_line.name = "FooterLine"
        
        # Footer text
        footer = slide.shapes.add_textbox(Inches(0.5), prs.slide_height - Inches(0.43), Inches(12.333), Inches(0.35))
        tf = footer.text_frame
        tf.text = "SOLUTIONS DESIGNED FOR CORRUGATED PLANTS"
        tf.paragraphs[0].font.size = Pt(9)
        tf.paragraphs[0].font.color.rgb = MID_GRAY
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        footer.name = "FooterText"

def add_title_slide(prs, main_title, subtitle, project_title=""):
    """Premium cover slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Dark blue background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BLUE
    bg.line.fill.background()
    
    # Top orange bar
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.12))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = ORANGE
    top_bar.line.fill.background()
    
    # Bottom blue bar
    bottom_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, prs.slide_height - Inches(0.08), prs.slide_width, Inches(0.08))
    bottom_bar.fill.solid()
    bottom_bar.fill.fore_color.rgb = MEDIUM_BLUE
    bottom_bar.line.fill.background()
    
    # Main title
    title = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(12.333), Inches(1.2))
    tf = title.text_frame
    tf.text = main_title
    tf.paragraphs[0].font.size = Pt(44)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Subtitle
    sub = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(12.333), Inches(0.6))
    sf = sub.text_frame
    sf.text = subtitle
    sf.paragraphs[0].font.size = Pt(22)
    sf.paragraphs[0].font.color.rgb = RGBColor(200, 200, 200)
    sf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    if project_title:
        proj = slide.shapes.add_textbox(Inches(0.5), Inches(4.4), Inches(12.333), Inches(0.5))
        pf = proj.text_frame
        pf.text = project_title
        pf.paragraphs[0].font.size = Pt(16)
        pf.paragraphs[0].font.color.rgb = ORANGE
        pf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Partners line
    partners = slide.shapes.add_textbox(Inches(0.5), Inches(6.3), Inches(12.333), Inches(0.4))
    lf = partners.text_frame
    lf.text = "INGECART  |  LINETEX USA  |  PACIFIC SOUTHWEST CONTAINER"
    lf.paragraphs[0].font.size = Pt(11)
    lf.paragraphs[0].font.color.rgb = RGBColor(150, 150, 150)
    lf.paragraphs[0].alignment = PP_ALIGN.CENTER

def add_section_header(slide, title, subtitle=""):
    """Elegant section header with orange accent"""
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.45), Inches(11), Inches(0.7))
    tf = title_box.text_frame
    tf.text = title
    tf.paragraphs[0].font.size = Pt(30)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = DARK_BLUE
    
    # Orange accent line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.05), Inches(2.0), Inches(0.06))
    line.fill.solid()
    line.fill.fore_color.rgb = ORANGE
    line.line.fill.background()
    
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12), Inches(0.4))
        sf = sub.text_frame
        sf.text = subtitle
        sf.paragraphs[0].font.size = Pt(14)
        sf.paragraphs[0].font.color.rgb = MID_GRAY

def add_bullet_slide(prs, title, bullets, highlight_bullet=None):
    """Standard content slide with professional bullet formatting"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_booth_background(slide)
    add_section_header(slide, title)
    
    content = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(12.333), Inches(5.2))
    tf = content.text_frame
    tf.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        p.text = bullet
        p.space_after = Pt(10)
        p.font.size = Pt(14)
        
        if bullet.startswith("•") or bullet.startswith("✓") or bullet.startswith("→"):
            p.font.color.rgb = DARK_GRAY
            p.level = 0
        elif bullet.startswith("BENEFITS:") or bullet.startswith("KEY") or bullet.startswith("INVESTMENT:"):
            p.font.size = Pt(15)
            p.font.bold = True
            p.font.color.rgb = MEDIUM_BLUE
        elif bullet.startswith('"'):
            p.font.size = Pt(16)
            p.font.italic = True
            p.font.color.rgb = ORANGE
        elif highlight_bullet and bullet == highlight_bullet:
            p.font.bold = True
            p.font.color.rgb = GREEN
        else:
            p.font.color.rgb = DARK_GRAY

def add_two_column_slide(prs, title, left_title, left_items, right_title, right_items):
    """Professional before/after comparison slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_booth_background(slide)
    add_section_header(slide, title)
    
    # Left column header
    left_header = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(5.8), Inches(0.5))
    lh = left_header.text_frame
    lh.text = left_title
    lh.paragraphs[0].font.size = Pt(16)
    lh.paragraphs[0].font.bold = True
    lh.paragraphs[0].font.color.rgb = MEDIUM_BLUE
    
    # Left column content
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(5.8), Inches(4.5))
    lf = left_box.text_frame
    lf.word_wrap = True
    for i, item in enumerate(left_items):
        p = lf.paragraphs[i] if i < len(lf.paragraphs) else lf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(12)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(6)
    
    # Right column header
    right_header = slide.shapes.add_textbox(Inches(7.0), Inches(1.6), Inches(5.8), Inches(0.5))
    rh = right_header.text_frame
    rh.text = right_title
    rh.paragraphs[0].font.size = Pt(16)
    rh.paragraphs[0].font.bold = True
    rh.paragraphs[0].font.color.rgb = GREEN
    
    # Right column content
    right_box = slide.shapes.add_textbox(Inches(7.0), Inches(2.2), Inches(5.8), Inches(4.5))
    rf = right_box.text_frame
    rf.word_wrap = True
    for i, item in enumerate(right_items):
        p = rf.paragraphs[i] if i < len(rf.paragraphs) else rf.add_paragraph()
        p.text = f"✓ {item}"
        p.font.size = Pt(12)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(6)

def add_pillars_slide(prs):
    """The 4 pillars of Ingecart – signature brand slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_booth_background(slide, with_footer=True)
    
    # Main tagline
    tagline = slide.shapes.add_textbox(Inches(0.5), Inches(0.7), Inches(12.333), Inches(0.7))
    tf = tagline.text_frame
    tf.text = "Solving Real Bottlenecks in Corrugated Plants"
    tf.paragraphs[0].font.size = Pt(28)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = DARK_BLUE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Subtagline
    subtag = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12.333), Inches(0.5))
    sf = subtag.text_frame
    sf.text = "Industrial Intelligence Applied to Performance"
    sf.paragraphs[0].font.size = Pt(16)
    sf.paragraphs[0].font.color.rgb = MID_GRAY
    sf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    pillars = [
        ("FASTER THAN OEMs", "Agile engineering and\nrapid implementation", "01"),
        ("MORE RELIABLE", "Robust solutions built\nto perform", "02"),
        ("MEASURABLE ROI", "Focus on what impacts\nyour OEE", "03"),
        ("KNOWLEDGE DRIVEN", "Every project makes\nus smarter", "04")
    ]
    
    positions = [(0.8, 2.4, 2.7, 1.9), (3.8, 2.4, 2.7, 1.9), 
                 (6.8, 2.4, 2.7, 1.9), (9.8, 2.4, 2.7, 1.9)]
    
    for i, (p_title, p_desc, num) in enumerate(pillars):
        # Card
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Inches(positions[i][0]), Inches(positions[i][1]),
                                      Inches(positions[i][2]), Inches(positions[i][3]))
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT_GRAY
        card.line.color.rgb = MEDIUM_BLUE
        card.line.width = Pt(1.5)
        
        # Number
        num_box = slide.shapes.add_textbox(Inches(positions[i][0] + 0.12), Inches(positions[i][1] + 0.1), 
                                           Inches(0.5), Inches(0.5))
        nf = num_box.text_frame
        nf.text = num
        nf.paragraphs[0].font.size = Pt(18)
        nf.paragraphs[0].font.bold = True
        nf.paragraphs[0].font.color.rgb = ORANGE
        
        # Title
        pt = slide.shapes.add_textbox(Inches(positions[i][0] + 0.65), Inches(positions[i][1] + 0.12),
                                      Inches(2.0), Inches(0.5))
        ptf = pt.text_frame
        ptf.text = p_title
        ptf.paragraphs[0].font.size = Pt(11)
        ptf.paragraphs[0].font.bold = True
        ptf.paragraphs[0].font.color.rgb = DARK_BLUE
        
        # Description
        pd = slide.shapes.add_textbox(Inches(positions[i][0] + 0.15), Inches(positions[i][1] + 0.7),
                                      Inches(2.4), Inches(1.0))
        pdf = pd.text_frame
        pdf.text = p_desc
        pdf.paragraphs[0].font.size = Pt(10)
        pdf.paragraphs[0].font.color.rgb = DARK_GRAY
        pdf.paragraphs[0].alignment = PP_ALIGN.CENTER

def add_stats_slide(prs):
    """Company statistics from website"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_booth_background(slide)
    add_section_header(slide, "Ingecart by the Numbers", "Proven track record in corrugated packaging")
    
    stats_items = [
        (STATS["projects"], "Projects Completed", 1.5),
        (STATS["experience"], "Years of Experience", 4.5),
        (STATS["agreements"], "International Agreements", 7.5),
        (STATS["installations"], "Installations", 10.5)
    ]
    
    for value, label, x in stats_items:
        # Value
        val = slide.shapes.add_textbox(Inches(x), Inches(2.5), Inches(1.8), Inches(0.9))
        vf = val.text_frame
        vf.text = value
        vf.paragraphs[0].font.size = Pt(48)
        vf.paragraphs[0].font.bold = True
        vf.paragraphs[0].font.color.rgb = ORANGE
        vf.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Label
        lbl = slide.shapes.add_textbox(Inches(x - 0.2), Inches(3.5), Inches(2.2), Inches(0.6))
        lf = lbl.text_frame
        lf.text = label
        lf.paragraphs[0].font.size = Pt(12)
        lf.paragraphs[0].font.color.rgb = DARK_GRAY
        lf.paragraphs[0].alignment = PP_ALIGN.CENTER

def add_closing_slide(prs):
    """Professional closing slide matching cover"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Dark blue background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BLUE
    bg.line.fill.background()
    
    # Top orange bar
    top = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.12))
    top.fill.solid()
    top.fill.fore_color.rgb = ORANGE
    top.line.fill.background()
    
    # Bottom blue bar
    bottom = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, prs.slide_height - Inches(0.08), prs.slide_width, Inches(0.08))
    bottom.fill.solid()
    bottom.fill.fore_color.rgb = MEDIUM_BLUE
    bottom.line.fill.background()
    
    # Main message
    msg = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(12.333), Inches(1.0))
    mf = msg.text_frame
    mf.text = "Solving Real Bottlenecks\nin Corrugated Plants"
    mf.paragraphs[0].font.size = Pt(42)
    mf.paragraphs[0].font.bold = True
    mf.paragraphs[0].font.color.rgb = WHITE
    mf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Sub message
    sub = slide.shapes.add_textbox(Inches(0.5), Inches(3.7), Inches(12.333), Inches(0.6))
    sf = sub.text_frame
    sf.text = "Industrial Intelligence Applied to Performance"
    sf.paragraphs[0].font.size = Pt(20)
    sf.paragraphs[0].font.color.rgb = RGBColor(200, 200, 200)
    sf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Thanks
    thanks = slide.shapes.add_textbox(Inches(0.5), Inches(4.8), Inches(12.333), Inches(0.6))
    tf = thanks.text_frame
    tf.text = "Thank You | Q&A"
    tf.paragraphs[0].font.size = Pt(26)
    tf.paragraphs[0].font.color.rgb = ORANGE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Contact
    contact = slide.shapes.add_textbox(Inches(0.5), Inches(6.3), Inches(12.333), Inches(0.6))
    cf = contact.text_frame
    cf.text = "Diego Garcia | diego@ingecart.eu | ingecart.eu\nLinetex USA | Local Support"
    cf.paragraphs[0].font.size = Pt(12)
    cf.paragraphs[0].font.color.rgb = RGBColor(180, 180, 180)
    cf.paragraphs[0].alignment = PP_ALIGN.CENTER

# ============================================================
# BUILD THE PRESENTATION – 20 SLIDES
# ============================================================

print("=" * 60)
print("BUILDING ULTIMATE PROFESSIONAL PRESENTATION")
print("For: Pacific Southwest Container, LLC")
print("=" * 60)

# Slide 1: Cover
add_title_slide(prs, 
    "Solving Real Bottlenecks\nin Corrugated Plants",
    "Industrial Intelligence Applied to Performance",
    "PSC VISALIA PROJECT | Complete Relocation + Modernization + Automation")

# Slide 2: The 4 Pillars (Signature Ingecart Slide)
add_pillars_slide(prs)

# Slide 3: Executive Summary
add_bullet_slide(prs, "Executive Summary – Turnkey Project", [
    "",
    "PROJECT OVERVIEW:",
    "• Relocate complete BHS corrugator (1997, Turkey → Visalia, California)",
    "• Complete obsolescence elimination (Siemens S5 → S7 unified architecture)",
    "• Full downstream automation: Ingetrans, Conveyors, AMR, WIP buffer",
    "",
    "KEY METRICS:",
    "• Total investment: $6.92M USD",
    "• Expected payback: <18 months",
    "• Speed increase: 250 → 300 mpm (+20%)",
    "• Downtime reduction: -70% from obsolescence elimination",
    "",
    '"Not just moving a machine – transforming the operation"'
])

# Slide 4: Strategic Necessity
add_bullet_slide(prs, "Strategic Necessity – Why This Project", [
    "",
    "ELIMINATE CRITICAL OBSOLESCENCE:",
    "• S5 PLC – No factory support since 2015",
    "• Critical spare parts unavailable globally",
    "• Weeks of unplanned downtime risk",
    "",
    "INCREASE PRODUCTIVITY & OEE:",
    "• Sustained 300 mpm vs current 250 mpm",
    "• Real-time diagnostics and predictive maintenance",
    "• Estimated 30% setup time reduction",
    "",
    "PREPARE FOR INDUSTRY 4.0:",
    "• Open Siemens S7 architecture",
    "• MES/ERP integration ready",
    "• Foundation for full AMR/WIP automation"
])

# Slide 5: Ingecart by the Numbers
add_stats_slide(prs)

# Slide 6: Obsolescence Elimination (Before/After)
add_two_column_slide(prs, "Obsolescence Elimination – Siemens S7 Migration",
    "BEFORE (Current State in Turkey)",
    ["Mixed S5/S7 control architecture",
     "No factory support for S5 PLCs",
     "DC drives with obsolete components",
     "Legacy HMI with limited diagnostics",
     "Spare parts unavailable globally",
     "High risk of extended downtime",
     "No remote diagnostics capability"],
    "AFTER (Installed in Visalia)",
    ["Unified Siemens S7 platform – 15+ year future",
     "Full factory support and updates",
     "Modern AC drives + motion control",
     "Touchscreen HMI + SCADA ready",
     "Global spare parts availability",
     "Predictive maintenance capable",
     "Remote diagnostics and support"])

# Slide 7: Project Scope – 4 Phases
add_bullet_slide(prs, "Project Scope – 4 Integrated Phases", [
    "",
    "PHASE 1: DISASSEMBLY (Turkey)",
    "→ Complete teardown + component tagging + preservation for transport",
    "",
    "PHASE 2: TRANSPORT (Turkey → Visalia)",
    "→ 11 Open top + 2 Flat rack + 16 HQ 40' containers",
    "→ Full packaging and export documentation",
    "",
    "PHASE 3: REINSTALLATION (Visalia)",
    "→ Mechanical assembly + laser alignment + peripheral connection",
    "",
    "PHASE 4: MODERNIZATION + INTEGRATION",
    "→ Siemens S7 + new drives + HMI + SCADA",
    "→ Ingetrans + Conveyors + AMR + WIP buffer"
])

# Slide 8: Ingetrans – Automated Roll Loading (from Excel)
add_bullet_slide(prs, "Ingetrans – Automated Roll Loading System", [
    "",
    "SYSTEM CAPABILITIES:",
    "• Automatic roll loading up to 3,700 mm diameter",
    "• Zero forklift intervention – improved safety",
    "• Integrated loading/unloading ramp system",
    "• Double transfer car + rails + catenary system",
    "• Complete electrical cabinets + proprietary software",
    "• Assembly + 4 x 40' containers transport included",
    "",
    f"INVESTMENT: ${FINANCIALS['ingetrans']:,} USD (sale price)",
    "",
    '"Reduced manual handling + increased operational safety"'
])

# Slide 9: Conveyors & AMR System (from Excel)
add_bullet_slide(prs, "Conveyors & AMR – Intelligent Material Flow", [
    "",
    "SYSTEM COMPONENTS:",
    f"• Double output line with integrated strappers → ${FINANCIALS['conveyors_amr'] - 550000:,}K",
    f"• AMR reception conveyor + waste movement → $550,000",
    "• Intelligent WIP buffer between corrugator and converting",
    "• Full integration with existing AAR stacker",
    "",
    "OPERATIONAL BENEFITS:",
    "• Continuous flow to converting area – no accumulation stoppages",
    "• Real-time production tracking",
    "• Reduced manual pallet movement",
    "",
    '"The corrugator never stops due to downstream accumulation"'
])

# Slide 10: Peripherals & Engineering Services (from Excel)
add_bullet_slide(prs, "Peripherals & Engineering Services", [
    "",
    "PERIPHERALS INCLUDED:",
    "• Flexamix adhesive kitchen",
    "• Complete boiler system + steam tubes",
    "• Full steam tube assembly and integration",
    "",
    "ENGINEERING SERVICES INCLUDED:",
    "• Complete machine disassembly (Turkey)",
    "• Full installation and commissioning (USA)",
    "• Electrical modernization to Siemens S7",
    "• Project management + on-site supervision",
    "",
    f"TOTAL SERVICES INVESTMENT: ${FINANCIALS['services']:,} USD",
    f"TOTAL PERIPHERALS: ${FINANCIALS['peripherals']:,} USD"
])

# Slide 11: Measurable ROI – OEE Impact
add_bullet_slide(prs, "Measurable ROI – Focus on OEE", [
    "",
    "QUANTIFIABLE BENEFITS:",
    "",
    "• Downtime reduction from obsolescence: -70%",
    "• Sustained speed increase: 250 → 300 mpm (+20%)",
    "• Setup time reduction: -30%",
    "• Unplanned stops from parts unavailability: ELIMINATED",
    "• Maintenance cost reduction: -25%",
    "• Energy efficiency improvement: -15% (new AC drives)",
    "",
    f"ESTIMATED PAYBACK: <18 MONTHS",
    "",
    f"TOTAL INVESTMENT: ${FINANCIALS['total']:,} USD",
    "",
    '"Focus on what impacts your OEE"'
])

# Slide 12: Project Timeline – 9 Months
add_bullet_slide(prs, "Project Timeline – ~9 Months", [
    "",
    "Engineering + FAT: ████████ 8 weeks",
    "→ Detailed design, software development, factory acceptance test",
    "",
    "Disassembly (Turkey): ████ 4 weeks",
    "→ On-site teardown, tagging, container loading",
    "",
    "Shipping (Turkey → Visalia): ██████ 6 weeks",
    "→ Ocean freight, customs clearance, delivery to site",
    "",
    "Mechanical Reinstallation: ████████ 8 weeks",
    "→ Assembly, laser alignment, mechanical integration",
    "",
    "Electrical Modernization + Software: ██████ 6 weeks",
    "→ S7 installation, drives, HMI, system commissioning",
    "",
    "Site Acceptance Test + Ramp-up: ████ 4 weeks",
    "→ Validation, training, production handover",
    "",
    "TOTAL ESTIMATED: 36 weeks (9 months)"
])

# Slide 13: Quality Assurance – FAT & SAT
add_bullet_slide(prs, "Quality Assurance – FAT & SAT Process", [
    "",
    "FACTORY ACCEPTANCE TEST (FAT):",
    "• Conducted at Ingecart partner facilities before shipping",
    "• Full software and control system validation",
    "• Customer witness opportunity available",
    "",
    "SITE ACCEPTANCE TEST (SAT):",
    "• Conducted at PSC Visalia after installation",
    "• Performance validation against agreed KPIs",
    "• Joint sign-off with customer operations team",
    "",
    "POST-COMMISSIONING SUPPORT:",
    "• 2 weeks intensive on-site support",
    "• Remote diagnostics and monitoring via secure VPN",
    "• 24/7 emergency response availability"
])

# Slide 14: Investment Summary (from Excel)
add_bullet_slide(prs, "Investment Summary – Total Cost Breakdown", [
    "",
    f"BHS CORRUGATOR (reused equipment): ${FINANCIALS['bhs_corrugator']:,} USD",
    f"SERVICES (disassembly, transport, installation, PM): ${FINANCIALS['services']:,} USD",
    f"PERIPHERALS (Flexamix, boiler, steam piping): ${FINANCIALS['peripherals']:,} USD",
    f"CONVEYORS + AMR SYSTEM: ${FINANCIALS['conveyors_amr']:,} USD",
    f"INGETRANS (automated roll loading): ${FINANCIALS['ingetrans']:,} USD",
    f"CIVIL WORK + PERMITS (estimated): ${FINANCIALS['civil_work']:,} USD",
    "",
    f"TOTAL ESTIMATED INVESTMENT: ${FINANCIALS['total']:,} USD",
    "",
    "PROPOSED PAYMENT STRUCTURE:",
    "• 30% with Purchase Order",
    "• 30% upon installation start",
    "• 30% upon successful commissioning",
    "• 10% upon SAT completion and handover"
])

# Slide 15: Why Ingecart + Linetex – Differentiators
add_bullet_slide(prs, "Why Ingecart + Linetex?", [
    "",
    "✓ REAL CORRUGATED EXPERIENCE",
    "→ Dozens of projects across Europe and USA",
    f"→ {STATS['installations']}+ corrugated plants in Europe",
    "",
    "✓ TRUE TURNKEY APPROACH",
    "→ Single point of responsibility from disassembly to production",
    "→ No finger-pointing between multiple vendors",
    "",
    "✓ PROVEN AMR INTEGRATION",
    "→ Not theoretical – we make it work in production environments",
    "",
    "✓ LOCAL USA SUPPORT",
    "→ Linetex as local agent + available field engineers",
    "→ Rapid response times – no overseas delays",
    "",
    "✓ VERIFIED REFERENCES",
    "→ International Paper, Saica, DS Smith, Smurfit WestRock, Hinojosa",
    "",
    '"Faster than OEMs · More reliable than low-cost"'
])

# Slide 16: Technical Specifications Summary
add_bullet_slide(prs, "Technical Specifications – Key Parameters", [
    "",
    "CORRUGATOR SPECIFICATIONS:",
    "• BHS Jets 300-2500II, 2500 mm working width",
    "• Maximum speed: 300 mpm after modernization",
    "• Flute capabilities: B, C, E configurations",
    "",
    "MODERNIZATION SCOPE:",
    "• Siemens S7-1500 series PLC",
    "• SINAMICS drives with motion control",
    "• SCADA-ready HMI with remote access",
    "",
    "AUTOMATION SCOPE:",
    "• Ingetrans: 3,700 mm roll capacity, fully automated",
    "• Conveyors: Double output line with strappers",
    "• AMR: Fleet-ready infrastructure",
    "",
    "UTILITIES REQUIRED:",
    "• Electrical: 480V, 3-phase, 60Hz",
    "• Steam: Integration with existing boiler system",
    "• Compressed air: Standard plant supply"
])

# Slide 17: Risk Mitigation & Project Guarantees
add_bullet_slide(prs, "Risk Mitigation & Project Guarantees", [
    "",
    "OBSOLESCENCE RISK:",
    "• Eliminated via Siemens S7 migration",
    "• 15+ years of guaranteed spare parts availability",
    "",
    "SCHEDULE RISK:",
    "• Dedicated project manager assigned",
    "• Weekly progress reporting",
    "• Contingency built into timeline",
    "",
    "PERFORMANCE RISK:",
    "• FAT before shipping ensures software quality",
    "• SAT with customer sign-off on KPIs",
    "",
    "OPERATIONAL RISK:",
    "• Comprehensive training for your technicians",
    "• 2 weeks on-site support after handover",
    "• 24/7 remote support available"
])

# Slide 18: Next Steps
add_bullet_slide(prs, "Next Steps – Path to Project Start", [
    "",
    "IMMEDIATE ACTIONS:",
    "",
    "1. Layout Validation",
    "   → Onsite or remote technical meeting at Visalia",
    "",
    "2. Final Scope Definition",
    "   → Detailed engineering confirmation",
    "",
    "3. Final Economic Proposal",
    "   → Updated with any scope changes, ready for signatures",
    "",
    "4. Engineering Start",
    "   → 8 weeks before scheduled disassembly",
    "",
    "5. Financing Plan",
    "   → Discussion if lease or extended payment needed",
    "",
    "CONTACT INFORMATION:",
    "",
    "Diego Garcia | Technical & Commercial Director | Ingecart",
    "Email: diego@ingecart.eu | Web: ingecart.eu",
    "",
    "Linetex USA | Local Agent & Contractual Support"
])

# Slide 19: References – Previous Success
add_bullet_slide(prs, "References – Proven Track Record", [
    "",
    "INTERNATIONAL PAPER (USA/Europe):",
    "• Multiple Ingetrans installations",
    "• Ongoing service contracts",
    "",
    "SAICA PACK (Spain):",
    "• AMR Systems deployment",
    "• Complete intralogistics integration",
    "",
    "DS SMITH (Europe):",
    "• Easy Pack installations across multiple plants",
    "• Process optimization consulting",
    "",
    "SMURFIT WESTROCK (Spain/USA):",
    "• Palletizer modernization",
    "• Control system upgrades",
    "",
    "HINOJOSA PACKAGING GROUP (Spain):",
    "• Retal scrap management systems",
    "• Ongoing maintenance support",
    "",
    '"Every project makes us smarter"'
])

# Slide 20: Closing
add_closing_slide(prs)

# ============================================================
# SAVE PRESENTATION
# ============================================================

output_path = r"C:\Users\Inaki Senar\Documents\GitHub\IS-BACKOFFICE\PSC_VISALIA_ULTIMATE_PROFESSIONAL.pptx"
prs.save(output_path)

print("\n" + "=" * 60)
print("✅ PRESENTATION GENERATED SUCCESSFULLY")
print("=" * 60)
print(f"\n📁 File: {output_path}")
print(f"📊 Total slides: {len(prs.slides)}")
print("\n🎨 DESIGN ELEMENTS:")
print("   • Dark Blue primary (#002850)")
print("   • Medium Blue secondary (#0070C0)")
print("   • Orange accent (#FF8C00)")
print("   • 4 Pillars slide (FASTER, MORE RELIABLE, MEASURABLE ROI, KNOWLEDGE DRIVEN)")
print("   • Consistent footer: 'SOLUTIONS DESIGNED FOR CORRUGATED PLANTS'")
print("   • Professional before/after comparison layout")
print("\n📝 CONTENT INCLUDED:")
print(f"   • Company stats: {STATS['projects']} projects, {STATS['experience']} years")
print(f"   • Financial data: ${FINANCIALS['total']:,} total investment")
print("   • ROI analysis: <18 months payback")
print("   • 9-month timeline with 6 phases")
print("   • Technical specifications")
print("   • Risk mitigation section")
print("   • References section")
print("\n🚀 READY FOR CUSTOMER PRESENTATION")
print("   Pacific Southwest Container, LLC")
print("=" * 60)
