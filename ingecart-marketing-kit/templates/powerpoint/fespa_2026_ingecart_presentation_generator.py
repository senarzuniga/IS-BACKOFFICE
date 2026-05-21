from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt


def add_title(slide, text, left=0.8, top=0.35, width=11.8, height=0.8):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = "Montserrat"
    run.font.size = Pt(34)
    run.font.bold = True
    run.font.color.rgb = RGBColor(244, 245, 247)


def add_subtitle(slide, text, left=0.8, top=1.15, width=11.8, height=0.5):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = "Inter"
    run.font.size = Pt(16)
    run.font.color.rgb = RGBColor(126, 132, 142)


def add_bullets(slide, bullets, left=0.95, top=1.9, width=11.2, height=4.8):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, txt in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = txt
        p.level = 0
        p.font.name = "Inter"
        p.font.size = Pt(23)
        p.font.color.rgb = RGBColor(244, 245, 247)
        p.space_after = Pt(12)


def add_footer(slide, text="INGECART | FESPA 2026"):
    box = slide.shapes.add_textbox(Inches(0.8), Inches(6.95), Inches(8.0), Inches(0.3))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = "Inter"
    run.font.size = Pt(11)
    run.font.color.rgb = RGBColor(126, 132, 142)


def add_orange_line(slide):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.62), Inches(1.4), Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(255, 106, 0)
    line.line.color.rgb = RGBColor(255, 106, 0)


def add_logo_if_exists(slide, logo_path: Path):
    if logo_path.exists():
        slide.shapes.add_picture(str(logo_path), Inches(10.95), Inches(0.2), height=Inches(0.75))


def add_cover_art_if_exists(slide, art_path: Path):
    if art_path.exists():
        slide.shapes.add_picture(str(art_path), Inches(7.9), Inches(1.8), width=Inches(5.2), height=Inches(4.9))


def new_dark_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(5, 7, 11)
    return slide


def build_fespa_kit_presentation(output_path: Path, logo_path: Path, art_path: Path) -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1: portada
    s1 = new_dark_slide(prs)
    add_title(s1, "Menos friccion operativa. Mas produccion util.")
    add_subtitle(s1, "Kit comercial de stand - FESPA 2026")
    add_orange_line(s1)
    add_cover_art_if_exists(s1, art_path)
    add_logo_if_exists(s1, logo_path)
    add_footer(s1)

    # Slide 2: dolor
    s2 = new_dark_slide(prs)
    add_title(s2, "El problema real en planta")
    add_subtitle(s2, "No falta tecnologia; falta integracion con impacto")
    add_orange_line(s2)
    add_bullets(s2, [
        "Paradas no planificadas que erosionan OEE",
        "Trafico interno y riesgo en zonas criticas",
        "Coste oculto de retal, merma y retrabajo",
        "Mucha data y poca accion operativa",
    ])
    add_logo_if_exists(s2, logo_path)
    add_footer(s2)

    # Slide 3: propuesta
    s3 = new_dark_slide(prs)
    add_title(s3, "Propuesta Ingecart")
    add_subtitle(s3, "Ingenieria independiente + automatizacion + intralogistica + inteligencia")
    add_orange_line(s3)
    add_bullets(s3, [
        "Disenamos sobre su realidad de planta, no sobre catalogo",
        "Implantamos por fases para minimizar riesgo y parada",
        "Medimos impacto en KPI: OEE, MTTR, energia, coste por tonelada",
        "Acompanamos desde auditoria hasta operacion estable",
    ])
    add_logo_if_exists(s3, logo_path)
    add_footer(s3)

    # Slide 4: soluciones
    s4 = new_dark_slide(prs)
    add_title(s4, "5 Soluciones Base para FESPA")
    add_subtitle(s4, "Final de linea e intralogistica con resultado medible")
    add_orange_line(s4)
    add_bullets(s4, [
        "Paletizer + Easy Pack",
        "Sistema Retal SR1400",
        "Ingetrans 2800",
        "AMR",
        "Ing_PRO (copiloto industrial)",
    ], top=2.0)
    add_logo_if_exists(s4, logo_path)
    add_footer(s4)

    # Slide 5: cierre
    s5 = new_dark_slide(prs)
    add_title(s5, "Siguiente paso")
    add_subtitle(s5, "Diagnostico express + ROI preliminar en 48h")
    add_orange_line(s5)
    add_bullets(s5, [
        "1) Identificamos su principal cuello de botella",
        "2) Proponemos la fase de mayor impacto y menor riesgo",
        "3) Definimos KPI objetivo y plan de implantacion",
        "Agende reunion tecnica post-feria",
    ])
    add_logo_if_exists(s5, logo_path)
    add_footer(s5)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    return output_path


if __name__ == "__main__":
    repo_root = Path(__file__).resolve().parents[3]
    output = repo_root / "ingecart-marketing-kit" / "assets" / "presentations" / "Ingecart_FESPA_2026_Stand_Deck.pptx"

    # Rutas que indico el usuario
    logo = Path(r"C:\Users\Inaki Senar\Documents\INGECART\MARKETING\LOGOS\ingeeniering.png")
    art = Path(r"C:\Users\Inaki Senar\Documents\INGECART\MARKETING\ARTWORK\Imagen Slogan principal Ingecart.png")

    saved = build_fespa_kit_presentation(output, logo, art)
    print(f"Presentacion generada: {saved}")
