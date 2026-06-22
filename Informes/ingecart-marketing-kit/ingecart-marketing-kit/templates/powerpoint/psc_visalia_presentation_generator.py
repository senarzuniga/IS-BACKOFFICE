#!/usr/bin/env python3
"""
PSC VISALIA Executive Presentation Generator
Converts markdown executive report into branded PowerPoint with Ingecart corporate identity.

Color Palette (Ingecart Corporate):
- Dark Blue Primary: #002850
- Medium Blue Secondary: #0070C0
- Orange Accent: #FF8C00
- Light Background: #F2F6FC
- Text Gray: #6B7A8F
- White: #FFFFFF
"""

from pathlib import Path
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN


class IngeCartPSCPresentation:
    """Generate branded PowerPoint for PSC VISALIA investment case."""
    
    # =========================================
    # INGECART CORPORATE COLORS
    # =========================================
    DARK_BLUE = RGBColor(0, 40, 80)        # #002850
    MEDIUM_BLUE = RGBColor(0, 112, 192)   # #0070C0
    ORANGE = RGBColor(255, 140, 0)        # #FF8C00
    LIGHT_BG = RGBColor(242, 246, 252)    # #F2F6FC
    TEXT_GRAY = RGBColor(107, 122, 143)   # #6B7A8F
    WHITE = RGBColor(255, 255, 255)
    BLACK = RGBColor(0, 0, 0)
    
    def __init__(self):
        self.prs = Presentation()
        # 16:9 aspect ratio
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        
    def add_title_slide(self, title: str, subtitle: str, meta: str = ""):
        """Add title slide with Ingecart branding."""
        blank_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_layout)
        
        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.DARK_BLUE
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(2.0), Inches(11.7), Inches(1.5)
        )
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        p = title_frame.paragraphs[0]
        p.text = title
        p.font.name = "Montserrat"
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = self.WHITE
        
        # Subtitle
        subtitle_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(3.7), Inches(11.7), Inches(2.0)
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.word_wrap = True
        p_sub = subtitle_frame.paragraphs[0]
        p_sub.text = subtitle
        p_sub.font.name = "Segoe UI"
        p_sub.font.size = Pt(28)
        p_sub.font.color.rgb = self.ORANGE
        
        # Meta
        if meta:
            meta_box = slide.shapes.add_textbox(
                Inches(0.8), Inches(6.5), Inches(11.7), Inches(0.8)
            )
            meta_frame = meta_box.text_frame
            p_meta = meta_frame.paragraphs[0]
            p_meta.text = meta
            p_meta.font.name = "Segoe UI"
            p_meta.font.size = Pt(14)
            p_meta.font.color.rgb = self.LIGHT_BG
        
        return slide
    
    def add_content_slide(self, title: str, content_items: list, use_bullets: bool = True):
        """Add content slide with header and bullet points."""
        blank_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_layout)
        
        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.LIGHT_BG
        
        # Header bar
        header_shape = slide.shapes.add_shape(
            1,  # rectangle
            Inches(0), Inches(0), Inches(13.333), Inches(1.0)
        )
        header_shape.fill.solid()
        header_shape.fill.fore_color.rgb = self.DARK_BLUE
        header_shape.line.color.rgb = self.DARK_BLUE
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(0.15), Inches(12), Inches(0.7)
        )
        title_frame = title_box.text_frame
        p_title = title_frame.paragraphs[0]
        p_title.text = title
        p_title.font.name = "Montserrat"
        p_title.font.size = Pt(32)
        p_title.font.bold = True
        p_title.font.color.rgb = self.WHITE
        
        # Orange accent line
        line = slide.shapes.add_shape(
            1,  # rectangle
            Inches(0.8), Inches(0.95), Inches(1.5), Inches(0.04)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = self.ORANGE
        line.line.color.rgb = self.ORANGE
        
        # Content
        content_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(1.3), Inches(11.7), Inches(5.8)
        )
        text_frame = content_box.text_frame
        text_frame.word_wrap = True
        
        for i, item in enumerate(content_items):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            if use_bullets:
                p.level = 0
                p.text = item
            else:
                p.text = item
            
            p.font.name = "Segoe UI"
            p.font.size = Pt(16)
            p.font.color.rgb = self.BLACK
            p.space_after = Pt(10)
        
        return slide
    
    def add_two_column_slide(self, title: str, left_items: list, right_items: list, 
                             left_title: str = "", right_title: str = ""):
        """Add two-column content slide."""
        blank_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_layout)
        
        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.LIGHT_BG
        
        # Header bar
        header_shape = slide.shapes.add_shape(
            1,  # rectangle
            Inches(0), Inches(0), Inches(13.333), Inches(1.0)
        )
        header_shape.fill.solid()
        header_shape.fill.fore_color.rgb = self.DARK_BLUE
        header_shape.line.color.rgb = self.DARK_BLUE
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(0.15), Inches(12), Inches(0.7)
        )
        title_frame = title_box.text_frame
        p_title = title_frame.paragraphs[0]
        p_title.text = title
        p_title.font.name = "Montserrat"
        p_title.font.size = Pt(32)
        p_title.font.bold = True
        p_title.font.color.rgb = self.WHITE
        
        # Orange accent line
        line = slide.shapes.add_shape(
            1,  # rectangle
            Inches(0.8), Inches(0.95), Inches(1.5), Inches(0.04)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = self.ORANGE
        line.line.color.rgb = self.ORANGE
        
        # LEFT COLUMN
        if left_title:
            left_title_box = slide.shapes.add_textbox(
                Inches(0.8), Inches(1.3), Inches(5.5), Inches(0.4)
            )
            left_title_frame = left_title_box.text_frame
            p_lt = left_title_frame.paragraphs[0]
            p_lt.text = left_title
            p_lt.font.name = "Montserrat"
            p_lt.font.size = Pt(16)
            p_lt.font.bold = True
            p_lt.font.color.rgb = self.DARK_BLUE
        
        left_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(1.8) if left_title else Inches(1.3), 
            Inches(5.5), Inches(5.3)
        )
        left_frame = left_box.text_frame
        left_frame.word_wrap = True
        
        for i, item in enumerate(left_items):
            if i == 0:
                p = left_frame.paragraphs[0]
            else:
                p = left_frame.add_paragraph()
            p.text = item
            p.font.name = "Segoe UI"
            p.font.size = Pt(14)
            p.font.color.rgb = self.BLACK
            p.space_after = Pt(8)
        
        # RIGHT COLUMN
        if right_title:
            right_title_box = slide.shapes.add_textbox(
                Inches(7.0), Inches(1.3), Inches(5.5), Inches(0.4)
            )
            right_title_frame = right_title_box.text_frame
            p_rt = right_title_frame.paragraphs[0]
            p_rt.text = right_title
            p_rt.font.name = "Montserrat"
            p_rt.font.size = Pt(16)
            p_rt.font.bold = True
            p_rt.font.color.rgb = self.DARK_BLUE
        
        right_box = slide.shapes.add_textbox(
            Inches(7.0), Inches(1.8) if right_title else Inches(1.3), 
            Inches(5.5), Inches(5.3)
        )
        right_frame = right_box.text_frame
        right_frame.word_wrap = True
        
        for i, item in enumerate(right_items):
            if i == 0:
                p = right_frame.paragraphs[0]
            else:
                p = right_frame.add_paragraph()
            p.text = item
            p.font.name = "Segoe UI"
            p.font.size = Pt(14)
            p.font.color.rgb = self.BLACK
            p.space_after = Pt(8)
        
        return slide
    
    def add_table_slide(self, title: str, rows: list, headers: list):
        """Add slide with simple table (3 columns max)."""
        blank_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_layout)
        
        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.LIGHT_BG
        
        # Header bar
        header_shape = slide.shapes.add_shape(
            1,  # rectangle
            Inches(0), Inches(0), Inches(13.333), Inches(1.0)
        )
        header_shape.fill.solid()
        header_shape.fill.fore_color.rgb = self.DARK_BLUE
        header_shape.line.color.rgb = self.DARK_BLUE
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(0.15), Inches(12), Inches(0.7)
        )
        title_frame = title_box.text_frame
        p_title = title_frame.paragraphs[0]
        p_title.text = title
        p_title.font.name = "Montserrat"
        p_title.font.size = Pt(32)
        p_title.font.bold = True
        p_title.font.color.rgb = self.WHITE
        
        # Orange accent line
        line = slide.shapes.add_shape(
            1,  # rectangle
            Inches(0.8), Inches(0.95), Inches(1.5), Inches(0.04)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = self.ORANGE
        line.line.color.rgb = self.ORANGE
        
        # Add table using text boxes (simpler than python-pptx tables)
        num_cols = len(headers)
        col_width = 11.7 / num_cols
        
        # Headers
        for col, header in enumerate(headers):
            header_box = slide.shapes.add_textbox(
                Inches(0.8 + col * col_width),
                Inches(1.3),
                Inches(col_width - 0.1),
                Inches(0.4)
            )
            header_frame = header_box.text_frame
            header_frame.word_wrap = True
            p_header = header_frame.paragraphs[0]
            p_header.text = header
            p_header.font.name = "Montserrat"
            p_header.font.size = Pt(13)
            p_header.font.bold = True
            p_header.font.color.rgb = self.DARK_BLUE
        
        # Rows
        for row_idx, row in enumerate(rows):
            for col_idx, cell in enumerate(row):
                cell_box = slide.shapes.add_textbox(
                    Inches(0.8 + col_idx * col_width),
                    Inches(1.8 + row_idx * 0.5),
                    Inches(col_width - 0.1),
                    Inches(0.45)
                )
                cell_frame = cell_box.text_frame
                cell_frame.word_wrap = True
                p_cell = cell_frame.paragraphs[0]
                p_cell.text = str(cell)
                p_cell.font.name = "Segoe UI"
                p_cell.font.size = Pt(11)
                p_cell.font.color.rgb = self.BLACK
        
        return slide
    
    def generate(self) -> Path:
        """Generate complete PSC VISALIA presentation."""
        
        # SLIDE 1: TITLE
        self.add_title_slide(
            "PSC VISALIA",
            "BHS JETS 300-2500 II\nStrategic Relocation & Modernization",
            "Investment USD 6.5M | May 2026 - Feb 2027"
        )
        
        # SLIDE 2: EXECUTIVE SUMMARY
        self.add_content_slide(
            "Decision: Proceed with Conditions",
            [
                "✓ Strong value case: capacity recovery + risk reduction + modernization",
                "✓ Phased execution model supports cash control and milestone tracking",
                "✓ Technology platform is AMR-ready and automation-ready",
                "",
                "⏳ GATES TO CLOSE:",
                "  • Technical: Site survey + utility confirmation + final layout",
                "  • Commercial: Reconcile CAPEX basis + freeze scope",
                "  • Execution: Secure BHS support letter + logistics + commissioning ownership"
            ]
        )
        
        # SLIDE 3: EXECUTIVE THESIS
        self.add_content_slide(
            "Not a Second-Hand Purchase",
            [
                "🔄 Recover and redeploy an industrial asset with known technical history",
                "🔧 Rebuild around current control and handling standards",
                "🛡️  Reduce manual traffic and operational risk at corrugator",
                "🚀 Create plant flow ready for future automation and serviceability",
                "",
                "This is controlled relocation + modernization = operational reset"
            ]
        )
        
        # SLIDE 4: WHAT PSC IS BUYING
        self.add_content_slide(
            "Value Layers in This Investment",
            [
                "1. Rebuilt corrugated line with proven BHS base",
                "2. Cleaner and safer handling model",
                "3. More stable flow between corrugator, WIP, and conversion",
                "4. Serviceable architecture for maintenance and future upgrades",
                "5. Commercial project framework (not isolated equipment buys)"
            ]
        )
        
        # SLIDE 5: PROJECT SCOPE
        self.add_content_slide(
            "Six Blocks of Scope",
            [
                "① Core Line: BHS corrugator 2500mm, duplex configuration",
                "② Services: Disassembly, transport, assembly, commissioning",
                "③ Peripherals: Glue kitchen, boiler, steam tubes, utilities",
                "④ Intralogistics: Conveyors, WIP handling, output flow redesign",
                "⑤ Ingetrans: Automated reel handling and internal transport",
                "⑥ Civil & Permits: Site preparation for installation and utilities"
            ]
        )
        
        # SLIDE 6: WHY THIS INVESTMENT MAKES SENSE
        self.add_two_column_slide(
            "Strategic Rationale",
            [
                "TECHNICAL DE-RISKING:",
                "• Modernize legacy base",
                "• Simplify controls",
                "• Reduce spare-parts complexity",
                "",
                "SAFETY & TRAFFIC:",
                "• Reduce forklift traffic",
                "• Decrease uncontrolled circulation",
                "• Lower reel handling risk"
            ],
            [
                "PLANT PRODUCTIVITY:",
                "• More real line capacity",
                "• Fewer interruptions",
                "• Better production rhythm",
                "",
                "SCALABILITY:",
                "• AMR-ready architecture",
                "• Automation-ready layout",
                "• Prevents locked design"
            ]
        )
        
        # SLIDE 7: INVESTMENT FIGURES
        self.add_content_slide(
            "Financial Foundation",
            [
                "OBSERVED FIGURES (to reconcile):",
                "  • Internal working sheet: USD 7,873,480 (with tax)",
                "  • Cost base: USD 5,439,600",
                "  • Earlier deck reference: ~USD 6.92M",
                "  • Target scope: USD 6.5M",
                "",
                "⚠️  MUST RECONCILE before final approval",
                "Commercial team must present ONE approved number + scope + payment structure"
            ]
        )
        
        # SLIDE 8: FINANCIAL MODEL
        self.add_content_slide(
            "Phased Execution for Cash Control",
            [
                "Advantage: Avoids single large lump-sum purchase",
                "Advantage: Ties payments to visible milestones",
                "",
                "Phases:",
                "  • Engineering → Disassembly → Transport",
                "  • Upgrade → Installation → Commissioning → Training",
                "",
                "This structure should be maintained in final agreement"
            ]
        )
        
        # SLIDE 9: MAIN BENEFITS
        self.add_content_slide(
            "Expected Benefits for PSC Visalia",
            [
                "✓ Lower manual traffic around corrugator area",
                "✓ Better control of reels, WIP, and line-out flow",
                "✓ Reduced dependency on obsolete equipment behavior",
                "✓ Cleaner maintenance and serviceability model",
                "✓ Stronger plant safety posture",
                "✓ Better scalability for future automation layers"
            ]
        )
        
        # SLIDE 10: MAIN RISKS (Part 1)
        self.add_content_slide(
            "Critical Risks to Close - Priority 1",
            [
                "🚨 Risk 1: Site survey missing or incomplete",
                "  → Impact: Layout confirmation delays entire timeline",
                "  → Mitigation: Commission survey by Q2 2026",
                "",
                "🚨 Risk 2: Utility capacity not confirmed",
                "  → Impact: Could require infrastructure upgrades",
                "  → Mitigation: Final electrical/steam/civil assessment now",
                "",
                "🚨 Risk 3: BHS support letter not formalized",
                "  → Impact: Post-install service at risk",
                "  → Mitigation: Obtain written BHS commitment before release"
            ]
        )
        
        # SLIDE 11: MAIN RISKS (Part 2)
        self.add_content_slide(
            "Critical Risks to Close - Priority 2",
            [
                "🚨 Risk 4: CAPEX figure not reconciled",
                "  → Impact: Approval gateway failure",
                "  → Mitigation: Align PDF, Excel, and user target by next review",
                "",
                "🚨 Risk 5: Baseline OEE not visible",
                "  → Impact: Cannot validate hard ROI claim",
                "  → Mitigation: Establish pre-install metrics for comparison",
                "",
                "Key: Address all gates before executive sign-off"
            ]
        )
        
        # SLIDE 12: APPROVAL PATH
        self.add_content_slide(
            "Three-Gate Approval Process",
            [
                "GATE 1: TECHNICAL VALIDATION",
                "  ✓ Site survey approved",
                "  ✓ Utility check completed",
                "  ✓ Final layout confirmed",
                "",
                "GATE 2: COMMERCIAL VALIDATION",
                "  ✓ Scope frozen",
                "  ✓ Investment amount reconciled",
                "  ✓ Payment schedule closed",
                "",
                "GATE 3: EXECUTION READINESS",
                "  ✓ BHS support formalized",
                "  ✓ Logistics plan documented",
                "  ✓ Commissioning ownership assigned"
            ]
        )
        
        # SLIDE 13: BOTTOM LINE
        self.add_content_slide(
            "Strategic Industrial Transformation",
            [
                "TREAT THIS AS: Strategic project (not commodity purchase)",
                "",
                "VALUE CASE IS STRONG because it combines:",
                "  • Capacity recovery from proven asset",
                "  • Risk reduction through modernization",
                "  • Layout modernization in one integrated move",
                "",
                "ONLY BLOCKER: Reconcile capex basis + validate site readiness",
                "",
                "RECOMMENDATION: Conditional approval pending three-gate closure"
            ]
        )
        
        # SLIDE 14: PROJECT SNAPSHOT (Reference)
        rows = [
            ["BHS corrugator 2500mm class", "Up to 300 m/min", "Turnkey + modernization"],
            ["Ingecart + Linetex (USA)", "May 2026 - Feb 2027", "9 months"]
        ]
        headers = ["Asset & Scope", "Performance", "Timeline"]
        self.add_table_slide("Project Snapshot", rows, headers)
        
        # SLIDE 15: NEXT STEPS
        self.add_content_slide(
            "Immediate Actions Required",
            [
                "Week 1-2:",
                "  • Commission site survey + schedule utility assessment",
                "  • Request formal BHS support letter + relocation terms",
                "",
                "Week 3-4:",
                "  • Reconcile CAPEX (align PDF, Excel, target)",
                "  • Freeze scope definition",
                "  • Establish payment schedule",
                "",
                "Week 5-6:",
                "  • Present reconciled case to Capex Committee",
                "  • Obtain three-gate approvals"
            ]
        )
        
        # SLIDE 16: CLOSING
        self.add_title_slide(
            "Ready for Approval",
            "Strong value case | Controlled execution | Clear pathway forward",
            "PSC VISALIA BHS 98 | Investment Decision 2026"
        )
        
        # Save
        output_path = Path(__file__).resolve().parents[3] / "informes" / "PSC_VISALIA_EXECUTIVE_PRESENTATION_2026-05-11.pptx"
        output_path.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(str(output_path))
        
        return output_path


if __name__ == "__main__":
    generator = IngeCartPSCPresentation()
    saved_to = generator.generate()
    print(f"✓ PowerPoint presentation created: {saved_to}")
    print(f"  Slides: 16 | Format: 16:9 | Branding: Ingecart Corporate")
