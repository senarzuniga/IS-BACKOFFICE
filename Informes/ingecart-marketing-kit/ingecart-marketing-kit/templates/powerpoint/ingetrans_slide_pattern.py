from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt


# =========================================
# REUSABLE PATTERN: INGETRANS CORPORATE SLIDE
# =========================================

def build_ingetrans_corporate_slide(output_path: Path) -> Path:
    prs = Presentation()

    # Slide size (16:9)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_layout = prs.slide_layouts[6]  # blank slide
    slide = prs.slides.add_slide(slide_layout)

    # =========================================
    # COLORS
    # =========================================
    BACKGROUND = RGBColor(244, 243, 241)  # #F4F3F1
    BLACK = RGBColor(21, 21, 21)  # #151515
    GREY = RGBColor(111, 115, 120)  # #6F7378
    ORANGE = RGBColor(255, 106, 0)  # #FF6A00

    # =========================================
    # BACKGROUND
    # =========================================
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BACKGROUND

    # =========================================
    # TITLE BLOCK
    # =========================================
    title_box = slide.shapes.add_textbox(
        Inches(1.0),
        Inches(0.8),
        Inches(11),
        Inches(1.5),
    )

    title_tf = title_box.text_frame
    title_tf.clear()

    p = title_tf.paragraphs[0]
    run = p.add_run()
    run.text = "More meters per minute.\nZero chaos at the corrugator."
    run.font.name = "Montserrat"
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = BLACK

    # =========================================
    # ORANGE DIVIDER LINE
    # =========================================
    line = slide.shapes.add_shape(
        1,  # rectangle
        Inches(1.0),
        Inches(2.1),
        Inches(1.3),
        Inches(0.03),
    )

    line.fill.solid()
    line.fill.fore_color.rgb = ORANGE
    line.line.color.rgb = ORANGE

    # =========================================
    # PRODUCT NAME + SUBTITLE
    # =========================================
    subtitle_box = slide.shapes.add_textbox(
        Inches(1.0),
        Inches(2.4),
        Inches(11),
        Inches(1.0),
    )

    subtitle_tf = subtitle_box.text_frame
    subtitle_tf.clear()

    p2 = subtitle_tf.paragraphs[0]

    run1 = p2.add_run()
    run1.text = "INGETRANS 2800"
    run1.font.name = "Montserrat"
    run1.font.size = Pt(24)
    run1.font.bold = True
    run1.font.color.rgb = ORANGE

    run2 = p2.add_run()
    run2.text = ": safe productivity for corrugated plants."
    run2.font.name = "Inter"
    run2.font.size = Pt(20)
    run2.font.color.rgb = BLACK

    # =========================================
    # BODY TEXT
    # =========================================
    body_box = slide.shapes.add_textbox(
        Inches(1.0),
        Inches(3.2),
        Inches(8.8),
        Inches(2.0),
    )

    body_tf = body_box.text_frame
    body_tf.word_wrap = True

    body_text = [
        "More real line capacity: continuous, production-synchronized supply.",
        "Automates reel movement, delivery, and return from warehouse to corrugator.",
        "Creates safer workspaces by eliminating uncontrolled circulation.",
    ]

    for i, txt in enumerate(body_text):
        if i == 0:
            paragraph = body_tf.paragraphs[0]
        else:
            paragraph = body_tf.add_paragraph()

        paragraph.text = txt
        paragraph.font.name = "Inter"
        paragraph.font.size = Pt(18)
        paragraph.font.color.rgb = GREY
        paragraph.space_after = Pt(12)

    # =========================================
    # FINAL CLAIM
    # =========================================
    claim_box = slide.shapes.add_textbox(
        Inches(1.0),
        Inches(5.7),
        Inches(10),
        Inches(1.2),
    )

    claim_tf = claim_box.text_frame
    claim_tf.clear()

    p3 = claim_tf.paragraphs[0]

    claim_run1 = p3.add_run()
    claim_run1.text = "The Corrugated Plant\nof the "
    claim_run1.font.name = "Montserrat"
    claim_run1.font.size = Pt(30)
    claim_run1.font.bold = True
    claim_run1.font.color.rgb = BLACK

    claim_run2 = p3.add_run()
    claim_run2.text = "Future"
    claim_run2.font.name = "Montserrat"
    claim_run2.font.size = Pt(30)
    claim_run2.font.bold = True
    claim_run2.font.color.rgb = ORANGE

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    return output_path


if __name__ == "__main__":
    project_root = Path(__file__).resolve().parents[2]
    default_output = project_root / "assets" / "presentations" / "Ingecart_Corporate_Slide.pptx"
    saved_to = build_ingetrans_corporate_slide(default_output)
    print(f"PowerPoint created successfully: {saved_to}")
