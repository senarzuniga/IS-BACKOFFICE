from __future__ import annotations

from pathlib import Path
from typing import Iterable

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt


REPO_ROOT = Path(__file__).resolve().parents[1]
KIT_DIR = REPO_ROOT / "ingecart-marketing-kit" / "fespa-2026-kit-contenidos"
DEST_DIR = Path(r"C:\Users\Inaki Senar\Documents\INGECART\MARKETING\CONTENT\CONTENIDOS INGECART FESPA 2026")
LOGO_PATH = DEST_DIR / "ingeeniering.png"
HERO_PATH = DEST_DIR / "imagen_slogan_principal_ingecart.png"

BLACK = RGBColor(5, 7, 11)
DARK = RGBColor(26, 29, 36)
ORANGE = RGBColor(255, 106, 0)
WHITE = RGBColor(244, 245, 247)
GREY = RGBColor(126, 132, 142)
LIGHT = RGBColor(232, 235, 240)


class BrandDeck:
    def __init__(self, title: str, subtitle: str):
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        self.title = title
        self.subtitle = subtitle

    def new_slide(self, dark: bool = True):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = BLACK if dark else WHITE
        return slide

    def _add_logo(self, slide):
        if LOGO_PATH.exists():
            slide.shapes.add_picture(str(LOGO_PATH), Inches(10.8), Inches(0.2), height=Inches(0.7))

    def _header(self, slide, title: str, subtitle: str | None = None, dark: bool = True):
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.0))
        bar.fill.solid()
        bar.fill.fore_color.rgb = BLACK if dark else WHITE
        bar.line.color.rgb = BLACK if dark else WHITE

        tbox = slide.shapes.add_textbox(Inches(0.8), Inches(0.12), Inches(10.6), Inches(0.45))
        tf = tbox.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = title
        run.font.name = "Montserrat"
        run.font.size = Pt(28)
        run.font.bold = True
        run.font.color.rgb = WHITE if dark else BLACK

        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.92), Inches(1.6), Inches(0.04))
        line.fill.solid()
        line.fill.fore_color.rgb = ORANGE
        line.line.color.rgb = ORANGE

        if subtitle:
            sbox = slide.shapes.add_textbox(Inches(0.85), Inches(1.18), Inches(11.5), Inches(0.32))
            sf = sbox.text_frame
            sf.clear()
            p2 = sf.paragraphs[0]
            run2 = p2.add_run()
            run2.text = subtitle
            run2.font.name = "Inter"
            run2.font.size = Pt(13)
            run2.font.color.rgb = GREY if dark else RGBColor(90, 90, 90)

        self._add_logo(slide)
        return slide

    def _cover(self, title: str, subtitle: str, kicker: str, with_hero: bool = True):
        slide = self.new_slide(dark=True)
        if with_hero and HERO_PATH.exists():
            slide.shapes.add_picture(str(HERO_PATH), Inches(7.7), Inches(1.45), width=Inches(5.15), height=Inches(5.25))

        tbox = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(6.4), Inches(2.1))
        tf = tbox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = title
        run.font.name = "Montserrat"
        run.font.size = Pt(28)
        run.font.bold = True
        run.font.color.rgb = WHITE

        sbox = slide.shapes.add_textbox(Inches(0.8), Inches(3.45), Inches(6.1), Inches(1.1))
        sf = sbox.text_frame
        sp = sf.paragraphs[0]
        run2 = sp.add_run()
        run2.text = subtitle
        run2.font.name = "Inter"
        run2.font.size = Pt(18)
        run2.font.color.rgb = ORANGE

        kbox = slide.shapes.add_textbox(Inches(0.8), Inches(5.9), Inches(6.1), Inches(0.6))
        kf = kbox.text_frame
        kp = kf.paragraphs[0]
        run3 = kp.add_run()
        run3.text = kicker
        run3.font.name = "Inter"
        run3.font.size = Pt(12)
        run3.font.color.rgb = GREY

        if LOGO_PATH.exists():
            slide.shapes.add_picture(str(LOGO_PATH), Inches(0.85), Inches(0.35), height=Inches(0.65))

        return slide

    def _bullet_slide(self, title: str, bullets: Iterable[str], subtitle: str | None = None):
        slide = self.new_slide(dark=True)
        self._header(slide, title, subtitle)
        box = slide.shapes.add_textbox(Inches(0.9), Inches(1.6), Inches(11.4), Inches(5.5))
        tf = box.text_frame
        tf.word_wrap = True
        tf.clear()
        for i, bullet in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = bullet
            p.level = 0
            p.font.name = "Inter"
            p.font.size = Pt(19)
            p.font.color.rgb = WHITE
            p.space_after = Pt(10)
        return slide

    def _two_column_slide(self, title: str, left_title: str, left: Iterable[str], right_title: str, right: Iterable[str]):
        slide = self.new_slide(dark=True)
        self._header(slide, title, "Uso interno y reuniones formales")

        for x, ttl, items in [(0.9, left_title, left), (7.0, right_title, right)]:
            hbox = slide.shapes.add_textbox(Inches(x), Inches(1.55), Inches(5.2), Inches(0.35))
            htf = hbox.text_frame
            htf.clear()
            hp = htf.paragraphs[0]
            hr = hp.add_run()
            hr.text = ttl
            hr.font.name = "Montserrat"
            hr.font.size = Pt(15)
            hr.font.bold = True
            hr.font.color.rgb = ORANGE

            bbox = slide.shapes.add_textbox(Inches(x), Inches(1.95), Inches(5.25), Inches(4.95))
            btf = bbox.text_frame
            btf.word_wrap = True
            btf.clear()
            for i, item in enumerate(items):
                p = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
                p.text = item
                p.font.name = "Inter"
                p.font.size = Pt(16)
                p.font.color.rgb = WHITE
                p.space_after = Pt(8)

        return slide

    def _kpi_slide(self, title: str, metrics: list[tuple[str, str]]):
        slide = self.new_slide(dark=True)
        self._header(slide, title, "Datos de referencia para reforzar credibilidad")
        w = 5.5
        positions = [(0.9, 1.8), (6.75, 1.8), (0.9, 4.0), (6.75, 4.0)]
        for (label, value), (x, y) in zip(metrics, positions):
            rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(1.6))
            rect.fill.solid()
            rect.fill.fore_color.rgb = DARK
            rect.line.color.rgb = GREY
            lbox = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.12), Inches(w - 0.4), Inches(0.35))
            tf1 = lbox.text_frame
            tf1.clear()
            p1 = tf1.paragraphs[0]
            r1 = p1.add_run()
            r1.text = label
            r1.font.name = "Inter"
            r1.font.size = Pt(12)
            r1.font.color.rgb = GREY
            vbox = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.45), Inches(w - 0.4), Inches(0.7))
            tf2 = vbox.text_frame
            tf2.clear()
            p2 = tf2.paragraphs[0]
            r2 = p2.add_run()
            r2.text = value
            r2.font.name = "Montserrat"
            r2.font.size = Pt(22)
            r2.font.bold = True
            r2.font.color.rgb = ORANGE
        return slide

    def save(self, path: Path):
        path.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(str(path))
        return path


def corporate_deck() -> Path:
    d = BrandDeck(
        "Ingecart",
        "Engineering & Auditing | Tailored Automation & Intralogistics",
    )
    d._cover(
        "Ingecart",
        "Engineering & Auditing | Tailored Automation & Intralogistics",
        "Presentacion oficial para iniciar reuniones formales en oficina",
    )
    d._bullet_slide(
        "Quien somos",
        [
            "Empresa de ingenieria industrial especializada en carton corrugado y plantas de packaging.",
            "28 anos de experiencia, 1.268 proyectos, 26 acuerdos internacionales y 194 instalaciones activas.",
            "Independencia tecnica: no representamos fabricantes; elegimos la mejor tecnologia para cada planta.",
        ],
        "Vision 360 de planta y responsabilidad sobre el resultado.",
    )
    d._bullet_slide(
        "Que problema resolvemos",
        [
            "Paradas no planificadas, merma, energia elevada y logistica interna ineficiente.",
            "Proyectos complejos sin unico interlocutor ni criterio tecnico claro.",
            "Necesidad de acelerar digitalizacion sin parar produccion.",
        ],
        "El problema no es falta de tecnologia; es integracion con ROI.",
    )
    d._two_column_slide(
        "Nuestra propuesta de valor",
        "Lo que hacemos",
        [
            "Auditoria tecnica independiente.",
            "Automatizacion e intralogistica a medida.",
            "Robotica, bobinas, retal y fin de linea.",
            "Ing_PRO: copiloto industrial con IA.",
        ],
        "Que gana el cliente",
        [
            "Menos paradas y menos riesgo operativo.",
            "Mayor productividad y menor coste por tonelada.",
            "Implantacion por fases con control de CAPEX.",
            "Acompanamiento hasta operacion estable.",
        ],
    )
    d._kpi_slide(
        "Cifras de credibilidad",
        [
            ("Proyectos ejecutados", "1.268"),
            ("Anios de trayectoria", "28"),
            ("Acuerdos internacionales", "26"),
            ("Instalaciones activas", "194"),
        ],
    )
    d._bullet_slide(
        "Como trabajamos",
        [
            "1. Diagnostico y auditoria tecnica.",
            "2. Diseño de solucion y roadmap por fases.",
            "3. Integracion, instalacion y puesta en marcha.",
            "4. Formacion, soporte y optimizacion continua.",
        ],
        "Proceso de menor riesgo para reuniones de direccion.",
    )
    d._bullet_slide(
        "Por que Ingecart",
        [
            "Independencia de fabricante.",
            "Experiencia especifica en corrugado y packaging.",
            "Capacidad de integrar hardware, software y operaciones.",
            "Resultados medibles: OEE, MTTR, energia, merma y seguridad.",
        ],
        "Un solo equipo para llevar una planta de la idea al arranque.",
    )
    d._bullet_slide(
        "Siguiente paso",
        [
            "Reunion inicial de 30 minutos.",
            "Lectura de pain points y priorizacion de oportunidades.",
            "Propuesta de diagnostico preliminar y siguiente accion.",
        ],
        "Ideal para reunion formal de presentacion en oficina.",
    )
    return d.save(DEST_DIR / "01_Ingecart_Oficial_Reunion_Formal.pptx")


def product_deck(name: str, title: str, subtitle: str, overview: list[str], value: list[str], metrics: list[tuple[str, str]], next_step: list[str], filename: str) -> Path:
    d = BrandDeck(title, subtitle)
    d._cover(title, subtitle, "Deck comercial de producto para reuniones y visitas tecnicas")
    d._bullet_slide(f"Que es {name}", overview, "Mensaje para presentacion de 30 segundos")
    d._bullet_slide(f"Valor para el cliente - {name}", value, "Beneficios de negocio y operacion")
    d._kpi_slide(f"Cifras clave - {name}", metrics)
    d._bullet_slide("Objecciones tipicas y respuesta", [
        "Es muy caro -> modelamos ROI sobre sus datos y decidimos con criterio.",
        "No quiero parar produccion -> implantamos por fases y ventanas controladas.",
        "Mi planta es distinta -> precisamente por eso adaptamos el layout y la integracion.",
    ], "Frases puente para cerrar sin friccion")
    d._bullet_slide("Siguiente paso", next_step, "Cierre comercial recomendado")
    return d.save(DEST_DIR / filename)


def ing_pro_deck() -> Path:
    return product_deck(
        "Ing_PRO",
        "Ing_PRO",
        "Copiloto industrial de Ingecart",
        [
            "Ing_PRO traduce datos complejos de planta en acciones automáticas y claras.",
            "Combina agentes de IA, soporte operacional y automatizacion de flujos.",
            "No es una dashboard: es un copiloto industrial desde el primer dia.",
        ],
        [
            "Reduce paradas no planificadas y acelera la respuesta del equipo.",
            "Crea ordenes de trabajo y prioriza acciones por impacto real.",
            "Se integra con CMMS, ERP, PLC y dashboards existentes.",
        ],
        [
            ("OEE", "+8 a +15 pp"),
            ("Paradas", "-20% a -35%"),
            ("MTTR", "-30% a -50%"),
            ("Implantacion", "2 semanas"),
        ],
        [
            "Proponer demo operacional de 5 minutos con un caso real de planta.",
            "Cerrar POC corto con KPI definido.",
            "Escalar a mantenimiento, calidad y operacion tras validacion.",
        ],
        "03_Ing_PRO_Industrial_Copilot.pptx",
    )


def in_getrans_deck() -> Path:
    return product_deck(
        "Ingetrans 2800",
        "Ingetrans 2800",
        "Transporte automatizado de bobinas para corrugadoras",
        [
            "Sistema de alimentacion, movimiento, entrega y retorno de bobinas en planta.",
            "Elimina la dependencia de carretillas elevadoras en zona critica.",
            "Diseñado para sincronizar logistica interna con la corrugadora.",
        ],
        [
            "Menos trafico de carretillas, menos riesgo y menos parada por abastecimiento.",
            "Integracion adaptada al layout real del cliente.",
            "Disponible con enfoque de implantacion por fases y minimo impacto.",
        ],
        [
            ("Capacidad", "Hasta 2.800 kg por bobina"),
            ("Beneficio", "-40% a -70% trafico de carretillas"),
            ("Entorno", "Compatible con corrugadoras existentes"),
            ("KPI", "ROI acelerado por menor downtime"),
        ],
        [
            "Solicitar pre-layout de flujo de bobinas.",
            "Validar layout y puntos de entrada/salida.",
            "Definir fase 1 con ROI preliminar y ventana de parada.",
        ],
        "02_Ingetrans_2800_Bobina_Transport.pptx",
    )


def retal_deck() -> Path:
    return product_deck(
        "Sistema Retal SR1400",
        "Sistema Retal SR1400",
        "Gestion integrada del desperdicio en corrugado",
        [
            "Sistema de recogida, transporte y gestion del retal producido en proceso.",
            "Reduce energia y elimina intervencion manual innecesaria.",
            "Diseñado para integrarse al layout y proceso de cada planta.",
        ],
        [
            "Menor coste energetico y mejor orden operativo.",
            "Menos acumulacion de scrap y mejor seguridad en planta.",
            "Contribuye a objetivos de sostenibilidad y eficiencia.",
        ],
        [
            ("Energia", "Hasta -93% vs alternativas"),
            ("Intervencion manual", "Drastica reduccion"),
            ("Aplicacion", "Troquelado, impresion, laminado"),
            ("ROI", "Primer ano en caso tipo"),
        ],
        [
            "Hacer benchmark de consumo actual de retal.",
            "Medir puntos de generacion y salida.",
            "Definir instalacion por fases con ahorro cuantificado.",
        ],
        "04_Sistema_Retal_SR1400.pptx",
    )


def palletizer_deck() -> Path:
    return product_deck(
        "Paletizer + Easy Pack",
        "Paletizer + Easy Pack",
        "Fin de linea automatizado y estable",
        [
            "Solucion de paletizado y embalado para fin de linea con cadencia constante.",
            "Reduce error de apilado, variabilidad manual y coste de retrabajo.",
            "Pensado para clientes que necesitan un ultimo paso de produccion mas fiable.",
        ],
        [
            "Mejora la estabilidad de expedicion y la imagen de salida.",
            "Reduce trabajo manual en la fase final de la linea.",
            "Permite explicar ROI con ahorro de personal, tiempo y material.",
        ],
        [
            ("Cadencia", "Hasta 15 paquetes/min"),
            ("Carga", "240 kg"),
            ("Precision", "+/- 2 mm"),
            ("Altura", "> 2.500 mm"),
        ],
        [
            "Pedir mix de producto y volumen diario.",
            "Validar capacidad actual de paletizado manual.",
            "Proponer simulacion de ROI y ruta de implantacion.",
        ],
        "05_Paletizer_EasyPack_Fin_Line.pptx",
    )


def amr_deck() -> Path:
    return product_deck(
        "AMR Industrial",
        "AMR Industrial",
        "Movilidad autonoma para plantas flexibles",
        [
            "Robot autonomo que mueve cargas sin rieles ni marcas fijas.",
            "Reduce dependencia de sistemas rigidos y facilita cambios de layout.",
            "Ideal para operaciones con variacion de rutas o expansion por fases.",
        ],
        [
            "Instalacion rapida y reconfiguracion sencilla.",
            "Mejor adaptacion a cambios de planta o demanda.",
            "Escalable por flota y compatible con un futuro automatizado.",
        ],
        [
            ("Infraestructura", "Sin obra fija"),
            ("Flexibilidad", "Reubica rutas en horas"),
            ("Uso", "Logistica interna multipunto"),
            ("KPI", "Menor tiempo de cambio"),
        ],
        [
            "Identificar puntos de carga/descarga y rutas actuales.",
            "Diseñar un mapa inicial y casos de uso priorizados.",
            "Lanzar piloto con KPI de tiempos y seguridad.",
        ],
        "06_AMR_Industrial_Flex_Transport.pptx",
    )


def build_all():
    DEST_DIR.mkdir(parents=True, exist_ok=True)
    files = []
    files.append(corporate_deck())
    files.append(palletizer_deck())
    files.append(retal_deck())
    files.append(in_getrans_deck())
    files.append(amr_deck())
    files.append(ing_pro_deck())

    index = DEST_DIR / "11_INDICE_PRESENTACIONES_Ingecart_FESPA2026.txt"
    index.write_text(
        "PRESENTACIONES GENERADAS\n"
        "========================\n\n"
        "01_Ingecart_Oficial_Reunion_Formal.pptx\n"
        "02_Ingetrans_2800_Bobina_Transport.pptx\n"
        "03_Ing_PRO_Industrial_Copilot.pptx\n"
        "04_Sistema_Retal_SR1400.pptx\n"
        "05_Paletizer_EasyPack_Fin_Line.pptx\n"
        "06_AMR_Industrial_Flex_Transport.pptx\n",
        encoding="utf-8",
    )
    return files


if __name__ == "__main__":
    generated = build_all()
    print("Generadas presentaciones:")
    for path in generated:
        print(path)
