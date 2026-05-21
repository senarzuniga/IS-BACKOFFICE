from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

BASE_DIR = Path(r"C:\Users\Inaki Senar\Documents\GitHub\IS-BACKOFFICE")
TEMPLATE_PATH = BASE_DIR / "ingecart-marketing-kit" / "assets" / "presentations" / "Ingecart_Corporate_Slide.pptx"
OUTPUT_PATH = Path(r"C:\Users\Inaki Senar\Documents\INGECART\MARKETING\CONTENT\ING_PRO_PRESENTACION_ESTRATEGICA.pptx")

RED = RGBColor(204, 0, 0)
DARK = RGBColor(35, 35, 35)
GRAY = RGBColor(95, 95, 95)
LIGHT = RGBColor(245, 245, 245)
SOFT_RED = RGBColor(247, 233, 233)
SOFT_BLUE = RGBColor(234, 241, 250)
SOFT_GREEN = RGBColor(233, 247, 238)
WHITE = RGBColor(255, 255, 255)

TITLE_FONT = "Montserrat"
BODY_FONT = "Arial"


def remove_slide(prs: Presentation, idx: int) -> None:
    sld = prs.slides._sldIdLst[idx]
    prs.part.drop_rel(sld.rId)
    del prs.slides._sldIdLst[idx]


def clear_presentation(prs: Presentation) -> None:
    for i in range(len(prs.slides) - 1, -1, -1):
        remove_slide(prs, i)


def set_bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, shape_type, left, top, width, height, fill, line=None):
    shp = slide.shapes.add_shape(shape_type, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line if line is not None else fill
    return shp


def add_text(
    slide,
    left,
    top,
    width,
    height,
    text,
    *,
    font=BODY_FONT,
    size=14,
    color=DARK,
    bold=False,
    italic=False,
    align=PP_ALIGN.LEFT,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.bold = bold
        r.font.italic = italic
    return box


def add_header(slide, title: str, subtitle: str = ""):
    add_text(slide, Inches(0.55), Inches(0.3), Inches(8.6), Inches(0.62), title, font=TITLE_FONT, size=24, bold=True)
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(0.95), Inches(1.3), Inches(0.06), RED)
    if subtitle:
        add_text(slide, Inches(0.57), Inches(1.05), Inches(8.8), Inches(0.35), subtitle, size=11, color=GRAY)


def add_note(slide, note: str):
    slide.notes_slide.notes_text_frame.text = note


def title_slide(prs: Presentation):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, DARK)
    add_shape(s, MSO_SHAPE.RECTANGLE, Inches(9.0), Inches(0), Inches(0.6), Inches(7.5), RED)
    add_shape(s, MSO_SHAPE.OVAL, Inches(7.55), Inches(0.65), Inches(1.7), Inches(1.7), RED)
    add_text(s, Inches(0.8), Inches(1.35), Inches(7.2), Inches(1.5), "Ing_PRO", font=TITLE_FONT, size=44, color=WHITE, bold=True)
    add_text(s, Inches(0.82), Inches(2.95), Inches(7.2), Inches(1.1), "El Copiloto Industrial de Ingecart\nDe la señal al impacto operativo en segundos", font=BODY_FONT, size=20, color=WHITE)
    add_text(s, Inches(0.82), Inches(5.05), Inches(7.8), Inches(0.45), '"From Insight to Action. On the factory floor. On day one."', size=14, color=WHITE, italic=True)
    add_text(s, Inches(7.2), Inches(6.7), Inches(2.0), Inches(0.25), "Mayo 2026", size=10, color=WHITE, align=PP_ALIGN.RIGHT)
    add_note(s, "Hoy vamos a mostrar por qué Ing_PRO no es una dashboard más: es un copiloto industrial accionable que convierte datos dispersos en decisiones y acciones automáticas con impacto medible en OEE, MTTR, merma y coste energético.")


def problem_slide(prs: Presentation):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, WHITE)
    add_header(s, "El problema operativo que más cuesta", "Muchas plantas digitalizan datos, pero no digitalizan decisiones.")

    cards = [
        ("Paradas no planificadas", "Pérdida de OEE y urgencias de mantenimiento."),
        ("Datos dispersos", "SCADA, ERP, CMMS y sensores sin una lógica única."),
        ("Mantenimiento reactivo", "Técnicos saturados analizando y no resolviendo."),
        ("Adopción lenta", "Herramientas complejas que la operación no usa."),
    ]
    pos = [(Inches(0.75), Inches(1.75)), (Inches(5.0), Inches(1.75)), (Inches(0.75), Inches(3.45)), (Inches(5.0), Inches(3.45))]
    for (t, b), (x, y) in zip(cards, pos):
        add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(4.05), Inches(1.4), LIGHT, RGBColor(220, 220, 220))
        add_text(s, x + Inches(0.2), y + Inches(0.2), Inches(3.65), Inches(0.3), t, font=TITLE_FONT, size=14, bold=True)
        add_text(s, x + Inches(0.2), y + Inches(0.62), Inches(3.65), Inches(0.5), b, size=11, color=GRAY)

    add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.2), Inches(5.45), Inches(5.2), Inches(0.6), RED)
    add_text(s, Inches(2.35), Inches(5.6), Inches(4.9), Inches(0.25), "Conclusión: sin acción automática, el dato no crea valor.", font=TITLE_FONT, size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_note(s, "El dolor principal no es falta de datos; es la brecha entre detectar una anomalía y ejecutar la acción correcta a tiempo. Ahí es donde Ing_PRO cambia el juego.")


def what_is_slide(prs: Presentation):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, WHITE)
    add_header(s, "¿Qué es Ing_PRO?", "Sistema de soporte a producción y mantenimiento orientado a acción.")

    add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(1.8), Inches(8.0), Inches(1.2), SOFT_BLUE, RGBColor(210, 220, 235))
    add_text(s, Inches(1.2), Inches(2.05), Inches(7.4), Inches(0.7), "Ing_PRO traduce datos complejos de planta en acciones automáticas y claras.\nNo es otra dashboard: es un copiloto industrial accionable.", font=TITLE_FONT, size=17, bold=True, align=PP_ALIGN.CENTER)

    bullets = [
        "Detección de anomalías en tiempo real",
        "Siguiente mejor acción para operario y jefe de turno",
        "Creación automática de órdenes en CMMS",
        "Ajustes automáticos de parámetros cuando aplica",
        "Despliegue en 2 semanas sobre infraestructura existente",
    ]
    y = 3.5
    for item in bullets:
        add_text(s, Inches(1.0), Inches(y), Inches(7.8), Inches(0.35), "• " + item, size=13, color=DARK)
        y += 0.52

    add_note(s, "La definición comercial correcta es clara: Ing_PRO es copiloto industrial de acción, no una capa visual. Su propuesta es operativa y medible desde las primeras semanas.")


def architecture_slide(prs: Presentation):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, WHITE)
    add_header(s, "Arquitectura de valor en 3 capas")

    xs = [Inches(0.75), Inches(3.55), Inches(6.35)]
    titles = ["Ingeniería Ingecart", "Agentes IA", "Automatización de flujos"]
    bodies = [
        "Conocimiento profundo de corrugadoras, automatización y puntos críticos de fallo.",
        "Operacional, Mantenimiento, Optimización, Comercial, Management y más.",
        "CMMS, alertas, escalación y ajustes en tiempo real con trazabilidad.",
    ]
    fills = [SOFT_BLUE, SOFT_GREEN, SOFT_RED]
    for x, t, b, f in zip(xs, titles, bodies, fills):
        add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.0), Inches(2.45), Inches(3.6), f, RGBColor(220, 220, 220))
        add_text(s, x + Inches(0.14), Inches(2.2), Inches(2.15), Inches(0.5), t, font=TITLE_FONT, size=16, bold=True, align=PP_ALIGN.CENTER)
        add_shape(s, MSO_SHAPE.RECTANGLE, x + Inches(0.25), Inches(2.95), Inches(1.95), Inches(0.05), RED)
        add_text(s, x + Inches(0.14), Inches(3.15), Inches(2.15), Inches(1.8), b, size=11, color=GRAY, align=PP_ALIGN.CENTER)

    add_text(s, Inches(2.2), Inches(6.0), Inches(5.2), Inches(0.3), "La ventaja estratégica surge de combinar proceso + IA + ejecución automática.", size=11, color=GRAY, italic=True, align=PP_ALIGN.CENTER)
    add_note(s, "Ninguna capa por sí sola crea la propuesta. Lo diferencial es la combinación de ingeniería aplicada, agentes especializados y automatización de ejecución.")


def agents_slide(prs: Presentation):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, WHITE)
    add_header(s, "Capacidades clave: Agentes especializados")

    agent_data = [
        ("Operacional", "Detecta anomalías de presión, vibración, temperatura y merma."),
        ("Mantenimiento", "Predice fallos y prioriza órdenes por impacto en OEE."),
        ("Optimización", "Propone ajustes de energía, material y throughput."),
        ("Management", "Resume riesgos, desvíos y decisiones ejecutivas diarias."),
        ("Comercial", "Identifica señales de upsell y oportunidades de mejora."),
        ("Customer Success", "Seguimiento de adopción, uso y valor conseguido."),
    ]

    x_left, x_right = Inches(0.75), Inches(5.0)
    y = Inches(1.75)
    for idx, (name, desc) in enumerate(agent_data):
        x = x_left if idx % 2 == 0 else x_right
        add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(4.1), Inches(0.95), LIGHT, RGBColor(225, 225, 225))
        add_text(s, x + Inches(0.18), y + Inches(0.12), Inches(3.65), Inches(0.25), name, font=TITLE_FONT, size=13, bold=True)
        add_text(s, x + Inches(0.18), y + Inches(0.42), Inches(3.75), Inches(0.35), desc, size=10, color=GRAY)
        if idx % 2 == 1:
            y += Inches(1.08)

    add_note(s, "Ing_PRO no depende de un único modelo. Opera con agentes especializados por rol y objetivo operativo, lo que acelera adopción y precisión en planta.")


def flow_slide(prs: Presentation):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, WHITE)
    add_header(s, "De anomalía a impacto: flujo operativo")

    steps = [
        ("1. Detección", "Sensor identifica evento en tiempo real."),
        ("2. Análisis", "IA cruza histórico, contexto y criticidad."),
        ("3. Acción", "Crea OT, alerta al operario y escala si no hay respuesta."),
        ("4. Impacto", "Parada evitada, OEE protegido, aprendizaje continuo."),
    ]
    x = Inches(0.9)
    for i, (t, b) in enumerate(steps):
        add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.4), Inches(1.95), Inches(2.2), SOFT_BLUE if i % 2 == 0 else SOFT_GREEN, RGBColor(220, 220, 220))
        add_text(s, x + Inches(0.12), Inches(2.65), Inches(1.7), Inches(0.35), t, font=TITLE_FONT, size=13, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, x + Inches(0.12), Inches(3.1), Inches(1.7), Inches(0.95), b, size=10, color=GRAY, align=PP_ALIGN.CENTER)
        if i < 3:
            add_shape(s, MSO_SHAPE.RIGHT_ARROW, x + Inches(1.95), Inches(3.08), Inches(0.5), Inches(0.32), RED)
        x += Inches(2.15)

    add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.15), Inches(5.4), Inches(5.2), Inches(0.6), RED)
    add_text(s, Inches(2.3), Inches(5.57), Inches(4.9), Inches(0.2), "Clave: Ing_PRO cierra el ciclo completo en segundos.", font=TITLE_FONT, size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_note(s, "El flujo demuestra por qué Ing_PRO es estratégico: convierte detección en acción y acción en aprendizaje, creando un sistema que mejora cada semana.")


def integration_slide(prs: Presentation):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, WHITE)
    add_header(s, "Integración técnica sin fricción", "Compatible con ecosistemas existentes de planta.")

    blocks = [
        ("Entrada", "PLCs, SCADA, sensores IoT\nModbus, OPC-UA, MQTT"),
        ("Núcleo Ing_PRO", "Agentes IA + reglas de negocio\nMotor de priorización por OEE"),
        ("Salida", "CMMS (SAP, Maximo, Infor)\nERP, alertas, dashboards"),
    ]
    x = Inches(0.8)
    for i, (t, b) in enumerate(blocks):
        add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.25), Inches(2.55), Inches(2.35), LIGHT, RGBColor(220, 220, 220))
        add_text(s, x + Inches(0.18), Inches(2.48), Inches(2.15), Inches(0.35), t, font=TITLE_FONT, size=14, bold=True, align=PP_ALIGN.CENTER)
        add_shape(s, MSO_SHAPE.RECTANGLE, x + Inches(0.25), Inches(2.9), Inches(2.05), Inches(0.05), RED)
        add_text(s, x + Inches(0.2), Inches(3.08), Inches(2.1), Inches(1.1), b, size=10, color=GRAY, align=PP_ALIGN.CENTER)
        if i < 2:
            add_shape(s, MSO_SHAPE.CHEVRON, x + Inches(2.55), Inches(3.05), Inches(0.45), Inches(0.45), RED)
        x += Inches(3.0)

    add_text(s, Inches(1.1), Inches(5.3), Inches(7.9), Inches(0.5), "No requiere sustituir la planta. Se monta sobre infraestructura existente y acelera resultados desde la semana 2.", size=12, color=DARK, align=PP_ALIGN.CENTER)
    add_note(s, "La barrera de entrada tecnológica es baja porque Ing_PRO se conecta a lo que la planta ya tiene. Esa capacidad de integración rápida es un factor decisivo en ventas.")


def kpi_slide(prs: Presentation):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, WHITE)
    add_header(s, "KPIs documentados y mejoras típicas")

    kpis = [
        ("+8 a +15 pp", "OEE de línea"),
        ("-20% a -35%", "Paradas no planificadas"),
        ("-30% a -50%", "MTTR"),
        ("-8% a -20%", "Merma / defectos"),
        ("-5% a -12%", "Consumo energético"),
        ("+25%", "Disponibilidad del equipo de mantenimiento"),
    ]
    xs = [Inches(0.75), Inches(3.55), Inches(6.35)]
    ys = [Inches(1.85), Inches(3.75)]
    idx = 0
    for y in ys:
        for x in xs:
            big, label = kpis[idx]
            add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(2.55), Inches(1.55), SOFT_RED if idx % 2 == 0 else LIGHT, RGBColor(220, 220, 220))
            add_text(s, x + Inches(0.15), y + Inches(0.2), Inches(2.2), Inches(0.5), big, font=TITLE_FONT, size=22, color=RED, bold=True, align=PP_ALIGN.CENTER)
            add_text(s, x + Inches(0.15), y + Inches(0.85), Inches(2.2), Inches(0.35), label, size=11, color=GRAY, align=PP_ALIGN.CENTER)
            idx += 1

    add_note(s, "La conversación con dirección debe centrarse en impacto financiero y operativo. Estos indicadores permiten demostrar valor en semanas y defender escalado en comité.")


def role_value_slide(prs: Presentation):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, WHITE)
    add_header(s, "Valor por rol en la empresa")

    role_data = [
        ("Operario", "Instrucciones claras y accionables en tiempo real."),
        ("Jefe de Mantenimiento", "Órdenes priorizadas y menos tiempo en diagnóstico."),
        ("Plant Manager", "Visibilidad diaria de OEE, paradas evitadas y riesgos."),
        ("Director de Operaciones", "Tendencias multi-línea y decisiones con impacto de negocio."),
    ]
    y = Inches(1.85)
    for i, (r, v) in enumerate(role_data):
        fill = SOFT_BLUE if i % 2 == 0 else SOFT_GREEN
        add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.85), y, Inches(8.35), Inches(0.95), fill, RGBColor(220, 220, 220))
        add_text(s, Inches(1.1), y + Inches(0.18), Inches(2.1), Inches(0.25), r, font=TITLE_FONT, size=13, bold=True)
        add_text(s, Inches(3.0), y + Inches(0.18), Inches(5.8), Inches(0.45), v, size=11, color=GRAY)
        y += Inches(1.1)

    add_note(s, "Ing_PRO acelera adopción porque aterriza valor por perfil. Cada rol recibe utilidad práctica en su día a día, reduciendo la resistencia al cambio.")


def case_study_slide(prs: Presentation):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, WHITE)
    add_header(s, "Caso de uso: parada evitada en corrugadora")

    add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(8.4), Inches(3.8), LIGHT, RGBColor(220, 220, 220))
    add_text(s, Inches(1.1), Inches(2.1), Inches(7.8), Inches(0.45), "Situación inicial", font=TITLE_FONT, size=15, bold=True)
    add_text(s, Inches(1.1), Inches(2.55), Inches(7.8), Inches(0.4), "Corrugadora de alta velocidad con OEE en 82% y riesgo de fallo en rodillo crítico.", size=12, color=GRAY)
    add_text(s, Inches(1.1), Inches(3.05), Inches(7.8), Inches(0.45), "Acción Ing_PRO", font=TITLE_FONT, size=15, bold=True)
    add_text(s, Inches(1.1), Inches(3.5), Inches(7.8), Inches(0.6), "Detección de vibración anómala 6 horas antes del fallo, creación automática de OT y ejecución preventiva.", size=12, color=GRAY)
    add_text(s, Inches(1.1), Inches(4.18), Inches(7.8), Inches(0.45), "Impacto", font=TITLE_FONT, size=15, bold=True)
    add_text(s, Inches(1.1), Inches(4.62), Inches(7.8), Inches(0.6), "Parada evitada, +2% OEE mensual y ahorro estimado de 12.000 USD por evento.", size=12, color=GRAY)

    add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.35), Inches(5.75), Inches(5.2), Inches(0.55), RED)
    add_text(s, Inches(2.55), Inches(5.9), Inches(4.8), Inches(0.2), "La IA crea valor cuando reduce eventos críticos reales.", font=TITLE_FONT, size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_note(s, "Este tipo de narrativa conecta con dirección: evento concreto, acción concreta, resultado concreto. Ing_PRO debe venderse con casos y no solo con funcionalidades.")


def comparison_slide(prs: Presentation):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, WHITE)
    add_header(s, "Ing_PRO vs alternativas")

    table = s.shapes.add_table(7, 4, Inches(0.55), Inches(1.65), Inches(8.9), Inches(4.7)).table
    table.columns[0].width = Inches(2.25)
    table.columns[1].width = Inches(2.2)
    table.columns[2].width = Inches(2.2)
    table.columns[3].width = Inches(2.25)

    headers = ["Aspecto", "Dashboard genérica", "CMMS tradicional", "Ing_PRO"]
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = RED if c == 3 else DARK
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        for r in p.runs:
            r.font.name = TITLE_FONT
            r.font.size = Pt(11)
            r.font.bold = True
            r.font.color.rgb = WHITE

    rows = [
        ("Detecta anomalías", "No", "No", "Sí, automático"),
        ("Sugiere acciones", "No", "No", "Sí, accionable"),
        ("Automatiza tareas", "No", "Parcial", "Sí: OT, alertas y ajustes"),
        ("Tiempo de puesta en marcha", "Meses", "Meses", "2 semanas"),
        ("Interfaz para operación", "Compleja", "Compleja", "Simple y guiada"),
        ("ROI comprobable", "Difícil", "Difícil", "Visible en semanas"),
    ]
    for r_idx, row in enumerate(rows, start=1):
        for c_idx, value in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = value
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if r_idx % 2 else LIGHT
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            for r in p.runs:
                r.font.name = BODY_FONT
                r.font.size = Pt(10)
                r.font.color.rgb = DARK if c_idx == 3 else GRAY

    add_note(s, "Comparar claramente con alternativas ayuda a cerrar oportunidades. Ing_PRO gana por acción automática, rapidez de implantación y resultados medibles.")


def implementation_slide(prs: Presentation):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, WHITE)
    add_header(s, "Modelo de implantación y escalado")

    phases = [
        ("Semana 1-2", "Despliegue rápido", "Conectores, agentes base, integración CMMS, formación inicial."),
        ("Semana 3-4", "Afinación", "Ajuste de umbrales, validación con operación, primeros impactos KPI."),
        ("Mes 2+", "Optimización continua", "Aprendizaje de agentes, expansión por líneas, nuevas integraciones."),
    ]
    y = Inches(1.95)
    fills = [SOFT_BLUE, SOFT_GREEN, SOFT_RED]
    for (time, phase, desc), fill in zip(phases, fills):
        add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), y, Inches(8.3), Inches(1.35), fill, RGBColor(220, 220, 220))
        add_text(s, Inches(1.05), y + Inches(0.18), Inches(1.4), Inches(0.25), time, font=TITLE_FONT, size=12, color=RED, bold=True)
        add_text(s, Inches(2.4), y + Inches(0.18), Inches(2.2), Inches(0.3), phase, font=TITLE_FONT, size=14, bold=True)
        add_text(s, Inches(2.4), y + Inches(0.55), Inches(6.35), Inches(0.45), desc, size=11, color=GRAY)
        y += Inches(1.53)

    add_note(s, "La venta estratégica de Ing_PRO combina entrada rápida y escalado progresivo. Esto baja riesgo percibido y facilita aprobación financiera.")


def strategic_slide(prs: Presentation):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, WHITE)
    add_header(s, "Por qué Ing_PRO es estratégico para la empresa")

    bullets = [
        "Protege margen operativo reduciendo paradas y merma.",
        "Mejora confiabilidad de entrega y servicio al cliente.",
        "Reduce dependencia de decisiones reactivas por turno.",
        "Convierte mantenimiento en ventaja competitiva y no en centro de coste.",
        "Aporta trazabilidad para dirección y comités de inversión.",
        "Crea base de escalado digital sin reemplazar activos existentes.",
    ]
    y = Inches(1.9)
    for b in bullets:
        add_shape(s, MSO_SHAPE.OVAL, Inches(0.95), y + Inches(0.07), Inches(0.2), Inches(0.2), RED)
        add_text(s, Inches(1.25), y, Inches(7.9), Inches(0.35), b, size=13, color=DARK)
        y += Inches(0.63)

    add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.3), Inches(5.85), Inches(7.0), Inches(0.5), RED)
    add_text(s, Inches(1.5), Inches(5.98), Inches(6.6), Inches(0.2), "Ing_PRO conecta operación diaria con estrategia de crecimiento sostenible.", font=TITLE_FONT, size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_note(s, "La tesis estratégica es simple: Ing_PRO no es solo eficiencia operativa; es una palanca de competitividad, resiliencia y escalado rentable.")


def roadmap_slide(prs: Presentation):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, WHITE)
    add_header(s, "Roadmap Ing_PRO 2026-2027")

    milestones = [
        ("Q2 2026", "Predicción de vida útil de componentes críticos"),
        ("Q3 2026", "Gemelo digital para simulación de parámetros"),
        ("Q4 2026", "Reportes automáticos con IA generativa"),
        ("Q1 2027", "Optimización energética con huella de carbono"),
        ("Q2 2027", "Agente comercial avanzado para upsell industrial"),
    ]
    x = Inches(0.9)
    for i, (q, d) in enumerate(milestones):
        add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.6), Inches(1.65), Inches(2.5), LIGHT, RGBColor(220, 220, 220))
        add_text(s, x + Inches(0.12), Inches(2.85), Inches(1.4), Inches(0.3), q, font=TITLE_FONT, size=12, color=RED, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, x + Inches(0.1), Inches(3.25), Inches(1.45), Inches(1.45), d, size=10, color=GRAY, align=PP_ALIGN.CENTER)
        if i < len(milestones) - 1:
            add_shape(s, MSO_SHAPE.CHEVRON, x + Inches(1.66), Inches(3.55), Inches(0.3), Inches(0.32), RED)
        x += Inches(1.74)

    add_note(s, "Este roadmap posiciona Ing_PRO como plataforma viva y evolutiva. Es relevante para clientes que buscan socios de largo plazo, no soluciones puntuales.")


def closing_slide(prs: Presentation):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, DARK)
    add_text(s, Inches(2.1), Inches(1.6), Inches(5.4), Inches(1.0), "Ing_PRO", font=TITLE_FONT, size=40, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(1.7), Inches(2.65), Inches(6.2), Inches(0.8), "Inteligencia industrial que actúa desde el primer día.", font=TITLE_FONT, size=20, color=WHITE, align=PP_ALIGN.CENTER)
    add_shape(s, MSO_SHAPE.RECTANGLE, Inches(2.8), Inches(3.6), Inches(3.95), Inches(0.06), RED)
    add_text(s, Inches(1.4), Inches(4.1), Inches(6.9), Inches(0.55), "Siguiente paso recomendado: Workshop de diagnóstico + piloto de 4 semanas", size=14, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(1.7), Inches(5.15), Inches(6.2), Inches(0.4), "hablemos@ingecart.eu | +34 938 183 316 | www.ingecart.eu/ing-pro", size=12, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(1.9), Inches(6.2), Inches(5.8), Inches(0.4), '"From Insight to Action. On the factory floor. On day one."', size=12, color=WHITE, italic=True, align=PP_ALIGN.CENTER)
    add_note(s, "Cierre con llamada a la acción: propuesta concreta, bajo riesgo y foco en resultados rápidos. El objetivo es pasar de presentación a oportunidad comercial en la misma reunión.")


def build() -> Presentation:
    prs = Presentation(str(TEMPLATE_PATH))
    clear_presentation(prs)

    title_slide(prs)
    problem_slide(prs)
    what_is_slide(prs)
    architecture_slide(prs)
    agents_slide(prs)
    flow_slide(prs)
    integration_slide(prs)
    kpi_slide(prs)
    role_value_slide(prs)
    case_study_slide(prs)
    comparison_slide(prs)
    implementation_slide(prs)
    strategic_slide(prs)
    roadmap_slide(prs)
    closing_slide(prs)

    return prs


if __name__ == "__main__":
    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    presentation = build()
    presentation.save(str(OUTPUT_PATH))
    print(f"Presentation saved to {OUTPUT_PATH}")
