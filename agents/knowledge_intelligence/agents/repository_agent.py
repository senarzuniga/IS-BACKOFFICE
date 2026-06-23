"""
Agente 2: Repository Knowledge Agent
- Busca información en archivos locales (PDF, Word, Excel, emails, etc.)
- Prioriza conocimiento interno antes de buscar fuera
"""

import os
from pathlib import Path
from datetime import datetime
from typing import List, Dict, Any


class RepositoryKnowledgeAgent:
    def __init__(self, search_paths: List[str], memory):
        self.search_paths = [Path(p) for p in search_paths]
        self.memory = memory
        self.extensions = {
            '.pdf': self._parse_pdf,
            '.docx': self._parse_docx,
            '.xlsx': self._parse_excel,
            '.xls': self._parse_excel,
            '.txt': self._parse_text,
            '.md': self._parse_text,
            '.csv': self._parse_text,
            '.eml': self._parse_email,
            '.msg': self._parse_email,
            '.pptx': self._parse_pptx,
            '.ppt': self._parse_pptx,
        }
    
    def run(self, context: Dict[str, Any]) -> Dict[str, Any]:
        research_plan = context.get('research_plan', {})
        target = research_plan.get('target_company', '')
        areas = research_plan.get('areas', [])
        keywords = [target] + areas + ['INGECART', 'corrugado', 'corrugadora', 'cartón', 'automatización']
        
        results = []
        for base_path in self.search_paths:
            if not base_path.exists():
                continue
            for file_path in base_path.rglob("*"):
                if file_path.suffix.lower() in self.extensions:
                    try:
                        content = self.extensions[file_path.suffix.lower()](file_path)
                        if content and self._is_relevant(content, keywords):
                            results.append({
                                'file_path': str(file_path),
                                'name': file_path.name,
                                'content': content[:2000],
                                'type': file_path.suffix.lower()[1:],
                                'modified': datetime.fromtimestamp(file_path.stat().st_mtime).isoformat()
                            })
                    except Exception as e:
                        print(f"Error leyendo {file_path}: {e}")
        context['repository_knowledge'] = results
        return {'found_items': len(results), 'items': results[:50]}
    
    def _is_relevant(self, content: str, keywords: List[str]) -> bool:
        content_lower = content.lower()
        return any(kw.lower() in content_lower for kw in keywords if kw)
    
    def _parse_pdf(self, path):
        try:
            import PyPDF2
            with open(path, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                text = ''.join([page.extract_text() or '' for page in reader.pages[:20]])
                return text[:10000]
        except Exception:
            return ""
    
    def _parse_docx(self, path):
        try:
            from docx import Document
            doc = Document(path)
            return '\n'.join([p.text for p in doc.paragraphs[:100]])
        except Exception:
            return ""
    
    def _parse_excel(self, path):
        try:
            import pandas as pd
            df = pd.read_excel(path, nrows=50)
            return df.to_string()
        except Exception:
            return ""
    
    def _parse_text(self, path):
        try:
            with open(path, 'r', encoding='utf-8') as f:
                return f.read()[:10000]
        except Exception:
            try:
                with open(path, 'r', encoding='latin-1') as f:
                    return f.read()[:10000]
            except Exception:
                return ""
    
    def _parse_email(self, path):
        try:
            import email
            with open(path, 'rb') as f:
                msg = email.message_from_binary_file(f)
                subject = msg.get('Subject', '')
                body = ""
                if msg.is_multipart():
                    for part in msg.walk():
                        if part.get_content_type() == 'text/plain':
                            body += part.get_payload(decode=True).decode('utf-8', errors='ignore')
                else:
                    body = msg.get_payload(decode=True).decode('utf-8', errors='ignore')
                return f"Asunto: {subject}\n\n{body[:5000]}"
        except Exception:
            return ""
    
    def _parse_pptx(self, path):
        try:
            from pptx import Presentation
            prs = Presentation(path)
            text = []
            for slide in prs.slides[:20]:
                for shape in slide.shapes:
                    if hasattr(shape, 'text'):
                        text.append(shape.text)
            return '\n'.join(text)
        except Exception:
            return ""
