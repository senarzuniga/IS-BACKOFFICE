from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Datos de ejercicios por unidad (resumido y formateado)
ejercicios = [
    {
        "unidad": "Unidad 1: Welcome!",
        "ejercicios": [
            "Completa: I ___ a student. | She ___ happy. | ___ you at school?",
            "Elige la opción correcta: He (is / are) my friend. | They (am / are) at home.",
            "Corrige el error: She are my sister. → ______",
            "Traduce: Yo soy profesor. → ______"
        ]
    },
    {
        "unidad": "Unidad 2: My World",
        "ejercicios": [
            "Completa: She ___ (have got) a bike. | My brother ___ (be) tall.",
            "Elige la opción correcta: We (has / have) got a dog. | This is (my / mine) house.",
            "Traduce: Ellos tienen dos gatos. → ______"
        ]
    },
    {
        "unidad": "Unidad 3: Day by Day",
        "ejercicios": [
            "Completa: I ___ (not/play) tennis. | ___ she (go) to school?",
            "Elige la opción correcta: He (always / never) eats vegetables. | We go to school (at / on) Monday.",
            "Traduce: Yo nunca como pescado. → ______"
        ]
    },
    {
        "unidad": "Unidad 4: Out and About",
        "ejercicios": [
            "Completa: I ___ (can) swim. | ___ (imperative) your homework!",
            "Elige la opción correcta: There (is / are) some apples. | Is there (some / any) milk?",
            "Traduce: ¿Puedes ayudarme? → ______"
        ]
    },
    {
        "unidad": "Unidad 5A: Food, Glorious Food",
        "ejercicios": [
            "Completa: I would ___ a pizza. | There isn’t ___ milk.",
            "Elige la opción correcta: I (like / likes) apples. | Is there (some / any) bread?",
            "Traduce: Me gusta el helado. → ______"
        ]
    }
]

prs = Presentation()

# Portada
slide = prs.slides.add_slide(prs.slide_layouts[5])
title_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(2))
tf = title_box.text_frame
tf.text = "REPASO DE INGLÉS ESO 1\nEjercicios por Unidad"
tf.paragraphs[0].font.size = Pt(44)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = RGBColor(0, 51, 153)
tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# Instrucciones
slide = prs.slides.add_slide(prs.slide_layouts[5])
tf = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(2)).text_frame
tf.text = "Instrucciones:\n- Haz los ejercicios en orden, de fácil a difícil.\n- Marca los que te resulten complicados.\n- Usa colores, subraya o dibuja para ayudarte.\n- ¡Repasa y diviértete aprendiendo!"
tf.paragraphs[0].font.size = Pt(28)
tf.paragraphs[0].font.color.rgb = RGBColor(0, 102, 0)

# Ejercicios por unidad
theme_colors = [RGBColor(0,51,153), RGBColor(0,102,0), RGBColor(204,102,0), RGBColor(153,0,51), RGBColor(0,153,153)]
for idx, bloque in enumerate(ejercicios):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    color = theme_colors[idx % len(theme_colors)]
    # Título de unidad
    tf = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1)).text_frame
    p = tf.add_paragraph()
    p.text = bloque["unidad"]
    p.font.size = Pt(38)
    p.font.bold = True
    p.font.color.rgb = color
    # Ejercicios
    for i, ejercicio in enumerate(bloque["ejercicios"]):
        tf2 = slide.shapes.add_textbox(Inches(1), Inches(1.5 + i*1.1), Inches(8), Inches(1))
        p2 = tf2.text_frame.add_paragraph()
        p2.text = ejercicio
        p2.font.size = Pt(28)
        p2.font.color.rgb = RGBColor(40,40,40)
        p2.space_after = Pt(16)

# Guardar
prs.save("informes/1 ESO/REPASO_EJERCICIOS_PRESENTACION.pptx")
print("Presentación generada: informes/1 ESO/REPASO_EJERCICIOS_PRESENTACION.pptx")
