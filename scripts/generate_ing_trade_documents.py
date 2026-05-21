from __future__ import annotations

import copy
import csv
import json
import shutil
from pathlib import Path

from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Pt, RGBColor as DocxRGBColor
from pptx.dml.color import RGBColor as PptxRGBColor
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt as PptPt


REPO_ROOT = Path(__file__).resolve().parents[1]
REPORTS_DIR = REPO_ROOT / "reports"
TEMPLATES_DIR = REPORTS_DIR / "templates"
TEMPLATE_PPT = Path(r"C:\Users\Inaki Senar\Documents\INGECART\MARKETING\CONTENT\POWER POINT TEMPLATE ING_TRADE.pptx")
CSV_PATH = REPO_ROOT / "research" / "ingecart" / "machine_trading_company_view.csv"
BRAND_JSON = REPO_ROOT / "ingecart_assets_fusion" / "ingecart_branding.json"

PRESENTATION_NAME = "Ing_TRADE_Business_Presentation.pptx"
PRESENTATION_V2_NAME = "Ing_TRADE_Business_Presentation - version 2.pptx"
INGECART_TEMPLATE_NAME = "INGECART_Standard_Presentation_Template.potx"
INGECART_TEMPLATE_V2_NAME = "INGECART_Standard_Presentation_Template - version 2.potx"
ING_TRADE_TEMPLATE_NAME = "ING_TRADE_Presentation_Template.potx"
ING_TRADE_TEMPLATE_V2_NAME = "ING_TRADE_Presentation_Template - version 2.potx"

INGECART_OFFER_DOC = "INGECART_Offer_Template.docx"
INGECART_OFFER_DOC_V2 = "INGECART_Offer_Template - version 2.docx"
ING_TRADE_OFFER_DOC = "ING_TRADE_Offer_Template.docx"
ING_TRADE_OFFER_DOC_V2 = "ING_TRADE_Offer_Template - version 2.docx"


DEFAULT_BLUE = (17, 109, 255)
DEFAULT_DARK = (0, 0, 0)
DEFAULT_LIGHT = (255, 255, 255)


def load_branding() -> tuple[tuple[int, int, int], tuple[int, int, int], str, str]:
    primary = DEFAULT_DARK
    accent = DEFAULT_BLUE
    tagline = "Solving Real Bottlenecks in Corrugated Plants"
    subtagline = "Industrial Intelligence Applied to Performance"

    if BRAND_JSON.exists():
        data = json.loads(BRAND_JSON.read_text(encoding="utf-8"))
        p = data.get("primary_color") or list(DEFAULT_DARK)
        s = data.get("secondary_color") or list(DEFAULT_LIGHT)
        primary = tuple(p[:3]) if len(p) >= 3 else DEFAULT_DARK
        if s and len(s) >= 3 and tuple(s[:3]) != DEFAULT_LIGHT:
            accent = tuple(s[:3])
        else:
            accent = tuple(data.get("colors", {}).get("#116dff", DEFAULT_BLUE))
        tagline = data.get("tagline", tagline)
        subtagline = data.get("subtagline", subtagline)

    return primary, accent, tagline, subtagline


def read_accounts(limit: int = 5) -> list[dict[str, str]]:
    rows: list[dict[str, str]] = []
    with CSV_PATH.open("r", encoding="utf-8-sig") as f:
        reader = csv.DictReader(f)
        for row in reader:
            normalized = {str(k).strip(): (v or "").strip() for k, v in row.items()}
            if normalized.get("trading_priority", "").upper() == "HIGH":
                rows.append(normalized)
    return rows[:limit]


def duplicate_slide(prs: Presentation, source_idx: int) -> int:
    source = prs.slides[source_idx]
    new_slide = prs.slides.add_slide(source.slide_layout)

    for shape in source.shapes:
        el = copy.deepcopy(shape.element)
        new_slide.shapes._spTree.insert_element_before(el, "p:extLst")

    return len(prs.slides) - 1


def remove_slide(prs: Presentation, index: int) -> None:
    slide_id = prs.slides._sldIdLst[index]
    r_id = slide_id.rId
    prs.part.drop_rel(r_id)
    del prs.slides._sldIdLst[index]


def set_textbox_text(slide, text: str) -> None:
    for shape in slide.shapes:
        if not hasattr(shape, "text_frame") or shape.text_frame is None:
            continue
        if shape.name.lower().startswith("textbox"):
            shape.text_frame.clear()
            p = shape.text_frame.paragraphs[0]
            p.text = text
            p.alignment = PP_ALIGN.LEFT
            if p.runs:
                p.runs[0].font.bold = True
                p.runs[0].font.size = PptPt(36)
            return


def set_content_title(slide, title: str) -> None:
    for shape in slide.shapes:
        if not hasattr(shape, "text_frame") or shape.text_frame is None:
            continue
        if shape.name.lower().startswith("textbox"):
            shape.text_frame.clear()
            p = shape.text_frame.paragraphs[0]
            p.text = f"Ing_TRADE | {title}"
            p.alignment = PP_ALIGN.LEFT
            if p.runs:
                p.runs[0].font.bold = True
                p.runs[0].font.size = PptPt(22)
            return


def add_body_bullets(slide, bullets: list[str], rgb: tuple[int, int, int]) -> None:
    box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.6), Inches(4.7))
    tf = box.text_frame
    tf.clear()

    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.space_after = PptPt(10)
        if p.runs:
            p.runs[0].font.size = PptPt(18)
            p.runs[0].font.color.rgb = PptxRGBColor(*rgb)


def add_page_number(slide, page: int, rgb: tuple[int, int, int]) -> None:
    box = slide.shapes.add_textbox(Inches(11.8), Inches(6.8), Inches(1.0), Inches(0.4))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = str(page)
    p.alignment = PP_ALIGN.RIGHT
    if p.runs:
        p.runs[0].font.size = PptPt(12)
        p.runs[0].font.bold = True
        p.runs[0].font.color.rgb = PptxRGBColor(*rgb)


def build_presentation() -> None:
    primary, accent, tagline, subtagline = load_branding()
    accounts = read_accounts(limit=5)

    prs = Presentation(str(TEMPLATE_PPT))

    # Keep only title and content template slides; move ending slide to the end after content creation.
    ending_idx = 2
    ending_clone_idx = duplicate_slide(prs, ending_idx)
    remove_slide(prs, ending_idx)

    # Slide 1: title
    title_slide = prs.slides[0]
    set_textbox_text(title_slide, "Ing_TRADE Strategic Commercial Plan 2026")

    # Slide 2: executive summary
    slide2 = prs.slides[1]
    set_content_title(slide2, "Executive Summary")
    add_body_bullets(
        slide2,
        [
            "Business focus: used machinery trading for fruit lines and Flexo Folder Gluer (FFG) assets.",
            "Priority channels: direct dealers + global marketplaces with measurable listing velocity.",
            "Primary geographies: Europe and North America, with selective global export coverage.",
            f"Brand promise alignment: {tagline} | {subtagline}",
        ],
        primary,
    )

    # Additional content slides based on the content template slide.
    def new_content_slide(title: str, bullets: list[str]):
        idx = duplicate_slide(prs, 1)
        slide = prs.slides[idx]
        set_content_title(slide, title)
        add_body_bullets(slide, bullets, primary)

    priority_lines = [
        f"{r['company_name']}: {r['business_role']} | Fit: {r['machine_fit']} | Action: {r['next_action']}"
        for r in accounts
    ]

    new_content_slide(
        "Priority Accounts",
        priority_lines,
    )

    new_content_slide(
        "Go-To-Market Strategy",
        [
            "Acquire exclusive inventory from high-priority dealers using monthly feed agreements.",
            "Distribute demand capture via Exapro/Kitmondo/Machineseeker with targeted model alerts.",
            "Qualify every opportunity with machine-fit score, urgency, and expected margin band.",
            "Run weekly pricing intelligence updates from listings, auctions, and closed deal signals.",
        ],
    )

    new_content_slide(
        "Channels and Tools",
        [
            "Channels: dealer outreach, marketplace messaging, association directories, event ecosystems.",
            "Tools: CRM pipeline, web scraping automation, price benchmark dashboard, auction tracker.",
            "Execution rhythm: daily sourcing, weekly qualification, monthly portfolio optimization.",
            "Governance: standardized outreach scripts, SLAs, and partner scorecards.",
        ],
    )

    new_content_slide(
        "Sectorial Numbers",
        [
            "Estimated global used machinery market: ~USD 400B.",
            "Estimated annual growth: ~5% CAGR driven by capex efficiency and sustainability targets.",
            "Regional concentration: Europe + North America represent roughly 60% of active demand.",
            "Implication for Ing_TRADE: prioritize cross-border inventory with fastest turnover categories.",
        ],
    )

    new_content_slide(
        "Top Opportunities",
        [
            "Immediate inventory arbitrage in Iberia via Oversys high-turn listings.",
            "FFG-focused buyer development in North America through Alpine and sector associations.",
            "Consignment and remarketing model adaptation inspired by EquipNet workflows.",
            "Auction intelligence (Surplex/Troostwijk) to capture distressed-value acquisitions.",
        ],
    )

    new_content_slide(
        "90-Day Action Plan",
        [
            "Days 1-30: onboard suppliers, activate listings monitoring, launch outreach sequences.",
            "Days 31-60: close first pilot transactions, refine pricing model with win/loss data.",
            "Days 61-90: scale high-performing channels, formalize partner contracts, track KPIs.",
            "Core KPIs: lead response time, quote-to-deal ratio, inventory turnover, gross margin per asset.",
        ],
    )

    # Add ending slide at the end and remove temporary clone position.
    final_idx = duplicate_slide(prs, ending_clone_idx)
    remove_slide(prs, ending_clone_idx)

    # Ensure ending slide message is aligned with Ing_TRADE.
    ending_slide = prs.slides[final_idx - 1]
    for shape in ending_slide.shapes:
        if hasattr(shape, "text_frame") and shape.text_frame is not None and shape.name.lower().startswith("textbox"):
            shape.text_frame.clear()
            p = shape.text_frame.paragraphs[0]
            p.text = "Thank you\nIng_TRADE"
            p.alignment = PP_ALIGN.LEFT
            if p.runs:
                p.runs[0].font.size = PptPt(38)
                p.runs[0].font.bold = True
            break

    # Force page numbering in every slide.
    for idx, slide in enumerate(prs.slides, start=1):
        add_page_number(slide, idx, accent)

    REPORTS_DIR.mkdir(parents=True, exist_ok=True)
    TEMPLATES_DIR.mkdir(parents=True, exist_ok=True)

    output_ppt = REPORTS_DIR / PRESENTATION_NAME
    output_v2 = REPORTS_DIR / PRESENTATION_V2_NAME

    prs.save(str(output_ppt))
    shutil.copy2(output_ppt, output_v2)

    # Save reusable templates (.potx) for future decks.
    ingecart_template = TEMPLATES_DIR / INGECART_TEMPLATE_NAME
    ingecart_template_v2 = TEMPLATES_DIR / INGECART_TEMPLATE_V2_NAME
    shutil.copy2(TEMPLATE_PPT, ingecart_template)
    shutil.copy2(ingecart_template, ingecart_template_v2)

    ing_trade_template = TEMPLATES_DIR / ING_TRADE_TEMPLATE_NAME
    ing_trade_template_v2 = TEMPLATES_DIR / ING_TRADE_TEMPLATE_V2_NAME
    shutil.copy2(output_ppt, ing_trade_template)
    shutil.copy2(ing_trade_template, ing_trade_template_v2)


def set_table_borders(table):
    for row in table.rows:
        for cell in row.cells:
            tc_pr = cell._tc.get_or_add_tcPr()
            for side in ["top", "left", "bottom", "right"]:
                element = OxmlElement(f"w:{side}")
                element.set(qn("w:val"), "single")
                element.set(qn("w:sz"), "6")
                element.set(qn("w:space"), "0")
                element.set(qn("w:color"), "000000")
                tc_pr.append(element)


def add_offer_template(doc: Document, business_name: str, accent: tuple[int, int, int], accounts: list[dict[str, str]]):
    title = doc.add_heading(f"{business_name} - Commercial Offer Template", level=1)
    title.alignment = WD_ALIGN_PARAGRAPH.LEFT
    if title.runs:
        title.runs[0].font.color.rgb = DocxRGBColor(*accent)

    subtitle = doc.add_paragraph("Use this template for client-facing equipment offers and proposals.")
    subtitle.runs[0].font.size = Pt(11)

    doc.add_heading("1. Client and Opportunity", level=2)
    doc.add_paragraph("Client Name: ______________________________")
    doc.add_paragraph("Company: _________________________________")
    doc.add_paragraph("Country/Region: __________________________")
    doc.add_paragraph("Requested Machines: _______________________")
    doc.add_paragraph("Urgency Level: ____________________________")

    doc.add_heading("2. Offer Scope", level=2)
    doc.add_paragraph("Included equipment and technical condition summary.")
    doc.add_paragraph("Commercial terms, lead time, and logistics scope.")
    doc.add_paragraph("Installation, commissioning, and support assumptions.")

    doc.add_heading("3. Commercial Terms", level=2)
    table = doc.add_table(rows=5, cols=2)
    table.style = "Table Grid"
    labels = [
        "Quoted Price",
        "Incoterm",
        "Delivery Timeline",
        "Payment Milestones",
        "Offer Validity",
    ]
    for i, label in enumerate(labels):
        table.cell(i, 0).text = label
        table.cell(i, 1).text = "______________________________"
    set_table_borders(table)

    doc.add_heading("4. Value Proposition", level=2)
    doc.add_paragraph("Expected productivity impact and bottleneck reduction.")
    doc.add_paragraph("Estimated ROI and payback assumptions.")
    doc.add_paragraph("Risk mitigation and execution safeguards.")

    doc.add_heading("5. Priority Channel References", level=2)
    for account in accounts[:5]:
        doc.add_paragraph(
            f"- {account['company_name']} | {account['business_role']} | {account['why_relevant']}",
            style="List Bullet",
        )

    doc.add_heading("6. Approval and Sign-Off", level=2)
    doc.add_paragraph("Prepared By: ______________________________")
    doc.add_paragraph("Date: _____________________________________")
    doc.add_paragraph("Client Acceptance: _________________________")


def build_word_templates() -> None:
    _, accent, _, _ = load_branding()
    accounts = read_accounts(limit=5)
    TEMPLATES_DIR.mkdir(parents=True, exist_ok=True)

    ingecart_doc = Document()
    add_offer_template(ingecart_doc, "INGECART", accent, accounts)
    ingecart_path = TEMPLATES_DIR / INGECART_OFFER_DOC
    ingecart_doc.save(str(ingecart_path))
    shutil.copy2(ingecart_path, TEMPLATES_DIR / INGECART_OFFER_DOC_V2)

    ing_trade_doc = Document()
    add_offer_template(ing_trade_doc, "ING_TRADE", accent, accounts)
    ing_trade_path = TEMPLATES_DIR / ING_TRADE_OFFER_DOC
    ing_trade_doc.save(str(ing_trade_path))
    shutil.copy2(ing_trade_path, TEMPLATES_DIR / ING_TRADE_OFFER_DOC_V2)


def main() -> None:
    if not TEMPLATE_PPT.exists():
        raise FileNotFoundError(f"Template not found: {TEMPLATE_PPT}")
    if not CSV_PATH.exists():
        raise FileNotFoundError(f"CSV not found: {CSV_PATH}")

    build_presentation()
    build_word_templates()
    print("Generated presentation, version 2 copies, PPT templates, and Word offer templates.")


if __name__ == "__main__":
    main()
