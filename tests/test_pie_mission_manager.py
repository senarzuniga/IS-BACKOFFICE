import tempfile
import unittest
from pathlib import Path

from backoffice.pie import PresentationIntelligenceMissionManager


class TestPieMissionManager(unittest.TestCase):
    def _create_sample_pptx(self, out_file: Path) -> None:
        from pptx import Presentation

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "PIE Test Presentation"
        slide.placeholders[1].text = "- Bullet one\n- Bullet two\n- Bullet three"

        # Add a second slide with a simple table.
        slide2 = prs.slides.add_slide(prs.slide_layouts[5])
        title = slide2.shapes.add_textbox(left=100000, top=100000, width=6000000, height=500000)
        title.text_frame.text = "Data Table"
        rows, cols = 3, 2
        table = slide2.shapes.add_table(rows, cols, 100000, 600000, 5000000, 1200000).table
        table.cell(0, 0).text = "KPI"
        table.cell(0, 1).text = "Value"
        table.cell(1, 0).text = "OEE"
        table.cell(1, 1).text = "89%"
        table.cell(2, 0).text = "Orders"
        table.cell(2, 1).text = "3820"

        prs.save(str(out_file))

    def test_run_generates_dual_html_outputs(self):
        with tempfile.TemporaryDirectory() as tmp:
            tmp_path = Path(tmp)
            source = tmp_path / "sample.pptx"
            out_root = tmp_path / "pie_out"
            self._create_sample_pptx(source)

            manager = PresentationIntelligenceMissionManager()
            result = manager.run(str(source), str(out_root))

            self.assertTrue(Path(result.version_1_html).exists())
            self.assertTrue(Path(result.version_2_html).exists())
            self.assertTrue(Path(result.corporate_css).exists())
            self.assertTrue(Path(result.components_matrix).exists())
            self.assertTrue(Path(result.technical_report).exists())
            self.assertTrue(Path(result.differences_report).exists())
            self.assertTrue(Path(result.enterprise_memory_file).exists())
            self.assertTrue(Path(result.knowledge_hub_file).exists())


if __name__ == "__main__":
    unittest.main()
