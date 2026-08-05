import tempfile
import unittest
from pathlib import Path

from backoffice.dipc import DocumentIntelligencePublishingCenter


class TestDipcMissionManager(unittest.TestCase):
    def _create_sample_pptx(self, out_file: Path) -> None:
        from pptx import Presentation

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Industrial Roadmap"
        slide.placeholders[1].text = "Timeline\nPhase 1\nPhase 2\nPhase 3"

        slide2 = prs.slides.add_slide(prs.slide_layouts[5])
        title = slide2.shapes.add_textbox(left=100000, top=100000, width=6000000, height=500000)
        title.text_frame.text = "Risk Matrix"
        body = slide2.shapes.add_textbox(left=100000, top=600000, width=7000000, height=1300000)
        body.text_frame.text = "High impact, low probability\nLow impact, high probability"

        prs.save(str(out_file))

    def test_build_and_command_mission_generate_versions_and_outputs(self):
        with tempfile.TemporaryDirectory() as tmp:
            tmp_path = Path(tmp)
            source = tmp_path / "sample.pptx"
            output_root = tmp_path / "dipc"
            self._create_sample_pptx(source)

            center = DocumentIntelligencePublishingCenter()
            initial = center.build_from_powerpoint(str(source), str(output_root))

            self.assertTrue(Path(initial.document_model_path).exists())
            self.assertTrue(Path(initial.publication_outputs["html"]).exists())
            self.assertTrue(Path(initial.publication_outputs["pdf"]).exists())
            self.assertTrue(Path(initial.publication_outputs["odt"]).exists())
            self.assertTrue(Path(initial.preview_manifest_path).exists())

            updated = center.apply_mission(initial.document_model_path, "Hazlo más ejecutivo y añade gráficos")
            self.assertTrue(Path(updated.document_model_path).exists())
            self.assertTrue(Path(updated.publication_outputs["html"]).exists())

            model_text = Path(updated.document_model_path).read_text(encoding="utf-8")
            self.assertIn("Executive Summary", model_text)
            self.assertIn("version_history", model_text)


if __name__ == "__main__":
    unittest.main()
